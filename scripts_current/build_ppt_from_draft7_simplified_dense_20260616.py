# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BASE_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_20260616.py"
spec = importlib.util.spec_from_file_location("base_ppt", BASE_SCRIPT)
base = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(base)

WORK = ROOT / "reports" / "ppt_from_draft7_simplified_dense_20260616"
WORK.mkdir(parents=True, exist_ok=True)
OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_简化充实版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_简化充实版PPT_QA报告.md"


C = base.COLORS
DARK = RGBColor(18, 30, 44)
SLATE = RGBColor(55, 72, 88)
LIGHT = RGBColor(246, 249, 252)
BLUE = RGBColor(33, 91, 150)
TEAL = RGBColor(26, 124, 118)
AMBER = RGBColor(188, 131, 36)
RED = RGBColor(176, 69, 68)
GREEN = RGBColor(63, 142, 94)


def add_bg(slide, accent=TEAL):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.11))
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(7.23), Inches(12.48), Inches(0.018))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C["line"]
    bottom.line.fill.background()


def text(slide, body, x, y, w, h, size=11.5, color=DARK, bold=False, align="left", valign="top"):
    return base.add_textbox(
        slide,
        body,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        margin=0.05,
    )


def title(slide, main, sub="", accent=TEAL):
    text(slide, main, 0.52, 0.34, 11.7, 0.48, size=22, bold=True, color=DARK)
    if sub:
        text(slide, sub, 0.54, 0.84, 11.4, 0.34, size=9.4, color=C["muted"])
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.19), Inches(1.08), Inches(0.035))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()


def footer(slide, idx, total):
    text(slide, "FZYC-Mol | 简化充实版汇报", 0.52, 7.26, 3.2, 0.15, size=6.3, color=C["muted"])
    text(slide, f"{idx:02d}/{total}", 12.18, 7.26, 0.68, 0.15, size=6.3, color=C["muted"], align="right")


def bullet_block(slide, items, x, y, w, h, size=10.2, accent=TEAL, gap=0.04):
    top = y
    line_h = h / max(len(items), 1)
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(top + 0.095), Inches(0.07), Inches(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        text(slide, item, x + 0.17, top, w - 0.18, max(0.32, line_h - gap), size=size, color=DARK)
        top += line_h


def callout(slide, head, body, x, y, w, h, accent=TEAL, fill=C["white"]):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = C["line"]
    box.line.width = Pt(0.75)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    text(slide, head, x + 0.18, y + 0.09, w - 0.30, 0.22, size=9.0, color=accent, bold=True)
    text(slide, body, x + 0.18, y + 0.36, w - 0.30, h - 0.45, size=9.5, color=DARK)


def mini_metric(slide, value, label, x, y, w=1.52, accent=TEAL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.62))
    box.fill.solid()
    box.fill.fore_color.rgb = C["white"]
    box.line.color.rgb = C["line"]
    box.line.width = Pt(0.65)
    text(slide, value, x + 0.08, y + 0.08, w - 0.16, 0.20, size=13.0, color=accent, bold=True, align="center")
    text(slide, label, x + 0.08, y + 0.34, w - 0.16, 0.18, size=6.6, color=C["muted"], align="center")


def takeaway(slide, body, accent=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(6.74), Inches(12.28), Inches(0.37))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C["pale_teal"] if accent == TEAL else C["pale_blue"]
    bar.line.fill.background()
    text(slide, "Takeaway  " + body, 0.68, 6.82, 11.9, 0.18, size=9.2, color=DARK, bold=True)


def add_line(slide, x1, y1, x2, y2, color=C["line"], width=0.9):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_three_line_table(slide, rows, x, y, w, h, widths=None, font_size=7.8, header_size=8.1, accent=TEAL):
    if not rows:
        return
    cols = len(rows[0])
    widths = widths or [1 / cols] * cols
    total = sum(widths)
    widths = [ww / total * w for ww in widths]
    row_h = h / len(rows)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = C["white"]
    panel.line.fill.background()
    add_line(slide, x, y, x + w, y, color=DARK, width=1.0)
    add_line(slide, x, y + row_h, x + w, y + row_h, color=DARK, width=0.75)
    add_line(slide, x, y + h, x + w, y + h, color=DARK, width=1.0)
    cx = x
    for j in range(cols):
        cell_x = cx
        for i, row in enumerate(rows):
            size = header_size if i == 0 else font_size
            color = accent if i == 0 else DARK
            bold = i == 0
            body = row[j] if j < len(row) else ""
            text(slide, body, cell_x + 0.03, y + i * row_h + 0.035, widths[j] - 0.06, row_h - 0.055, size=size, color=color, bold=bold)
        cx += widths[j]


def process_strip(slide, items, x, y, w, h, accent=TEAL):
    gap = 0.13
    item_w = (w - gap * (len(items) - 1)) / len(items)
    for i, (head, body) in enumerate(items):
        bx = x + i * (item_w + gap)
        callout(slide, head, body, bx, y, item_w, h, accent=accent, fill=C["white"])
        if i < len(items) - 1:
            base.add_arrow(slide, bx + item_w + 0.02, y + h / 2, bx + item_w + gap - 0.04, y + h / 2, color=C["muted"])


def note(slide, body):
    base.add_notes(slide, body)


def compact_table1(src):
    table = src["tables"][0]
    rows = [["数据集", "任务", "规模/阳性率", "评价"]]
    for r in table[1:6]:
        rows.append([r[0], r[1].replace("MoleculeNet；", "").replace("；", "/"), r[2], r[3]])
    return rows


def compact_table2(src):
    table = src["tables"][1]
    rows = [["终点", "指标", "最终保留", "对照/观测最优", "结论"]]
    for r in table[1:6]:
        rows.append([r[0], r[1], r[3], r[4], r[5]])
    return rows


def compact_table3(src):
    return src["tables"][2][:4]


def compact_table4(src):
    return src["tables"][3][:5]


def compact_table5(src):
    return src["tables"][4][:6]


def build_deck(src, figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []
    total = 20

    # 1
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    text(s, "FZYC-Mol", 0.62, 0.46, 6.2, 0.55, size=28, bold=True, color=DARK)
    text(s, "面向可靠性审计的验证集治理分子性质预测框架", 0.66, 1.06, 9.6, 0.45, size=15, color=SLATE)
    text(s, "简化充实版汇报 | 基于《初稿-7》自动整理", 0.68, 1.52, 6.0, 0.28, size=9.6, color=C["muted"])
    callout(s, "一句话问题", "分子性质预测不只需要更高平均性能，还需要说明何时可用、何时应谨慎，以及验证集选择是否引入乐观偏差。", 0.72, 2.15, 3.9, 1.32, accent=BLUE)
    callout(s, "一句话方法", "FZYC-Mol 将多源表示、候选专家、验证集选择器、适用域门控、风险分数与证据输出合并为一套冻结测试流程。", 4.88, 2.15, 3.9, 1.32, accent=TEAL)
    callout(s, "一句话结论", "模型贡献集中在验证集治理、选择性增益、可靠性审计和负结果呈现，而不是宣称在所有终点上普遍超过强基线。", 9.04, 2.15, 3.18, 1.32, accent=AMBER)
    process_strip(
        s,
        [("输入", "SMILES、任务标签、划分协议"),
         ("候选", "多表示与强基线专家池"),
         ("治理", "验证集选择与门控冻结"),
         ("输出", "预测、风险、解释、边界")],
        0.82, 4.10, 11.7, 1.02, accent=BLUE,
    )
    takeaway(s, "这版 PPT 采用更少页数、更高单页信息量，适合 15-20 分钟完整汇报。", accent=BLUE)
    footer(s, 1, total)
    note(s, "开场时先说明本汇报不是把论文逐节复述，而是围绕问题、方法、证据和边界组织。强调 FZYC-Mol 的定位是可靠性治理框架，核心是验证集选择、冻结测试和风险输出。")

    # 2
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "汇报结构压缩为一条主线", "从问题到方法，从证据到边界，减少过场页但保留完整流程。", TEAL)
    process_strip(
        s,
        [("1. 为什么", "现有模型在数据划分、外推和可靠性上容易给出过度乐观结论"),
         ("2. 做什么", "构建验证集治理流程，保留有效候选并拒绝不稳定候选"),
         ("3. 怎么证", "MoleculeNet、TDC、OOD、保形、消融和负结果共同支撑"),
         ("4. 如何用", "输出预测、适用域、风险覆盖、解释和复现接口")],
        0.72, 1.48, 11.9, 1.18, accent=TEAL,
    )
    add_three_line_table(
        s,
        [["板块", "核心问题", "对应证据", "汇报页"],
         ["背景与创新", "为什么不能只看平均指标", "问题拆解、贡献声明", "3-4"],
         ["方法流程", "如何避免测试集调参", "workflow、架构、选择器公式", "5-8"],
         ["结果验证", "哪些终点有增益，哪些只保留", "主结果、TDC、风险覆盖", "9-14"],
         ["边界与复现", "失败在哪里，如何避免夸大", "负结果、解释、局限、复现", "15-20"]],
        0.82, 3.20, 11.7, 2.65, widths=[0.18, 0.32, 0.30, 0.20], font_size=8.6, header_size=8.9, accent=TEAL,
    )
    takeaway(s, "整套汇报遵循 problem-to-solution arc：先解释瓶颈，再展示框架、验证、边界和可复用价值。")
    footer(s, 2, total)
    note(s, "这一页作为路线图。汇报时告诉听众，后面不会每个实验单独铺开，而是把相近实验合并到同一证据页中；每一页都有主张、证据和谨慎解释。")

    # 3
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, AMBER)
    title(s, "研究瓶颈：性能、可迁移性和可信度常被混在一起讨论", "FZYC-Mol 把“预测得准”和“知道何时不可靠”分开评估。", AMBER)
    callout(s, "瓶颈 1：验证集选择", "候选模型多、调参路径长时，验证集排名可能不稳定；若测试集被反复查看，会形成隐性乐观偏差。", 0.70, 1.45, 3.85, 1.35, accent=AMBER)
    callout(s, "瓶颈 2：结构外推", "随机划分下的高分不一定能转移到 scaffold、低相似度或活性悬崖样本，尤其是 ADMET 和毒性任务。", 4.76, 1.45, 3.85, 1.35, accent=RED)
    callout(s, "瓶颈 3：失败解释", "若只报告胜利结果，读者无法判断方法的适用域；负结果和失败案例需要进入主文逻辑。", 8.82, 1.45, 3.50, 1.35, accent=TEAL)
    bullet_block(
        s,
        ["本文不是把所有模型简单平均，而是将候选专家登记、验证集门控和测试集冻结作为主流程。",
         "主张限定为“选择性增益与可靠性审计”，避免写成跨任务普遍最优。",
         "边界案例包括 FreeSolv、bRo5 化学空间、低相似度样本和活性悬崖。",
         "审稿风险主要来自：结果是否过度解释、验证是否泄漏、负结果是否充分呈现。"],
        0.90, 3.35, 11.05, 2.35, size=10.4, accent=AMBER,
    )
    takeaway(s, "因此，方法创新的重点是治理流程和可靠性输出，而不是单一模型结构的复杂化。", accent=AMBER)
    footer(s, 3, total)
    note(s, "这一页要帮助听众理解为什么论文选择了比较保守的写法。强调分子性质预测的关键挑战不是单个排行榜数字，而是不同划分、外推场景和失败样本下的可信度。")

    # 4
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    title(s, "核心贡献：把候选模型选择变成可审计的冻结流程", "贡献强调流程治理、可靠性、边界和复现，而非无条件性能宣称。", BLUE)
    mini_metric(s, "6", "主文数据表", 0.78, 1.35, accent=BLUE)
    mini_metric(s, "13", "内嵌结果图", 2.45, 1.35, accent=TEAL)
    mini_metric(s, "80/90/95%", "保形覆盖", 4.12, 1.35, w=1.75, accent=GREEN)
    mini_metric(s, "5/17/0", "外部 win/tie/loss", 6.04, 1.35, w=1.75, accent=AMBER)
    mini_metric(s, "冻结", "测试集只评一次", 7.96, 1.35, w=1.72, accent=RED)
    callout(s, "贡献 A：验证集治理", "所有候选模型先登记，再由验证集证据选择；测试集仅用于冻结策略后的最终评估，降低隐性调参风险。", 0.78, 2.35, 3.72, 1.25, accent=BLUE)
    callout(s, "贡献 B：选择性融合", "当融合、补救头或适配器在验证集上不稳定时，不进入主性能结论，只作为负结果或后续接口保留。", 4.80, 2.35, 3.72, 1.25, accent=TEAL)
    callout(s, "贡献 C：可靠性审计", "风险覆盖、保形预测、低相似度分层、活性悬崖和失败案例共同限制外推主张。", 8.82, 2.35, 3.34, 1.25, accent=AMBER)
    add_three_line_table(
        s,
        [["输出类型", "论文中承担的作用", "汇报时如何解释"],
         ["预测性能", "展示主任务是否有效", "只在验证接受后报告最终测试"],
         ["风险分数", "识别错误或高误差风险", "作为审计信号，不替代性能指标"],
         ["负结果", "限制夸大结论", "说明何处不能声称改进"],
         ["解释模块", "提供关联性证据", "不写成因果机制证明"]],
        0.86, 4.20, 11.35, 1.72, widths=[0.18, 0.38, 0.44], font_size=8.4, header_size=8.7, accent=BLUE,
    )
    takeaway(s, "这篇小论文的安全写法是：可验证、可审计、可复现、不过度外推。", accent=BLUE)
    footer(s, 4, total)
    note(s, "这一页可以作为创新点页。需要避免说“全面优于所有方法”，改说“在部分终点实现选择性增益，并系统报告不增益或不稳定的候选”。")

    # 5
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "完整 workflow：从输入分子到证据输出", "简化后保留一张总流程图，旁边集中解释每一步的输入、决策和输出。", TEAL)
    base.add_picture_fit(s, figs[1], 0.58, 1.32, 8.10, 4.95, frame=True)
    callout(s, "输入层", "SMILES 与标签首先进入标准化、去重、划分和任务类型识别；该步骤决定后续评价协议是否可解释。", 8.95, 1.34, 3.55, 0.96, accent=TEAL)
    callout(s, "候选层", "多视图表示和专家池同时保留：传统指纹、图神经网络、强表格基线与可控补救模块。", 8.95, 2.48, 3.55, 0.96, accent=BLUE)
    callout(s, "治理层", "选择器只读取验证集；通过者进入冻结测试，不通过者进入负结果记录或后续验证接口。", 8.95, 3.62, 3.55, 0.96, accent=AMBER)
    callout(s, "输出层", "除最终预测外，还输出 AD gate、uncertainty、risk-coverage、保形覆盖和解释证据。", 8.95, 4.76, 3.55, 0.96, accent=GREEN)
    takeaway(s, "workflow 的核心不是堆模型，而是让每一次模型进入最终策略都有验证集证据。")
    footer(s, 5, total)
    note(s, "讲图时从左到右说明：数据进入、表示生成、候选专家、验证选择、冻结测试和证据输出。强调每一类候选如果不能通过验证集，不会在测试集结果中强行使用。")

    # 6
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    title(s, "模型结构：多源表示、专家矩阵与证据输出并行", "结构图展示预测模块与审计模块并列，而不是把可靠性作为事后附加。", BLUE)
    base.add_picture_fit(s, figs[0], 0.58, 1.32, 8.05, 4.88, frame=True)
    bullet_block(
        s,
        ["多源表示：SMILES、指纹、图结构和任务相关描述共同进入候选专家池。",
         "专家预测矩阵：不同模型家族产生可比较候选，选择器在验证集上判断保留或拒绝。",
         "不确定性与 AD：适用域门控和风险分数用于提示预测边界，而不是直接改写标签。",
         "证据输出：最终报告同时呈现性能、风险覆盖、校准、解释和失败案例。"],
        8.98, 1.48, 3.40, 3.10, size=9.6, accent=BLUE,
    )
    text(s, "简式公式", 9.06, 4.80, 1.4, 0.22, size=8.4, color=BLUE, bold=True)
    text(s, "m* = argmin_m L_val(m)\nŷ(x) = Σ_k w_k(x) f_k(x)\nr(x) = αu(x) + βd_AD(x) + γc(x)", 9.06, 5.12, 3.26, 0.78, size=10.4, color=DARK)
    takeaway(s, "架构设计服务于“验证选择 + 风险审计”，因此结构图必须和实验设计一起解释。", accent=BLUE)
    footer(s, 6, total)
    note(s, "这一页重点解释模型结构图。公式只作为汇报用简式：m* 表示验证集选出的候选；ŷ 表示融合预测；r 表示由不确定性、适用域距离和校准误差组成的风险信号。不要把这些公式讲成额外的新实验。")

    # 7
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "数据与评价协议：同时覆盖内部基准和外部 ADMET 任务", "数据表压缩为 PPT 三线表，突出任务类型、规模、指标和划分。", TEAL)
    add_three_line_table(
        s,
        compact_table1(src),
        0.66, 1.38, 7.50, 2.75, widths=[0.18, 0.35, 0.22, 0.25], font_size=8.0, header_size=8.4, accent=TEAL,
    )
    callout(s, "划分策略", "随机划分、scaffold 划分、结构分离和低相似度分层共同构成压力测试；不同划分结果不能混用。", 8.52, 1.40, 3.82, 1.02, accent=TEAL)
    callout(s, "评价指标", "回归终点采用 RMSE/MAE；分类终点采用 ROC-AUC、PR-AUC、Brier、ECE 和 fixed precision recall 等。", 8.52, 2.62, 3.82, 1.02, accent=BLUE)
    callout(s, "统计逻辑", "多 seed、配对比较、rank audit 与 optimism gap 用于区分真实增益和验证集偶然优势。", 8.52, 3.84, 3.82, 1.02, accent=AMBER)
    text(s, "报告原则：测试集只在策略冻结后评估；所有未通过门控的模块作为负结果或补充接口记录。", 0.80, 4.72, 7.20, 0.46, size=10.2, color=SLATE, bold=True)
    takeaway(s, "数据与划分页要让审稿人相信：性能比较来自同一评价协议，而非后验挑选。")
    footer(s, 7, total)
    note(s, "这一页压缩了原稿的数据集表。汇报时说明表格只展示代表性数据集，完整表格在论文中；重点不是背数据集数量，而是说明为什么需要多种划分和外部 ADMET 任务。")

    # 8
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, AMBER)
    title(s, "验证集治理：候选登记、选择、冻结测试和负结果归档", "这一页合并 selector、gate、test-freeze 和审稿风险控制。", AMBER)
    base.add_picture_fit(s, figs[2], 0.56, 1.32, 7.95, 4.80, frame=True)
    callout(s, "候选登记", "所有强基线、融合候选、补救头和适配器先在候选池登记，避免只展示成功模型。", 8.82, 1.35, 3.55, 0.92, accent=AMBER)
    callout(s, "验证选择", "selector 根据 L_val、校准、风险和稳定性决定是否保留；测试集不参与该决策。", 8.82, 2.42, 3.55, 0.92, accent=TEAL)
    callout(s, "冻结评估", "策略冻结后一次性在测试集评估，rank audit 用来检查验证排名与测试排名是否错位。", 8.82, 3.49, 3.55, 0.92, accent=BLUE)
    callout(s, "负结果处理", "未通过门控的模块不写入主性能结论，只作为边界、失败案例或后续验证接口。", 8.82, 4.56, 3.55, 0.92, accent=RED)
    text(s, "决策式：Accept(m) = I[Δval ≥ τ 且 calibration/risk 未恶化]", 8.90, 5.78, 3.35, 0.34, size=9.2, color=DARK, bold=True)
    takeaway(s, "验证集治理的价值在于减少后验选择，而不是保证每个候选都提升。", accent=AMBER)
    footer(s, 8, total)
    note(s, "这一页可以讲审稿人最关心的“有没有测试集调参”。要明确测试集只作为最终评估，不用于选择模型；失败候选也被记录，因此结果更可信。")

    # 9
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    title(s, "MoleculeNet：模型家族排名显示终点间差异明显", "rank map 用来说明没有单一模型在所有任务上稳定占优。", BLUE)
    base.add_picture_fit(s, figs[3], 0.56, 1.28, 6.18, 4.58, frame=True)
    add_three_line_table(
        s,
        compact_table2(src),
        6.98, 1.36, 5.45, 3.12, widths=[0.16, 0.14, 0.22, 0.32, 0.16], font_size=6.8, header_size=7.0, accent=BLUE,
    )
    bullet_block(
        s,
        ["ESOL、Lipo、BBBP、BACE、ClinTox 等终点呈现不同最优模型家族。",
         "FreeSolv 被保留为边界案例：重构缩小差距，但不夸大为超过 Chemprop。",
         "最终保留策略只由验证集证据决定，测试表现用于冻结后报告。"],
        7.10, 4.72, 5.15, 1.10, size=8.8, accent=BLUE,
    )
    takeaway(s, "主结果页的语气应是“选择性保留与边界清楚”，而不是“一招通吃”。", accent=BLUE)
    footer(s, 9, total)
    note(s, "讲图时先看热图颜色和排名，再转向右侧主结果表。强调不同终点最优模型不同，因此选择器的存在有必要。FreeSolv 要谨慎讲，是边界案例，不是成功案例。")

    # 10
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "MoleculeNet 主结果：验证保留策略带来选择性增益", "将性能比较和验证保留逻辑放在同一页，避免把结果讲成孤立数字。", TEAL)
    base.add_picture_fit(s, figs[4], 0.56, 1.30, 6.55, 4.28, frame=True)
    callout(s, "回归终点", "ESOL 和 Lipophilicity 的最终保留策略与验证证据一致；FreeSolv 不写作确定性提升，而作为物理相互作用相关边界。", 7.44, 1.35, 4.72, 0.92, accent=TEAL)
    callout(s, "分类终点", "BBBP、BACE 与 ClinTox 需要同时看 ROC-AUC、PR-AUC、阳性率和校准指标，避免稀有阳性任务被 ROC-AUC 掩盖。", 7.44, 2.42, 4.72, 0.92, accent=BLUE)
    callout(s, "统计表述", "建议报告均值 ± 标准差、seed 配对、fixed precision recall 和 rank audit；结论应围绕稳定性而非单次最高值。", 7.44, 3.49, 4.72, 0.92, accent=AMBER)
    callout(s, "审稿解释", "当最强单模型已经足够强时，FZYC-Mol 的优势体现在拒绝不稳定候选和输出可靠性证据。", 7.44, 4.56, 4.72, 0.92, accent=GREEN)
    takeaway(s, "主结果支持“验证治理有效”，但不支持“每个终点绝对最优”的夸张写法。")
    footer(s, 10, total)
    note(s, "这一页要把主结果讲细一些。解释柱状/点图中的最终保留策略，与右侧四个解释框对应。尤其强调分类终点不要只说 ROC-AUC。")

    # 11
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, AMBER)
    title(s, "补救头与融合门控：只接受验证集支持的局部改进", "rescue 和 fusion 是可控候选，不是默认进入最终模型的装饰模块。", AMBER)
    base.add_picture_fit(s, figs[5], 0.58, 1.30, 5.70, 2.24, frame=True)
    base.add_picture_fit(s, figs[6], 0.58, 3.82, 5.70, 2.20, frame=True)
    callout(s, "Targeted rebuild gate", "补救头只针对验证集暴露的问题；若 Δval 未达到阈值或校准/风险恶化，则保持原策略。", 6.68, 1.30, 5.60, 0.92, accent=AMBER)
    callout(s, "Multimethod fusion gate", "融合不是简单平均，而是在每个终点上判断是否带来验证集正增益。BBBP 和 ClinTox 体现选择性增益。", 6.68, 2.38, 5.60, 0.92, accent=TEAL)
    callout(s, "Uncertainty weighting", "不确定性权重用于减少高风险候选对最终预测的影响，但其收益必须由校准和风险覆盖共同验证。", 6.68, 3.46, 5.60, 0.92, accent=BLUE)
    callout(s, "Negative control", "若模块未通过门控，它的存在不是失败遮掩，而是说明选择器确实会拒绝候选。", 6.68, 4.54, 5.60, 0.92, accent=RED)
    takeaway(s, "补救与融合的合理写法是“验证集接受的候选进入最终策略，未接受者进入边界报告”。", accent=AMBER)
    footer(s, 11, total)
    note(s, "这一页把两张门控图合并。上图讲定向重构，下图讲多方法融合。核心语言是“candidate gate”，不要讲成每个模块必然提升。")

    # 12
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    title(s, "外部 TDC ADMET：强基线惩罚与融合门控共同验证外推", "外部任务强调评价协议真实性，而不是只追求更多数据集数量。", BLUE)
    base.add_picture_fit(s, figs[7], 0.56, 1.28, 5.55, 2.38, frame=True)
    base.add_picture_fit(s, figs[8], 0.56, 3.88, 5.55, 2.18, frame=True)
    add_three_line_table(
        s,
        compact_table3(src),
        6.52, 1.34, 5.72, 1.64, widths=[0.24, 0.38, 0.38], font_size=7.3, header_size=7.8, accent=BLUE,
    )
    bullet_block(
        s,
        ["Caco2、HIA、Pgp 等外部终点出现选择性增益，但多数终点保留原策略。",
         "22 个外部终点的 win/tie/loss = 5/17/0，说明方法避免了普遍提升的过度主张。",
         "官方划分与 scaffold 划分差异明显，划分本身就是评价协议的一部分。",
         "TDC 结果用于支持可迁移性审计，不应写成独立临床有效性证据。"],
        6.64, 3.35, 5.35, 1.95, size=8.8, accent=BLUE,
    )
    takeaway(s, "外部验证的关键结论是“有选择性增益且无系统性退化”，不是“全部外部任务显著提升”。", accent=BLUE)
    footer(s, 12, total)
    note(s, "这一页用上下两张外部图加一个摘要表。先讲强基线在 scaffold 下受到惩罚，再讲融合门控在 TDC 上的接受与保留。win/tie/loss 是一个非常适合答辩的数字。")

    # 13
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, GREEN)
    title(s, "可靠性审计：risk-coverage 与保形覆盖限制高风险预测", "这页把风险曲线、保形预测和校准指标合并成可信度证据。", GREEN)
    base.add_picture_fit(s, figs[9], 0.56, 1.28, 6.25, 4.70, frame=True)
    add_three_line_table(
        s,
        compact_table4(src),
        7.15, 1.38, 5.05, 1.70, widths=[0.38, 0.25, 0.37], font_size=7.6, header_size=7.9, accent=GREEN,
    )
    callout(s, "risk-coverage", "随着 coverage retained 降低，高风险样本被优先剔除；分类终点的错误风险识别更稳定。", 7.15, 3.36, 5.05, 0.76, accent=GREEN)
    callout(s, "conformal", "保形覆盖按 80%、90%、95% 报告；它提供覆盖率语境，但不等同于性能提升。", 7.15, 4.28, 5.05, 0.76, accent=BLUE)
    callout(s, "calibration", "Brier、ECE 和固定精度召回用于补充 ROC-AUC，尤其适合 ClinTox 这类低阳性率任务。", 7.15, 5.20, 5.05, 0.76, accent=AMBER)
    takeaway(s, "可靠性结果应写成“帮助识别风险与限定适用域”，而不是替代主性能指标。", accent=GREEN)
    footer(s, 13, total)
    note(s, "这一页是可靠性核心。解释 risk-coverage 曲线时，说明模型给出风险分数后，可以选择保留低风险样本进行更稳健预测。保形预测用于覆盖率保障，不能讲成准确率提升。")

    # 14
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, RED)
    title(s, "OOD、低相似度与活性悬崖：压力测试揭示外推边界", "该页强调从随机划分到结构分离划分的性能落差。", RED)
    base.add_picture_fit(s, figs[10], 0.58, 1.28, 7.02, 4.75, frame=True)
    bullet_block(
        s,
        ["从 random 到 scaffold 或 structure-separated split 后，部分终点性能明显下降，说明随机划分不能代表真实外推。",
         "低相似度样本应按互斥 Tanimoto bin 报告：>0.7、0.5-0.7、<0.5，并同时给出性能、校准和风险富集。",
         "MoleculeACE 活性悬崖需要区分整体 RMSE/MAE 与 cliff subset 误差；预测差异与真实差异相关性可作为补充证据。",
         "失败案例不能只列成功样本，应包含低相似度失败和活性悬崖失败，避免选择性呈现。"],
        7.92, 1.38, 4.55, 3.35, size=9.2, accent=RED,
    )
    callout(s, "审稿风险", "若 OOD 结果只作为补充材料，主文中仍需明确说明适用域边界；否则性能主张容易被认为外推过度。", 7.96, 5.05, 4.46, 0.82, accent=RED)
    takeaway(s, "OOD 页不是削弱论文，而是证明作者知道模型在哪些化学空间里不可靠。", accent=RED)
    footer(s, 14, total)
    note(s, "这一页要主动承认边界。答辩时可以说：压力测试不是为了证明模型永远有效，而是给出什么时候不应信任模型的证据。")

    # 15
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, AMBER)
    title(s, "消融、selector appendix 与负结果：保留失败证据能增强可信度", "把消融图、selector 结果和负结果表压缩到一页说明。", AMBER)
    base.add_picture_fit(s, figs[11], 0.58, 1.28, 5.72, 2.62, frame=True)
    add_three_line_table(
        s,
        compact_table5(src),
        0.66, 4.28, 5.62, 1.65, widths=[0.28, 0.28, 0.44], font_size=6.8, header_size=7.1, accent=AMBER,
    )
    callout(s, "消融矩阵", "Full / best single / simple mean / w/o selector / w/o fusion / w/o AD gate / w/o uncertainty weighting 等用于定位贡献来源。", 6.70, 1.30, 5.55, 0.90, accent=AMBER)
    callout(s, "负结果写法", "FreeSolv、bRo5、轻量适配器和 3D-lite 均需写成边界或待验证方向，不能包装为确定性成功。", 6.70, 2.36, 5.55, 0.90, accent=RED)
    callout(s, "selector appendix", "附录中的验证保留策略用于追踪候选接受/拒绝路径，帮助审稿人复核选择是否透明。", 6.70, 3.42, 5.55, 0.90, accent=TEAL)
    callout(s, "主文取舍", "主文只保留能支撑核心论证的消融摘要，完整长表和失败样本放补充材料更清晰。", 6.70, 4.48, 5.55, 0.90, accent=BLUE)
    takeaway(s, "负结果不是扣分项；对高影响力写法而言，它是防止学术不端风险的保护层。", accent=AMBER)
    footer(s, 15, total)
    note(s, "这一页回应用户此前特别担心的学术不端风险。强调所有没有跑齐或没有显著收益的实验都不能写成完成结果；能证实的写结果，不能证实的写限制。")

    # 16
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "可解释性：基序与片段证据只能支持关联解释", "化学解释页要保守，避免把富集或归因写成因果机制。", TEAL)
    base.add_picture_fit(s, figs[12], 0.58, 1.28, 7.18, 4.78, frame=True)
    callout(s, "可以说", "模型关注的基序或片段与任务相关信号存在一致性；该证据可辅助发现候选风险区域。", 8.10, 1.35, 4.00, 0.94, accent=TEAL)
    callout(s, "不宜说", "不能仅凭 motif attribution 或 fragment enrichment 断言化学因果机制，也不能替代湿实验验证。", 8.10, 2.45, 4.00, 0.94, accent=RED)
    callout(s, "统计要求", "若主文保留解释结果，应报告最小支持度、效应量、p 值或 FDR，并说明多重检验控制。", 8.10, 3.55, 4.00, 0.94, accent=AMBER)
    callout(s, "汇报定位", "解释模块用于增强透明度和错误分析，不作为主性能提升的直接证据。", 8.10, 4.65, 4.00, 0.94, accent=BLUE)
    takeaway(s, "解释结果服务于模型审计和假设生成，不能被写成已证明的化学机制。")
    footer(s, 16, total)
    note(s, "这一页必须谨慎讲。建议用“关联”“提示”“辅助审计”等词，不使用“证明”“决定”“因果”等词。")

    # 17
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    title(s, "复现与提交准备：把流程、数据和代码证据打包清楚", "投稿前的 PPT 需要能解释“别人如何复查这套流程”。", BLUE)
    add_three_line_table(
        s,
        [["项目", "需要保留的材料", "投稿/答辩作用"],
         ["数据", "原始来源、划分种子、清洗规则、低相似度 bin", "证明不是后验挑选数据"],
         ["模型", "候选登记表、超参数范围、冻结模型权重", "证明选择器前后一致"],
         ["结果", "主表、补充长表、seed 配对、失败样本", "支持主张并呈现边界"],
         ["代码", "训练脚本、评估脚本、图表脚本、环境文件", "允许审稿复核和二次运行"],
         ["声明", "Data availability、Code availability、负结果说明", "降低可重复性和合规风险"]],
        0.78, 1.35, 11.70, 3.45, widths=[0.16, 0.42, 0.42], font_size=8.4, header_size=8.8, accent=BLUE,
    )
    bullet_block(
        s,
        ["演示时应主动说明：图表来自《初稿-7》内嵌结果，PPT 未修改科学数据。",
         "对于仍需后续验证的模块，应在论文与答辩中保持同一谨慎表述。",
         "补充材料建议放完整长表、低相似度分层、活性悬崖案例和失败样本列表。"],
        0.92, 5.20, 11.10, 0.92, size=9.3, accent=BLUE,
    )
    takeaway(s, "复现页的作用是让审稿人看到所有关键选择都有记录、可追踪、可复核。", accent=BLUE)
    footer(s, 17, total)
    note(s, "这一页是提交前自查页。说明每类材料要保存在哪里，为什么能降低学术不端风险。")

    # 18
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, RED)
    title(s, "局限性与风险控制：主动限定主张比回避问题更稳妥", "这页用于答辩和投稿前自审。", RED)
    callout(s, "局限 1：外部数据", "部分外部任务依赖公开数据和官方划分，不能等同于独立湿实验或前瞻性盲测。", 0.82, 1.34, 3.66, 1.00, accent=RED)
    callout(s, "局限 2：低相似度空间", "模型在低 Tanimoto 或结构分离样本上仍可能出现误差放大，需要分层报告和风险提示。", 4.82, 1.34, 3.66, 1.00, accent=AMBER)
    callout(s, "局限 3：解释证据", "基序和片段结果是关联性解释，不能替代化学机制实验或毒理学验证。", 8.82, 1.34, 3.36, 1.00, accent=TEAL)
    add_three_line_table(
        s,
        [["审稿人可能问", "推荐回答方向", "证据页"],
         ["是否测试集调参？", "候选先登记，验证集选择，冻结后测试", "8"],
         ["是否普遍优于强基线？", "不是；报告选择性增益和保留策略", "9-12"],
         ["不可靠样本如何处理？", "risk-coverage、AD gate、保形覆盖和失败案例", "13-14"],
         ["解释是否因果？", "仅作为关联解释和假设生成", "16"],
         ["失败结果在哪里？", "主文摘要 + 补充长表完整保留", "15"]],
        0.92, 3.03, 11.22, 2.55, widths=[0.30, 0.48, 0.22], font_size=8.5, header_size=8.9, accent=RED,
    )
    takeaway(s, "最安全的 Nature 风格不是回避不足，而是把证据边界讲清楚。", accent=RED)
    footer(s, 18, total)
    note(s, "这一页可以在正式答辩中用于 Q&A 前。强调所有局限都是已知并被纳入论文表达，而不是事后补救。")

    # 19
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, TEAL)
    title(s, "结论：FZYC-Mol 的价值在于可审计的分子预测流程", "用三句话收束，不把结论写成过度承诺。", TEAL)
    callout(s, "结论 1", "FZYC-Mol 将多源表示、候选专家、验证选择器和风险输出整合为一个冻结测试流程，减少后验选择风险。", 0.88, 1.50, 3.52, 1.24, accent=TEAL)
    callout(s, "结论 2", "MoleculeNet 和 TDC 结果支持选择性增益与稳定保留策略；负结果和边界案例限制了过度外推。", 4.86, 1.50, 3.52, 1.24, accent=BLUE)
    callout(s, "结论 3", "risk-coverage、保形覆盖、OOD 压力测试和解释分析让预测结果更适合被审计和复核。", 8.84, 1.50, 3.18, 1.24, accent=GREEN)
    process_strip(
        s,
        [("可信", "冻结测试与验证集治理"),
         ("稳健", "多划分、外部任务、风险覆盖"),
         ("透明", "解释、失败案例、负结果"),
         ("可复用", "代码、数据、流程可追踪")],
        0.98, 3.40, 11.0, 1.05, accent=TEAL,
    )
    text(s, "建议最终论文题眼：不是“更复杂的分子预测模型”，而是“面向可靠性审计的验证集治理框架”。", 1.10, 5.30, 10.65, 0.45, size=12.2, color=DARK, bold=True, align="center")
    takeaway(s, "结论应回到主张与证据一致：有限、透明、可复核，因此更可信。")
    footer(s, 19, total)
    note(s, "结论页不要再加入新结果。只重复三件事：流程治理、选择性增益、可靠性审计。")

    # 20
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, BLUE)
    text(s, "Q&A", 0.68, 0.58, 2.3, 0.55, size=30, color=DARK, bold=True)
    text(s, "建议优先准备这四类问题", 0.72, 1.18, 4.2, 0.32, size=12.0, color=C["muted"])
    callout(s, "1. 验证集治理是否充分？", "答：说明候选登记、验证集选择、冻结测试和 rank audit；强调测试集不参与策略选择。", 0.86, 1.86, 5.20, 0.98, accent=BLUE)
    callout(s, "2. 为什么不是所有终点提升？", "答：方法本身会拒绝不稳定候选；论文主张是选择性增益和可靠性审计，不是普遍最优。", 6.50, 1.86, 5.20, 0.98, accent=TEAL)
    callout(s, "3. 如何处理低相似度与活性悬崖？", "答：报告 Tanimoto 分层、结构分离、risk enrichment 和代表性失败案例，不回避边界。", 0.86, 3.20, 5.20, 0.98, accent=RED)
    callout(s, "4. 解释结果是否可作为机制？", "答：不能。基序/片段解释仅为关联证据和假设生成，需要外部实验验证。", 6.50, 3.20, 5.20, 0.98, accent=AMBER)
    text(s, "汇报结束语：本研究的目标不是替代实验验证，而是让分子性质预测在模型选择、外推和可靠性报告上更可审计。", 1.15, 5.34, 10.65, 0.46, size=12.0, color=DARK, bold=True, align="center")
    takeaway(s, "答辩时保持保守语气：能证明的讲结果，不能证明的讲边界。", accent=BLUE)
    footer(s, 20, total)
    note(s, "最后一页用于 Q&A。建议用户把这四类问题背熟，因为它们对应审稿人最可能追问的逻辑漏洞。")

    for i, slide in enumerate(slides, 1):
        # Footers are already added; this loop keeps future edits easy.
        pass

    prs.save(OUT_PPTX)
    return prs


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    text_lengths = []
    shape_counts = []
    sw = prs.slide_width
    sh = prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        total_text = 0
        textboxes = 0
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > sw + 1000 or shp.top + shp.height > sh + 1000:
                issues.append(f"Slide {si}: shape outside slide bounds")
            if hasattr(shp, "text") and shp.text.strip():
                textboxes += 1
                total_text += len(shp.text.strip())
                for token in shp.text.split():
                    cjk = sum(1 for ch in token if "\u4e00" <= ch <= "\u9fff")
                    if len(token) > 42 and cjk < len(token) * 0.45:
                        issues.append(f"Slide {si}: long non-CJK token may overflow: {token[:50]}")
        if total_text < 50:
            issues.append(f"Slide {si}: low text density")
        if total_text > 760:
            issues.append(f"Slide {si}: high text density ({total_text} chars)")
        if textboxes > 45:
            issues.append(f"Slide {si}: too many text boxes ({textboxes})")
        text_lengths.append(total_text)
        shape_counts.append(len(slide.shapes))
    with ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        notes_xml = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")]
    return {
        "slides": len(prs.slides),
        "media": len(media),
        "notes": notes,
        "slide_xml": len(slide_xml),
        "note_xml": len(notes_xml),
        "min_text": min(text_lengths),
        "max_text": max(text_lengths),
        "min_shapes": min(shape_counts),
        "max_shapes": max(shape_counts),
        "issues": issues,
    }


def write_report(stats):
    lines = [
        "# 初稿-7 简化充实版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 改进方向：从 34 页压缩为 20 页；每页增加解释、证据、表格或讲者备注，减少过场页。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 发现以下低/中风险自动检查提示，已用于人工复核：")
        for item in stats["issues"][:20]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 本版保留原稿内嵌科学图像，未修改原始数据图。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页渲染预览；已完成 PPTX 包结构、媒体、备注、边界与文本密度检查。",
            "- 为满足“页数更简化、单页内容更多”的要求，本版采用 20 页中等密度结构；完整细节仍建议在讲者备注和论文中展开。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    src = base.extract_source()
    figs = base.extract_figures()
    if len(figs) < 13:
        raise RuntimeError(f"Expected at least 13 figures, got {len(figs)}")
    build_deck(src, figs)
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
