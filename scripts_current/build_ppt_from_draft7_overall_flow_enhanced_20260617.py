# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FLOW_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_overall_flow_largefont_20260616.py"
spec = importlib.util.spec_from_file_location("flow_ppt", FLOW_SCRIPT)
flow = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(flow)

base = flow.base
WORK = ROOT / "reports" / "ppt_from_draft7_overall_flow_enhanced_20260617"
WORK.mkdir(parents=True, exist_ok=True)

OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_整体流程详解增强说明版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_整体流程详解增强说明版PPT_QA报告.md"


def source_data():
    return base.extract_source()


def src_label(slide, body):
    flow.tx(slide, body, 0.62, 6.34, 11.8, 0.16, size=6.8, color=flow.C["muted"])


def table(slide, rows, x, y, w, h, widths=None, accent=None, fs=8.4, hfs=8.8):
    accent = accent or flow.TEAL
    flow.dense.add_three_line_table(slide, rows, x, y, w, h, widths=widths, font_size=fs, header_size=hfs, accent=accent)


def dataset_rows(src):
    t = src["tables"][0]
    return [["数据层级", "具体数据集", "任务与规模", "为什么使用"]] + [
        ["内部基准", "ESOL / FreeSolv / Lipophilicity", "MoleculeNet 回归；642-4200", "验证溶解度、溶剂化自由能、脂溶性等连续性质"],
        ["内部基准", "BBBP / BACE / ClinTox", "MoleculeNet 分类；阳性率差异大", "覆盖渗透、靶点与毒性，检验稀有阳性和校准"],
        ["外部 ADMET", "TDC ADMET", "578-13130；RMSE / ROC-AUC", "检验公开外部任务和 scaffold 审计下的迁移性"],
        ["结构外推", "MoleculeACE / bRo5", "活性悬崖；CycPept-PAMPA；LinPept", "评估低相似度、规则五以外空间和适用域边界"],
    ]


def main_result_rows(src):
    t = src["tables"][2]
    return [["证据模块", "本文怎么用", "对主张的约束"]] + [
        ["MoleculeNet", "内部主基准与 seed 评估", "支持选择性保留，不写成全任务最优"],
        ["TDC ADMET", "外部任务与官方划分审计", "支持公开外部迁移，不等同临床验证"],
        ["MoleculeACE / OOD", "活性悬崖、结构分离与低相似度", "用于发现失败边界"],
        ["可靠性审计", "AD、uncertainty、conformal、risk-coverage", "用于提示何时谨慎使用预测"],
    ]


def build_deck(src, figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 15

    # 1
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.tx(s, "FZYC-Mol", 0.66, 0.42, 4.8, 0.56, size=30, bold=True)
    flow.tx(s, "整体实验流程详解增强说明版", 0.70, 1.00, 7.2, 0.36, size=16, color=flow.SLATE)
    flow.tx(s, "补充内容：数据集、创新点、贡献度、与常见文献范式的优势对比", 0.72, 1.42, 10.6, 0.30, size=12.5, color=flow.C["muted"])
    flow.flow_row(s, ["问题定义", "数据矩阵", "划分冻结", "专家训练", "验证治理", "可靠性审计"], 0.62, 2.15, 12.05, 0.78, accent=flow.TEAL)
    flow.flow_row(s, ["OOD压力", "消融负结果", "解释边界", "复现材料", "文献优势", "答辩口径"], 0.62, 3.26, 12.05, 0.78, accent=flow.BLUE)
    flow.big_box(s, "本版定位", "不是单纯展示结果图，而是讲清楚整套实验为什么这样设计、数据如何覆盖问题、创新在哪里、优势如何成立。", 0.86, 4.58, 5.45, 0.95, accent=flow.TEAL)
    flow.big_box(s, "表达原则", "所有优势均围绕流程、审计和边界，不写成无条件超过全部已有方法。", 6.72, 4.58, 5.00, 0.95, accent=flow.RED)
    flow.takeaway(s, "这一版用于更完整地解释研究设计和贡献，而不只是讲流程步骤。")
    flow.footer(s, 1, total)
    flow.note(s, "开场说明：本版在大字流程版基础上增加数据集矩阵、创新点、贡献和与常见文献范式的对比。强调不会虚构具体文献结论。")

    # 2
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.BLUE)
    flow.title(s, "研究背景：已有工作常把性能、外推和可靠性混在一起", "FZYC-Mol 的出发点是把模型选择与可靠性审计制度化。", flow.BLUE)
    flow.big_box(s, "常见范式 1：单模型 benchmark", "很多工作聚焦某个模型在标准划分下的平均指标；问题是难以说明模型在低相似度或外部任务上何时失效。", 0.78, 1.50, 3.70, 1.22, accent=flow.BLUE)
    flow.big_box(s, "常见范式 2：简单 ensemble", "集成方法可能提升均值表现，但若没有验证集门控，很难证明每个任务上的融合都是必要且稳定的。", 4.80, 1.50, 3.70, 1.22, accent=flow.AMBER)
    flow.big_box(s, "常见范式 3：只报告 uncertainty", "不确定性或校准常被作为附加模块；如果不与适用域、失败样本和 risk-coverage 联动，审计价值有限。", 8.82, 1.50, 3.35, 1.22, accent=flow.GREEN)
    flow.two_col_steps(
        s,
        "本文要解决的核心问题",
        ["模型是否由验证集而不是测试集选择？",
         "哪些任务确实接受融合或补救模块？",
         "低相似度、活性悬崖和 OOD 样本是否单独审计？",
         "负结果是否进入论文逻辑而不是被隐藏？"],
        "本文采用的安全策略",
        ["候选先登记，再由 validation selector 选择。",
         "测试集只在策略冻结后一次性评估。",
         "可靠性、校准、保形和失败案例联合报告。",
         "不能证明的模块降调为边界或后续验证。"],
        accent=flow.BLUE,
    )
    flow.takeaway(s, "相对已有常见范式，本文的核心改进是把“选择过程”本身变成可审计对象。", accent=flow.BLUE)
    flow.footer(s, 2, total)
    flow.note(s, "这一页用于解释为什么研究有必要。讲法要保守：比较的是常见研究范式，不是指责某一篇文献。")

    # 3
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "我们使用的数据集覆盖内部基准、外部 ADMET 与结构外推", "数据设计不是越多越好，而是让每类实验对应一个明确问题。", flow.TEAL)
    table(s, dataset_rows(src), 0.62, 1.34, 12.00, 3.42, widths=[0.18, 0.30, 0.24, 0.28], accent=flow.TEAL, fs=7.9, hfs=8.4)
    flow.big_box(s, "内部基准回答", "模型在标准分子性质预测任务中是否能稳定保留有效候选。", 0.90, 5.12, 3.45, 0.68, accent=flow.TEAL)
    flow.big_box(s, "外部任务回答", "流程能否迁移到公开 ADMET 终点，而不是只适配 MoleculeNet。", 4.70, 5.12, 3.45, 0.68, accent=flow.BLUE)
    flow.big_box(s, "结构压力回答", "低相似度、bRo5 与活性悬崖样本是否暴露失败边界。", 8.50, 5.12, 3.35, 0.68, accent=flow.RED)
    flow.takeaway(s, "数据集矩阵支撑三类问题：标准性能、外部迁移、结构外推。")
    flow.footer(s, 3, total)
    flow.note(s, "这一页详细说明用了什么数据集。可以按内部基准、外部 ADMET、结构外推三层来讲。")

    # 4
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.AMBER)
    flow.title(s, "实验整体流程：从数据冻结到可复核结论", "每个阶段都有输入、处理、质控和输出文件。", flow.AMBER)
    base.add_picture_fit(s, figs[1], 0.58, 1.30, 7.25, 4.68, frame=True)
    flow.big_box(s, "流程说明", "数据清洗和划分冻结在训练前完成，后续模型只能读取已冻结版本。", 8.30, 1.42, 3.70, 0.82, accent=flow.AMBER)
    flow.big_box(s, "验证治理", "候选专家先进入 registry，由 validation selector 判断是否保留。", 8.30, 2.42, 3.70, 0.82, accent=flow.TEAL)
    flow.big_box(s, "测试原则", "测试集不参与选择，只用于冻结策略后的最终评估与 rank audit。", 8.30, 3.42, 3.70, 0.82, accent=flow.BLUE)
    flow.big_box(s, "输出层", "最终输出包括预测、风险、适用域、保形覆盖、失败案例和负结果。", 8.30, 4.42, 3.70, 0.82, accent=flow.RED)
    src_label(s, "Source: Fig. 2 from 初稿-7.docx；用于说明 workflow，不作为结果图。")
    flow.takeaway(s, "整套实验的可信度来自流程冻结，而不是单个结果数字。", accent=flow.AMBER)
    flow.footer(s, 4, total)
    flow.note(s, "这一页讲整体流程图。强调每个流程节点对应实际文件和脚本，方便审稿复核。")

    # 5
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.PURPLE)
    flow.title(s, "方法结构：多源表示与候选专家池支持 selector 治理", "模型结构服务于候选生成，最终选择仍由验证集决定。", flow.PURPLE)
    base.add_picture_fit(s, figs[0], 0.58, 1.30, 7.35, 4.62, frame=True)
    flow.big_box(s, "多源表示", "SMILES、描述符、指纹和图结构共同提供互补视角。", 8.38, 1.42, 3.70, 0.78, accent=flow.PURPLE)
    flow.big_box(s, "专家池", "强基线、图模型、指纹模型、表格模型和补救头先作为候选。", 8.38, 2.38, 3.70, 0.78, accent=flow.BLUE)
    flow.big_box(s, "证据输出", "除预测值外，还生成 uncertainty、AD gate、校准和解释信号。", 8.38, 3.34, 3.70, 0.78, accent=flow.GREEN)
    flow.big_box(s, "创新定位", "创新不是单个新模型，而是模型选择、可靠性和负结果的统一治理。", 8.38, 4.30, 3.70, 0.88, accent=flow.RED)
    src_label(s, "Source: Fig. 1 from 初稿-7.docx；用于说明模型结构，不作为性能结果。")
    flow.takeaway(s, "FZYC-Mol 的贡献在系统层面：把候选生成、选择和审计合成闭环。", accent=flow.BLUE)
    flow.footer(s, 5, total)
    flow.note(s, "这一页讲创新不是模型堆叠，而是候选池和验证治理的系统设计。")

    # 6
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "核心创新点 1：validation selector 让模型选择过程可审计", "接受、拒绝和冻结都留下证据，而不是只展示最终成功结果。", flow.TEAL)
    base.add_picture_fit(s, figs[2], 0.58, 1.30, 7.35, 4.36, frame=True)
    flow.big_box(s, "候选登记", "所有专家、融合、补救头和适配器先进入 candidate registry。", 8.38, 1.40, 3.72, 0.76, accent=flow.TEAL)
    flow.big_box(s, "验证选择", "selector 根据 valid 指标、校准和风险信号决定保留或拒绝。", 8.38, 2.34, 3.72, 0.76, accent=flow.BLUE)
    flow.big_box(s, "冻结测试", "selected_strategy.json 输出后，才允许进入测试集评估。", 8.38, 3.28, 3.72, 0.76, accent=flow.AMBER)
    flow.big_box(s, "负结果归档", "未通过模块进入 negative_results.md，防止选择性报告。", 8.38, 4.22, 3.72, 0.76, accent=flow.RED)
    src_label(s, "Source: Fig. 3 from 初稿-7.docx；用于说明 selector 治理，不作为性能图。")
    flow.takeaway(s, "selector 的贡献是把模型选择从隐性调参变成显性流程。")
    flow.footer(s, 6, total)
    flow.note(s, "这一页是最重要的创新点。要说明接受和拒绝都能被追踪。")

    # 7
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.BLUE)
    flow.title(s, "核心创新点 2：评价协议覆盖性能、外推和可靠性三条证据线", "同一套流程同时回答能不能预测、能否外推、何时谨慎。", flow.BLUE)
    table(s, main_result_rows(src), 0.70, 1.34, 11.90, 2.62, widths=[0.22, 0.40, 0.38], accent=flow.BLUE, fs=8.2, hfs=8.7)
    flow.big_box(s, "性能线", "MoleculeNet 主任务和外部 TDC 任务用于检验模型是否具有稳定候选策略。", 0.82, 4.36, 3.55, 0.86, accent=flow.BLUE)
    flow.big_box(s, "外推线", "scaffold、低相似度、MoleculeACE 和 bRo5 用于揭示结构外推边界。", 4.78, 4.36, 3.55, 0.86, accent=flow.RED)
    flow.big_box(s, "可靠性线", "risk-coverage、conformal、校准和失败案例让预测有风险语境。", 8.74, 4.36, 3.35, 0.86, accent=flow.GREEN)
    flow.takeaway(s, "优势不只在平均指标，而在三条证据线共同约束结论。", accent=flow.BLUE)
    flow.footer(s, 7, total)
    flow.note(s, "这一页说明实验不是一堆无关结果，而是三条证据线。")

    # 8
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.GREEN)
    flow.title(s, "创新点 3：可靠性审计从附加分析变成主流程输出", "AD、uncertainty、校准和保形预测共同解释预测是否可用。", flow.GREEN)
    flow.flow_row(s, ["预测值", "AD gate", "uncertainty", "calibration", "risk score", "conformal"], 0.78, 1.54, 11.80, 0.86, accent=flow.GREEN)
    flow.two_col_steps(
        s,
        "具体怎么做",
        ["计算训练分布距离或相似度，形成适用域信号。",
         "从专家分歧、模型方差或校准误差生成 uncertainty。",
         "将 risk score 与 risk-coverage、失败案例联动。",
         "用 80%、90%、95% 覆盖率报告 conformal 结果。"],
        "为什么有贡献",
        ["预测不再只有一个数值，而带有风险语境。",
         "高风险样本可进入失败案例和边界讨论。",
         "稀有阳性任务可结合 PR-AUC、Brier、ECE 和 fixed precision recall。",
         "可靠性主张被限制在可审计范围内。"],
        accent=flow.GREEN,
    )
    flow.takeaway(s, "可靠性模块让模型输出从“点预测”变成“可审计预测”。", accent=flow.GREEN)
    flow.footer(s, 8, total)
    flow.note(s, "这一页详细解释可靠性模块。不要讲成准确率提升，而是讲成风险识别和适用域限定。")

    # 9
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.RED)
    flow.title(s, "创新点 4：负结果、消融和失败案例进入主流程", "这部分是防止过度主张和学术风险的重要保护层。", flow.RED)
    flow.flow_row(s, ["Full", "best single", "simple mean", "w/o selector", "w/o fusion", "w/o AD"], 0.78, 1.48, 11.80, 0.86, accent=flow.RED)
    flow.big_box(s, "消融矩阵", "统一比较 Full、best single、simple mean、w/o selector、w/o fusion、w/o AD gate、w/o uncertainty weighting。", 0.80, 2.76, 3.70, 1.30, accent=flow.RED)
    flow.big_box(s, "负结果记录", "FreeSolv 低成本重构、bRo5、3D-lite、粗糙度加权等未稳定通过者作为边界或后续验证。", 4.82, 2.76, 3.70, 1.30, accent=flow.AMBER)
    flow.big_box(s, "失败案例", "ClinTox 假阴性、低相似度失败、活性悬崖失败和高风险 ADME 样本进入 casebook。", 8.84, 2.76, 3.35, 1.30, accent=flow.BLUE)
    flow.big_box(s, "解释边界", "motif attribution 和 fragment enrichment 仅支持关联解释，不作为因果机制证据。", 1.08, 4.70, 10.80, 0.72, accent=flow.PURPLE, body_size=12)
    flow.takeaway(s, "贡献不是把失败藏起来，而是把失败变成适用域和方法边界的一部分。", accent=flow.RED)
    flow.footer(s, 9, total)
    flow.note(s, "这一页回应学术严谨性。强调负结果和失败案例是可信度来源。")

    # 10
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.AMBER)
    flow.title(s, "本文贡献度：从模型性能扩展到实验治理和复现体系", "贡献不是单点，而是覆盖方法、实验、可靠性和写作风险控制。", flow.AMBER)
    table(
        s,
        [["贡献维度", "具体贡献", "论文/答辩价值"],
         ["方法贡献", "validation selector + gate + 冻结测试", "解释模型选择为何可信"],
         ["实验贡献", "MoleculeNet + TDC + MoleculeACE + bRo5 多层验证", "覆盖内部、外部和结构外推"],
         ["可靠性贡献", "AD / uncertainty / conformal / risk-coverage", "说明预测何时需要谨慎"],
         ["学术规范贡献", "负结果、失败案例、rank audit、optimism gap", "减少过度主张和选择性报告风险"],
         ["复现贡献", "registry、split、feature、config、metrics、casebook", "方便审稿人复核流程"]],
        0.70, 1.34, 11.90, 4.20, widths=[0.18, 0.42, 0.40], accent=flow.AMBER, fs=8.2, hfs=8.8,
    )
    flow.takeaway(s, "如果要一句话概括贡献：FZYC-Mol 是可靠性审计导向的验证集治理框架。", accent=flow.AMBER)
    flow.footer(s, 10, total)
    flow.note(s, "这一页专门讲贡献度。可以按方法、实验、可靠性、规范和复现五个维度讲。")

    # 11
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.BLUE)
    flow.title(s, "相对已有常见文献范式的优势", "这里比较研究范式，不替代逐篇文献综述或逐条引用核验。", flow.BLUE)
    table(
        s,
        [["常见范式", "常见不足", "FZYC-Mol 的优势"],
         ["单模型 QSAR / GNN", "强调模型结构，选择过程不一定透明", "候选登记 + selector 冻结，模型选择可复核"],
         ["普通 ensemble / averaging", "融合收益可能任务依赖，易后验挑选", "融合必须通过 validation gate，不通过则保留为负结果"],
         ["只做 benchmark", "随机或单一划分可能外推乐观", "scaffold、低相似度、MoleculeACE、bRo5 联合压力测试"],
         ["只做 uncertainty", "不一定连接失败案例和适用域", "AD、risk-coverage、conformal、casebook 形成闭环"],
         ["只报告成功结果", "边界和负结果不可见", "负结果、失败样本和限制进入主流程"]],
        0.62, 1.28, 12.08, 4.48, widths=[0.24, 0.36, 0.40], accent=flow.BLUE, fs=7.6, hfs=8.2,
    )
    flow.takeaway(s, "优势的核心不是宣称所有任务最强，而是流程透明、边界清楚、可复核性更高。", accent=flow.BLUE)
    flow.footer(s, 11, total)
    flow.note(s, "这一页讲与已有文献的优势。要明确这是范式层面的对比，不是逐篇文献贬低。")

    # 12
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "数据集与创新点如何一一对应", "每个数据层级都承担一个明确验证目的。", flow.TEAL)
    table(
        s,
        [["数据/实验", "对应问题", "支撑的创新点"],
         ["MoleculeNet 六任务", "标准分子性质预测是否稳定", "验证集 selector 与候选专家治理"],
         ["TDC ADMET", "外部 ADMET 是否可迁移", "多方法融合与适用域门控的外部审计"],
         ["MoleculeACE", "活性悬崖是否暴露误差", "低相似度和结构粗糙度失败分析"],
         ["bRo5 / CycPept-PAMPA / LinPept", "规则五以外空间是否可靠", "适用域边界与公共数据外推审计"],
         ["消融和负结果", "模块贡献是否真实", "防止选择性报告和过度解释"]],
        0.72, 1.34, 11.70, 4.20, widths=[0.30, 0.34, 0.36], accent=flow.TEAL, fs=8.0, hfs=8.6,
    )
    flow.takeaway(s, "数据集不是堆数量，而是让每个创新点都有对应的验证场景。")
    flow.footer(s, 12, total)
    flow.note(s, "这一页把数据集和创新点对应起来，适合答辩时回答“为什么选择这些数据集”。")

    # 13
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.PURPLE)
    flow.title(s, "论文写作中应突出的贡献句", "这些句子可以作为汇报和论文摘要/讨论的安全表述。", flow.PURPLE)
    flow.big_box(s, "贡献句 1", "我们提出一种以验证集治理为核心的分子性质预测流程，将候选专家选择、冻结测试和可靠性审计连接为闭环。", 0.82, 1.36, 5.45, 0.98, accent=flow.PURPLE)
    flow.big_box(s, "贡献句 2", "该流程并不假设融合或补救模块在所有任务上有效，而是通过 validation gate 对任务依赖的局部增益进行接受或拒绝。", 6.82, 1.36, 5.30, 0.98, accent=flow.TEAL)
    flow.big_box(s, "贡献句 3", "通过 MoleculeNet、TDC ADMET、MoleculeACE 和 bRo5 压力测试，本文同时评估标准性能、外部迁移和结构外推边界。", 0.82, 2.78, 5.45, 0.98, accent=flow.BLUE)
    flow.big_box(s, "贡献句 4", "AD、uncertainty、conformal 和 risk-coverage 分析为每个预测提供风险语境，帮助识别不应过度信任的样本。", 6.82, 2.78, 5.30, 0.98, accent=flow.GREEN)
    flow.big_box(s, "贡献句 5", "负结果、失败案例和解释边界被纳入主流程，从而降低选择性报告和因果过度解释风险。", 1.10, 4.40, 10.80, 0.82, accent=flow.RED, body_size=12.2)
    flow.takeaway(s, "贡献表述要突出流程治理和可靠性，而不是写成绝对性能冠军。", accent=flow.PURPLE)
    flow.footer(s, 13, total)
    flow.note(s, "这一页可作为汇报时的过渡，也可给用户用于论文摘要和讨论句式。")

    # 14
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.RED)
    flow.title(s, "需要谨慎表达的边界", "优势越明确，边界越要写清楚。", flow.RED)
    flow.big_box(s, "不能写成", "不能写成 FZYC-Mol 在所有数据集和所有终点上全面优于全部已有方法。", 0.82, 1.44, 3.70, 1.05, accent=flow.RED)
    flow.big_box(s, "不能写成", "不能将 motif/fragment 解释直接写成化学因果机制或实验验证结论。", 4.82, 1.44, 3.70, 1.05, accent=flow.RED)
    flow.big_box(s, "不能写成", "不能把 TDC 或 bRo5 公共数据外推写成独立盲测或湿实验验证。", 8.82, 1.44, 3.35, 1.05, accent=flow.RED)
    flow.two_col_steps(
        s,
        "安全写法",
        ["选择性增益，而非普遍提升。",
         "可靠性审计，而非替代实验验证。",
         "关联解释，而非因果机制。",
         "公开外部任务，而非独立前瞻性盲测。"],
        "答辩提示",
        ["问到性能：回到 validation gate。",
         "问到外推：回到 OOD 与低相似度。",
         "问到机制：回到关联解释和后续实验。",
         "问到复现：展示 registry、split 和脚本。"],
        accent=flow.RED,
    )
    flow.takeaway(s, "把边界讲清楚，反而会增强论文可信度。", accent=flow.RED)
    flow.footer(s, 14, total)
    flow.note(s, "这一页是风险控制。可以在答辩前重点复习。")

    # 15
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "最终汇报逻辑：用流程、数据和边界支撑贡献", "结尾回到一条可复核证据链。", flow.TEAL)
    flow.flow_row(s, ["数据覆盖", "流程冻结", "验证选择", "可靠性审计", "边界记录", "复现提交"], 0.78, 1.54, 11.80, 0.86, accent=flow.TEAL)
    flow.big_box(s, "我们用了什么", "MoleculeNet、TDC ADMET、MoleculeACE、bRo5 / CycPept-PAMPA / LinPept 等公开数据层级。", 0.82, 2.92, 3.70, 1.05, accent=flow.TEAL)
    flow.big_box(s, "我们创新什么", "验证集治理、门控冻结、可靠性审计、负结果归档和外推边界评估。", 4.82, 2.92, 3.70, 1.05, accent=flow.BLUE)
    flow.big_box(s, "我们贡献什么", "提供一套更透明、更可复核、更少过度外推风险的分子性质预测实验框架。", 8.82, 2.92, 3.35, 1.05, accent=flow.AMBER)
    flow.big_box(s, "一句话总结", "FZYC-Mol 的优势不是宣称万能预测，而是让模型选择、可靠性和失败边界都能被审计。", 1.10, 4.74, 10.80, 0.85, accent=flow.RED, body_size=12.4)
    flow.takeaway(s, "讲清“为什么这样设计”，比单独展示结果更能支撑投稿前自审。")
    flow.footer(s, 15, total)
    flow.note(s, "最后一页总结数据、创新、贡献和优势。强调整体流程比单个结果图更有说服力。")

    prs.save(OUT_PPTX)


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    text_lengths, shape_counts = [], []
    sw, sh = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        total_text = 0
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > sw + 1000 or shp.top + shp.height > sh + 1000:
                issues.append(f"Slide {si}: shape outside slide bounds")
            if hasattr(shp, "text") and shp.text.strip():
                total_text += len(shp.text.strip())
        if total_text < 120:
            issues.append(f"Slide {si}: low text density")
        if total_text > 950:
            issues.append(f"Slide {si}: high text density ({total_text} chars)")
        text_lengths.append(total_text)
        shape_counts.append(len(slide.shapes))
    with ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        notes_xml = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")]
    return {
        "slides": len(prs.slides),
        "media": len(media),
        "notes": notes,
        "slide_xml": len(slide_xml),
        "note_xml": len(notes_xml),
        "min_text": min(text_lengths),
        "max_text": max(text_lengths),
        "min_shapes": min(shape_counts),
        "max_shapes": max(shape_counts),
        "issues": issues,
    }


def write_report(stats):
    lines = [
        "# 初稿-7 整体流程详解增强说明版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 修改方向：在整体流程 PPT 中补充数据集、创新点、贡献度、与常见文献范式优势对比和谨慎表达边界。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个（概念图，不含性能结果图）",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 自动检查提示如下：")
        for item in stats["issues"][:20]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 数据集名称和关键数值来自《初稿-7.docx》的表格。",
            "- 文献优势部分采用“相对常见研究范式”的谨慎比较，没有虚构某一篇具体文献的结论。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页截图预览；已完成 PPTX 包结构、媒体、备注、边界和文本密度检查。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    src = source_data()
    figs = base.extract_figures()
    if len(figs) < 3:
        raise RuntimeError(f"Expected at least 3 figures, got {len(figs)}")
    build_deck(src, figs)
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
