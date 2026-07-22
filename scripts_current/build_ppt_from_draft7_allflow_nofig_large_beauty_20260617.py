# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BASE_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_20260616.py"
spec = importlib.util.spec_from_file_location("base_ppt", BASE_SCRIPT)
base = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(base)

WORK = ROOT / "reports" / "ppt_from_draft7_allflow_nofig_large_beauty_20260617"
WORK.mkdir(parents=True, exist_ok=True)

OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_全流程详解无图大字美化版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_全流程详解无图大字美化版PPT_QA报告.md"


BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(19, 31, 45)
MUTED = RGBColor(82, 97, 113)
LINE = RGBColor(203, 212, 222)
BLUE = RGBColor(34, 91, 151)
TEAL = RGBColor(26, 124, 118)
AMBER = RGBColor(188, 131, 37)
RED = RGBColor(177, 70, 68)
GREEN = RGBColor(62, 143, 94)
PURPLE = RGBColor(96, 83, 156)
PALE_BLUE = RGBColor(228, 239, 249)
PALE_TEAL = RGBColor(224, 244, 241)
PALE_AMBER = RGBColor(250, 239, 216)
PALE_RED = RGBColor(250, 228, 228)
PALE_PURPLE = RGBColor(239, 237, 249)


def add_bg(slide, accent=TEAL):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.13))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(7.16), Inches(12.30), Inches(0.02))
    foot.fill.solid()
    foot.fill.fore_color.rgb = LINE
    foot.line.fill.background()


def txt(slide, body, x, y, w, h, size=13.0, color=DARK, bold=False, align="left", valign="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "mid": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def title(slide, main, sub="", accent=TEAL):
    txt(slide, main, 0.58, 0.32, 11.9, 0.52, size=26.0, bold=True)
    if sub:
        txt(slide, sub, 0.60, 0.88, 11.4, 0.30, size=12.8, color=MUTED)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(1.23), Inches(1.30), Inches(0.045))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()


def footer(slide, idx, total):
    txt(slide, "FZYC-Mol | 全流程详解无图大字版", 0.56, 7.20, 4.60, 0.18, size=9.0, color=MUTED)
    txt(slide, f"{idx:02d}/{total}", 12.04, 7.20, 0.80, 0.18, size=9.0, color=MUTED, align="right")


def box(slide, head, body, x, y, w, h, accent=TEAL, fill=WHITE, head_size=13.8, body_size=12.8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    shp.line.width = Pt(0.8)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    txt(slide, head, x + 0.20, y + 0.10, w - 0.36, 0.28, size=head_size, color=accent, bold=True)
    txt(slide, body, x + 0.20, y + 0.48, w - 0.36, h - 0.56, size=body_size, color=DARK)
    return shp


def node(slide, label, x, y, w, h, accent=TEAL, fill=WHITE, size=12.8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = accent
    shp.line.width = Pt(1.1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = DARK
    return shp


def arrow(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.55)
    line.line.end_arrowhead = True
    return line


def flow_row(slide, items, x, y, w, h, accent=TEAL, size=12.0):
    gap = 0.13
    item_w = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        nx = x + i * (item_w + gap)
        node(slide, item, nx, y, item_w, h, accent=accent, size=size)
        if i < len(items) - 1:
            arrow(slide, nx + item_w + 0.02, y + h / 2, nx + item_w + gap - 0.04, y + h / 2, color=accent)


def takeaway(slide, body, accent=TEAL):
    color = PALE_TEAL if accent == TEAL else PALE_BLUE if accent == BLUE else PALE_AMBER if accent == AMBER else PALE_RED if accent == RED else PALE_PURPLE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(6.54), Inches(12.20), Inches(0.44))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    txt(slide, "Takeaway  " + body, 0.75, 6.65, 11.75, 0.22, size=11.7, bold=True)


def step_columns(slide, left_title, left_items, right_title, right_items, accent=TEAL):
    box(slide, left_title, "\n".join(left_items), 0.78, 3.04, 5.65, 2.72, accent=accent, body_size=12.5)
    other = BLUE if accent != BLUE else TEAL
    box(slide, right_title, "\n".join(right_items), 6.90, 3.04, 5.55, 2.72, accent=other, body_size=12.5)


def table(slide, rows, x, y, w, h, widths=None, accent=TEAL, font_size=10.0, header_size=10.5):
    if not rows:
        return
    cols = len(rows[0])
    widths = widths or [1 / cols] * cols
    total = sum(widths)
    widths = [ww / total * w for ww in widths]
    row_h = h / len(rows)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.fill.background()
    for yy, width in [(y, 1.05), (y + row_h, 0.75), (y + h, 1.05)]:
        line = slide.shapes.add_connector(1, Inches(x), Inches(yy), Inches(x + w), Inches(yy))
        line.line.color.rgb = DARK
        line.line.width = Pt(width)
    cx = x
    for j in range(cols):
        for i, row in enumerate(rows):
            body = row[j] if j < len(row) else ""
            txt(slide, body, cx + 0.04, y + i * row_h + 0.04, widths[j] - 0.08, row_h - 0.06,
                size=header_size if i == 0 else font_size,
                color=accent if i == 0 else DARK,
                bold=i == 0)
        cx += widths[j]


def note(slide, body):
    base.add_notes(slide, body)


def data_summary_rows():
    return [
        ["层级", "数据集", "目的"],
        ["内部", "ESOL / FreeSolv / Lipo", "回归性质与边界"],
        ["内部", "BBBP / BACE / ClinTox", "分类、毒性与稀有阳性"],
        ["外部", "TDC ADMET", "公开外部迁移审计"],
        ["压力", "MoleculeACE / bRo5", "活性悬崖与规则外空间"],
    ]


def build_deck(src):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 18

    # 1
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    txt(s, "FZYC-Mol", 0.66, 0.44, 4.7, 0.56, size=31, bold=True)
    txt(s, "全流程详解无图大字美化版", 0.70, 1.04, 6.8, 0.38, size=17, color=MUTED)
    txt(s, "从研究问题、数据集、实验执行到创新贡献和投稿答辩口径", 0.72, 1.46, 10.8, 0.30, size=13.2, color=MUTED)
    flow_row(s, ["问题", "数据", "划分", "表示", "训练", "选择"], 0.72, 2.20, 11.85, 0.80, accent=TEAL)
    flow_row(s, ["测试", "可靠性", "OOD", "消融", "复现", "投稿"], 0.72, 3.36, 11.85, 0.80, accent=BLUE)
    box(s, "本版原则", "不插入任何图片或结果图；只用大字号流程图、步骤框和三线表风格说明完整实验链。", 0.92, 4.76, 5.40, 0.90, accent=TEAL)
    box(s, "汇报目标", "让听众听懂：用了什么数据、每个实验怎么跑、创新在哪里、贡献和边界是什么。", 6.82, 4.76, 5.10, 0.90, accent=AMBER)
    takeaway(s, "这版适合讲实验设计和整体流程，而不是展示性能结果。")
    footer(s, 1, total)
    note(s, "开场说明：本版本完全无图，强调流程、数据和贡献。每页字体更大，用阶段化讲解。")

    # 2
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "整体逻辑：先固定流程，再讨论结果是否可信", "FZYC-Mol 的主线是从隐性模型选择转向显性验证治理。", BLUE)
    flow_row(s, ["立题", "数据注册", "划分冻结", "候选登记", "验证选择", "冻结测试"], 0.78, 1.54, 11.78, 0.86, accent=BLUE)
    step_columns(
        s,
        "为什么这样设计",
        ["分子性质预测容易受到划分方式、候选数量和后验选择影响。",
         "单一测试指标不能回答模型何时可靠、何时失效。",
         "因此先让数据、划分、候选和指标固定，再进入测试。"],
        "流程最终输出",
        ["最终预测与主指标。",
         "AD gate、uncertainty、calibration 和 conformal 结果。",
         "OOD、低相似度、失败案例、消融和负结果记录。"],
        accent=BLUE,
    )
    takeaway(s, "整套流程的可信度来自“先冻结、后评估、再审计”。", accent=BLUE)
    footer(s, 2, total)
    note(s, "这一页讲总逻辑。重点说不是先看结果再改方法，而是先固定流程。")

    # 3
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "阶段 0：定义研究问题、创新点和允许主张", "先决定论文要证明什么，也决定哪些话不能说过头。", AMBER)
    flow_row(s, ["需求", "瓶颈", "假设", "证据", "边界"], 1.00, 1.56, 11.20, 0.88, accent=AMBER)
    box(s, "研究问题", "分子性质预测不仅要提高平均性能，还要说明模型选择是否透明、外推是否可靠、失败样本是否被记录。", 0.82, 3.00, 3.60, 1.18, accent=AMBER)
    box(s, "创新点", "validation selector、冻结测试、适用域门控、不确定性审计、保形覆盖、负结果归档。", 4.82, 3.00, 3.60, 1.18, accent=TEAL)
    box(s, "安全主张", "本文主张选择性增益和可靠性审计；不写成所有终点全面优于全部已有方法。", 8.82, 3.00, 3.38, 1.18, accent=RED)
    box(s, "阶段输出", "实验登记表、候选模块清单、评价指标清单、风险边界清单。", 1.05, 4.82, 11.05, 0.70, accent=BLUE, body_size=12.8)
    takeaway(s, "流程第一步是固定主张和证据标准，而不是马上跑模型。", accent=AMBER)
    footer(s, 3, total)
    note(s, "这一页用于讲立题和创新点。强调安全主张。")

    # 4
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "我们使用的数据集覆盖四类验证场景", "数据设计服务于问题：内部基准、外部迁移、结构外推和边界压力。", TEAL)
    table(s, data_summary_rows(), 0.82, 1.36, 11.55, 3.55, widths=[0.16, 0.38, 0.46], accent=TEAL, font_size=10.6, header_size=11.2)
    box(s, "数据贡献", "不是简单堆数据集，而是让每个数据层级都对应一个明确验证目的。", 0.98, 5.26, 3.35, 0.68, accent=TEAL)
    box(s, "评价范围", "回归、分类、外部 ADMET、活性悬崖和规则五以外化学空间同时覆盖。", 4.72, 5.26, 3.55, 0.68, accent=BLUE)
    box(s, "谨慎边界", "bRo5 与公共外部数据只能支持公开数据外推，不等同独立盲测。", 8.66, 5.26, 3.35, 0.68, accent=RED)
    takeaway(s, "数据集矩阵回答“为什么这些实验是必要的”。")
    footer(s, 4, total)
    note(s, "这一页详细讲数据集：MoleculeNet 六任务、TDC ADMET、MoleculeACE、bRo5。")

    # 5
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "阶段 1：数据登记与清洗", "所有数据先进入统一 registry，后续实验只引用冻结版本。", TEAL)
    flow_row(s, ["导入", "字段检查", "SMILES标准化", "重复处理", "标签审计", "版本冻结"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    step_columns(
        s,
        "执行流程",
        ["记录来源、任务类型、样本量、主指标和阳性率。",
         "检查 SMILES、标签、单位、缺失值和冲突标签。",
         "处理无效 SMILES、重复分子和标签矛盾样本。"],
        "输出文件",
        ["cleaned.csv：清洗后数据。",
         "dataset_registry.csv：版本、任务和指标登记。",
         "label_audit.md：剔除与修正理由。"],
        accent=TEAL,
    )
    takeaway(s, "数据登记保证每个结果都能追溯到同一清洗版本。")
    footer(s, 5, total)
    note(s, "讲清数据清洗。强调所有删除和修正都要有记录。")

    # 6
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "阶段 2：划分冻结与泄漏控制", "划分决定外推结论是否成立，必须在训练前固定。", AMBER)
    flow_row(s, ["registry", "seed", "random", "scaffold", "结构分离", "低相似度"], 0.78, 1.54, 11.80, 0.86, accent=AMBER)
    step_columns(
        s,
        "执行流程",
        ["生成 train / valid / test 三部分，并固定随机种子。",
         "MoleculeNet 同时保留 random、scaffold、结构分离和低相似度压力划分。",
         "低相似度按 Tanimoto >0.7、0.5-0.7、<0.5 建立互斥分层。"],
        "质控输出",
        ["split_index.json：固定划分。",
         "seed_list.txt：全部种子。",
         "scaffold_report.md：分布与潜在泄漏检查。",
         "不同划分的结论不混写。"],
        accent=AMBER,
    )
    takeaway(s, "划分冻结防止根据结果反向调整评价协议。", accent=AMBER)
    footer(s, 6, total)
    note(s, "强调测试集不能参与调参，划分一旦固定不能因结果不好而变动。")

    # 7
    s = prs.slides.add_slide(blank); add_bg(s, PURPLE)
    title(s, "阶段 3：分子表示与特征缓存", "同一分子生成多视图表示，供不同候选专家读取。", PURPLE)
    flow_row(s, ["SMILES", "描述符", "指纹", "图结构", "任务信息", "缓存"], 0.78, 1.54, 11.80, 0.86, accent=PURPLE)
    step_columns(
        s,
        "执行流程",
        ["从冻结 cleaned.csv 读取分子，不从临时表重新取数。",
         "生成 RDKit 描述符、Morgan/多指纹、图结构和任务上下文。",
         "记录每类特征的维度、失败样本、软件版本和生成时间。"],
        "输出文件",
        ["descriptor.parquet：描述符矩阵。",
         "fingerprint.npy：指纹矩阵。",
         "graph_cache.pt：图结构缓存。",
         "feature_manifest.md：特征版本说明。"],
        accent=PURPLE,
    )
    takeaway(s, "特征缓存让所有候选专家读取同一输入版本，减少不可复现差异。", accent=BLUE)
    footer(s, 7, total)
    note(s, "讲多源表示。注意不要说多视图必然提升，只说提供候选专家输入。")

    # 8
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "阶段 4：候选专家训练", "所有模型先作为候选登记，不能直接进入最终结论。", BLUE)
    flow_row(s, ["读取split", "读取特征", "训练专家", "valid预测", "日志保存", "候选登记"], 0.78, 1.54, 11.80, 0.86, accent=BLUE)
    step_columns(
        s,
        "候选专家",
        ["强基线模型：传统指纹、表格模型和已知强模型。",
         "深度模型：图模型或消息传递模型。",
         "扩展候选：fusion、rescue head、adapter、3D-lite 等。"],
        "输出文件",
        ["candidate_registry.csv：候选列表和配置。",
         "valid_predictions：验证集预测。",
         "test_predictions：仅供冻结后评估。",
         "training_log：失败训练也要记录。"],
        accent=BLUE,
    )
    takeaway(s, "训练阶段只产生候选和验证证据，最终策略由 selector 决定。", accent=BLUE)
    footer(s, 8, total)
    note(s, "解释专家池，不展示结果。强调候选越多越需要验证治理。")

    # 9
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "阶段 5：validation selector 与门控", "selector 是中枢：决定保留、拒绝或转入负结果。", TEAL)
    flow_row(s, ["候选登记", "valid指标", "校准风险", "门控判断", "策略冻结", "拒绝归档"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    box(s, "接受规则", "验证性能改善，并且校准、风险或适用域信号没有明显恶化。", 0.82, 2.92, 3.60, 1.02, accent=TEAL)
    box(s, "拒绝规则", "未通过门控的候选不进入主结论，只保留为负结果或后续验证接口。", 4.82, 2.92, 3.60, 1.02, accent=RED)
    box(s, "冻结规则", "输出 selected_strategy.json 后，才允许进入测试集评估。", 8.82, 2.92, 3.38, 1.02, accent=BLUE)
    box(s, "阶段输出", "selected_strategy.json、rejected_candidates.csv、selector_audit.md。", 1.05, 4.72, 11.05, 0.70, accent=AMBER, body_size=12.8)
    takeaway(s, "selector 把模型选择从隐性调参变成显性审计流程。")
    footer(s, 9, total)
    note(s, "这是核心页。强调接受和拒绝都要记录。")

    # 10
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "阶段 6：冻结测试与排名审计", "测试集只在策略冻结后使用，用于一次性检验。", AMBER)
    flow_row(s, ["冻结策略", "test预测", "主指标", "rank audit", "optimism gap", "主表锁定"], 0.78, 1.54, 11.80, 0.86, accent=AMBER)
    step_columns(
        s,
        "执行流程",
        ["读取 selected_strategy.json，不临时替换候选。",
         "在 test split 上一次性计算 RMSE、ROC-AUC、PR-AUC、Brier、ECE 等。",
         "比较 valid rank 与 test rank，检查验证最优是否等于测试最优。"],
        "输出文件",
        ["test_metrics.csv：冻结测试指标。",
         "rank_audit.csv：验证-测试排名审计。",
         "optimism_gap.md：乐观偏差分析。",
         "main_table_source.xlsx：主表来源。"],
        accent=AMBER,
    )
    takeaway(s, "冻结测试回答“有没有测试集调参”的关键质疑。", accent=AMBER)
    footer(s, 10, total)
    note(s, "讲冻结测试，只讲流程和输出文件，不展示数值图。")

    # 11
    s = prs.slides.add_slide(blank); add_bg(s, GREEN)
    title(s, "阶段 7：可靠性审计", "性能之外，还要说明模型何时可靠、何时需要谨慎。", GREEN)
    flow_row(s, ["预测值", "AD距离", "uncertainty", "calibration", "risk score", "conformal"], 0.78, 1.54, 11.80, 0.86, accent=GREEN)
    step_columns(
        s,
        "执行流程",
        ["基于训练分布计算适用域距离或相似度信号。",
         "从模型方差、专家分歧或校准误差生成 uncertainty。",
         "构建 risk score，并生成 risk-coverage 分析文件。",
         "按 80%、90%、95% 目标覆盖率报告 conformal。"],
        "输出文件",
        ["ad_score.csv：适用域分数。",
         "uncertainty.csv：不确定性分数。",
         "risk_coverage.csv：风险覆盖分析。",
         "conformal_report.md：覆盖率报告。"],
        accent=GREEN,
    )
    takeaway(s, "可靠性模块把点预测扩展为“预测值 + 风险 + 覆盖率”。", accent=GREEN)
    footer(s, 11, total)
    note(s, "讲可靠性，不要把风险分数说成性能提升。")

    # 12
    s = prs.slides.add_slide(blank); add_bg(s, RED)
    title(s, "阶段 8：OOD、低相似度与失败案例", "外推风险进入统一错误分析链，而不是被放在边角。", RED)
    flow_row(s, ["scaffold", "结构分离", "Tanimoto bin", "MoleculeACE", "失败样本", "风险富集"], 0.78, 1.54, 11.80, 0.86, accent=RED)
    step_columns(
        s,
        "执行流程",
        ["比较 random、scaffold、structure-separated 的性能压力。",
         "按 Tanimoto 分层输出低相似度性能、校准和风险富集。",
         "对 MoleculeACE 活性悬崖记录 cliff subset 与代表性失败对。",
         "将 ClinTox 假阴性、低相似度失败、高风险 ADME 归档。"],
        "输出文件",
        ["ood_metrics.csv：外推压力结果。",
         "tanimoto_bins.csv：互斥相似度分层。",
         "cliff_pairs.csv：活性悬崖对。",
         "failure_casebook.md：失败案例集。"],
        accent=RED,
    )
    takeaway(s, "失败分析不是削弱论文，而是限定模型适用域。", accent=RED)
    footer(s, 12, total)
    note(s, "讲低相似度、活性悬崖和失败案例。")

    # 13
    s = prs.slides.add_slide(blank); add_bg(s, PURPLE)
    title(s, "阶段 9：消融、负结果与解释边界", "模块贡献和失败候选都必须进入流程记录。", PURPLE)
    flow_row(s, ["Full", "best single", "simple mean", "w/o selector", "w/o fusion", "w/o AD"], 0.78, 1.54, 11.80, 0.86, accent=PURPLE)
    step_columns(
        s,
        "消融与负结果",
        ["统一矩阵：Full、best single、simple mean、w/o selector、w/o fusion、w/o uncertainty。",
         "FreeSolv、bRo5、3D-lite、粗糙度加权等不稳定候选写成边界或后续验证。",
         "未通过验证门控的模块不得写入主性能结论。"],
        "解释边界",
        ["motif attribution 和 fragment enrichment 只支持关联解释。",
         "需要报告最小支持度、效应量、p 值或 FDR。",
         "不能写成化学因果机制或湿实验验证。"],
        accent=PURPLE,
    )
    takeaway(s, "负结果和解释边界是学术安全层，不是可有可无的补充。", accent=PURPLE)
    footer(s, 13, total)
    note(s, "讲消融、负结果和解释边界。")

    # 14
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "创新点与贡献度", "贡献集中在流程治理、可靠性审计和复现规范。", BLUE)
    table(
        s,
        [["维度", "具体贡献", "价值"],
         ["方法", "validation selector + gate", "模型选择可审计"],
         ["实验", "内部/外部/OOD 多层验证", "覆盖标准性能和外推边界"],
         ["可靠性", "AD / uncertainty / conformal", "预测带风险语境"],
         ["规范", "负结果 + 失败案例 + rank audit", "降低选择性报告风险"],
         ["复现", "registry / split / config / casebook", "方便审稿复核"]],
        0.74, 1.34, 11.85, 4.30, widths=[0.16, 0.42, 0.42], accent=BLUE, font_size=9.5, header_size=10.2,
    )
    takeaway(s, "一句话贡献：FZYC-Mol 是面向可靠性审计的验证集治理框架。", accent=BLUE)
    footer(s, 14, total)
    note(s, "这一页讲创新点和贡献度，适合正式汇报时重点停留。")

    # 15
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "相对常见文献范式的优势", "这里比较研究范式，不替代逐篇文献综述。", AMBER)
    table(
        s,
        [["常见范式", "常见不足", "本文优势"],
         ["单模型 QSAR/GNN", "选择过程不一定透明", "候选登记与 selector 冻结"],
         ["普通 ensemble", "融合收益可能后验挑选", "必须通过 validation gate"],
         ["只做 benchmark", "随机划分可能外推乐观", "scaffold / 低相似度 / cliff 压力测试"],
         ["只做 uncertainty", "未必连接失败案例", "risk、AD、conformal、casebook 闭环"],
         ["只报告成功结果", "边界不可见", "负结果和失败样本进入主流程"]],
        0.62, 1.30, 12.08, 4.42, widths=[0.25, 0.35, 0.40], accent=AMBER, font_size=9.3, header_size=10.0,
    )
    takeaway(s, "优势不是绝对性能冠军，而是流程透明、边界清楚、复核性更强。", accent=AMBER)
    footer(s, 15, total)
    note(s, "讲文献范式优势，谨慎表达，不要逐篇贬低。")

    # 16
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "论文输出与复现包", "把流程产物整理成可投稿、可审稿、可答辩的证据链。", TEAL)
    flow_row(s, ["方法流程图", "主表", "补充长表", "数据声明", "代码声明", "答辩材料"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    step_columns(
        s,
        "论文中怎么写",
        ["方法部分按数据、划分、表示、专家、selector、冻结测试、可靠性审计顺序写。",
         "结果部分按内部基准、外部验证、可靠性、OOD、消融和负结果组织。",
         "讨论部分明确边界：公开数据外推、关联解释、未完成盲测。"],
        "复现包包含",
        ["dataset_registry.csv、split_index.json、feature_manifest.md。",
         "model_config.yaml、candidate_registry.csv、selected_strategy.json。",
         "metrics_source.xlsx、failure_casebook.md、negative_results.md。"],
        accent=TEAL,
    )
    takeaway(s, "复现材料让每个主张都能回指到一个流程产物。")
    footer(s, 16, total)
    note(s, "讲论文输出和复现包。")

    # 17
    s = prs.slides.add_slide(blank); add_bg(s, RED)
    title(s, "投稿与答辩中的谨慎表达", "越强调创新，越要清楚说明不能证明什么。", RED)
    box(s, "不能写成", "不能写成在所有终点上全面超过全部已有模型。", 0.82, 1.50, 3.55, 0.88, accent=RED)
    box(s, "不能写成", "不能把 motif/fragment 解释说成因果机制证明。", 4.82, 1.50, 3.55, 0.88, accent=RED)
    box(s, "不能写成", "不能把公开外部数据写成独立盲测或湿实验验证。", 8.82, 1.50, 3.35, 0.88, accent=RED)
    step_columns(
        s,
        "推荐口径",
        ["选择性增益，而非普遍提升。",
         "可靠性审计，而非替代实验验证。",
         "关联解释，而非因果机制。",
         "公开外部任务，而非前瞻性盲测。"],
        "答辩问题",
        ["是否测试集调参：回答冻结流程。",
         "是否优于全部基线：回答选择性保留。",
         "外推是否可靠：回答 OOD 和低相似度。",
         "解释是否因果：回答关联证据。"],
        accent=RED,
    )
    takeaway(s, "清楚讲边界，会让贡献更可信。", accent=RED)
    footer(s, 17, total)
    note(s, "这一页用于答辩准备。")

    # 18
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "最终汇报口径：全流程支撑可信结论", "用一条证据链收束：数据、流程、选择、审计、边界。", TEAL)
    flow_row(s, ["数据覆盖", "流程冻结", "验证选择", "一次测试", "可靠性审计", "复现提交"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    box(s, "我们用了什么", "MoleculeNet、TDC ADMET、MoleculeACE、bRo5 / CycPept-PAMPA / LinPept 等公开数据层级。", 0.82, 3.00, 3.60, 1.02, accent=TEAL)
    box(s, "我们创新什么", "验证集治理、门控冻结、可靠性审计、负结果归档和外推边界评估。", 4.82, 3.00, 3.60, 1.02, accent=BLUE)
    box(s, "我们贡献什么", "提供一套更透明、更可复核、更少过度外推风险的分子性质预测实验框架。", 8.82, 3.00, 3.38, 1.02, accent=AMBER)
    box(s, "一句话总结", "FZYC-Mol 的优势不是宣称万能预测，而是让模型选择、可靠性和失败边界都能被审计。", 1.05, 4.74, 11.05, 0.74, accent=RED, body_size=13.0)
    takeaway(s, "讲完整流程，才能让听众相信每个结果都有来源、每个边界都被记录。")
    footer(s, 18, total)
    note(s, "最终总结：数据、创新、贡献、优势和边界。")

    prs.save(OUT_PPTX)


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    text_lengths = []
    shape_counts = []
    sw, sh = prs.slide_width, prs.slide_height
    min_font = 99
    for si, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        total_text = 0
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > sw + 1000 or shp.top + shp.height > sh + 1000:
                issues.append(f"Slide {si}: shape outside slide bounds")
            if hasattr(shp, "text") and shp.text.strip():
                total_text += len(shp.text.strip())
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            min_font = min(min_font, r.font.size.pt)
        if total_text < 110:
            issues.append(f"Slide {si}: low text density")
        if total_text > 780:
            issues.append(f"Slide {si}: high text density ({total_text} chars)")
        text_lengths.append(total_text)
        shape_counts.append(len(slide.shapes))
    with ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        notes_xml = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")]
    if media:
        issues.append(f"Unexpected media files found: {len(media)}")
    return {
        "slides": len(prs.slides),
        "media": len(media),
        "notes": notes,
        "slide_xml": len(slide_xml),
        "note_xml": len(notes_xml),
        "min_text": min(text_lengths),
        "max_text": max(text_lengths),
        "min_shapes": min(shape_counts),
        "max_shapes": max(shape_counts),
        "min_font": min_font if min_font != 99 else None,
        "issues": issues,
    }


def write_report(stats):
    lines = [
        "# 初稿-7 全流程详解无图大字美化版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 按用户要求：不插入图片；详细讲全部流程；字体放大；使用美观的原生流程图、步骤框和三线表风格。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个（应为 0，表示无图）",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        f"- 检测到的最小字体：{stats['min_font']} pt",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 自动检查提示如下：")
        for item in stats["issues"][:20]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 本版未插入任何论文图片、结果图或性能图。",
            "- 数据集和实验信息来自《初稿-7.docx》的表格与正文结构。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页截图预览；已完成 PPTX 包结构、媒体、备注、边界和文本密度检查。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    src = base.extract_source()
    build_deck(src)
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
