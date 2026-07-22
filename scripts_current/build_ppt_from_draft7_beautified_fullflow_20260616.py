# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BASE_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_20260616.py"
spec = importlib.util.spec_from_file_location("base_ppt", BASE_SCRIPT)
base = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(base)

WORK = ROOT / "reports" / "ppt_from_draft7_beautified_fullflow_20260616"
WORK.mkdir(parents=True, exist_ok=True)
OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_完整流程美化版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_完整流程美化版PPT_QA报告.md"


C = base.COLORS
PURPLE = RGBColor(92, 78, 148)
SOFT_PURPLE = RGBColor(235, 232, 247)
DARK = RGBColor(16, 28, 42)
SOFT_GRAY = RGBColor(241, 244, 248)


def add_bg(slide, accent=C["teal"]):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C["bg"]
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    # small academic anchor line at bottom, not decorative clutter
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.34), Inches(13.333), Inches(0.02))
    foot.fill.solid()
    foot.fill.fore_color.rgb = C["line"]
    foot.line.fill.background()


def title(slide, main, sub="", section="", accent=C["teal"]):
    if section:
        base.add_textbox(slide, section, 0.56, 0.28, 1.9, 0.22, size=8.5, color=accent, bold=True)
    base.add_textbox(slide, main, 0.55, 0.56, 11.5, 0.50, size=22.5, bold=True, color=DARK)
    if sub:
        base.add_textbox(slide, sub, 0.58, 1.06, 10.9, 0.28, size=9.7, color=C["muted"])


def footer(slide, idx, total):
    base.add_textbox(slide, "FZYC-Mol | 完整流程汇报", 0.56, 7.17, 3.8, 0.16, size=6.5, color=C["muted"])
    base.add_textbox(slide, f"{idx:02d}/{total}", 12.15, 7.17, 0.72, 0.16, size=6.5, color=C["muted"], align="right")


def note(slide, text):
    base.add_notes(slide, text)


def small_pill(slide, text, x, y, fill, w=1.45):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(8.2)
    r.font.bold = True
    r.font.color.rgb = DARK
    return shape


def process_node(slide, num, head, body, x, y, w=2.2, h=0.9, fill=C["white"], accent=C["teal"]):
    base.add_box(slide, f"{num}  {head}\n{body}", x, y, w, h, fill, accent, size=10.5, bold=True, align="left")


def section_slide(prs, slides, section_no, heading, subtitle, accent):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, accent)
    base.add_textbox(s, f"PART {section_no}", 0.75, 1.25, 2.2, 0.35, size=13, color=accent, bold=True)
    base.add_textbox(s, heading, 0.75, 1.80, 9.6, 0.75, size=30, bold=True, color=DARK)
    base.add_textbox(s, subtitle, 0.78, 2.78, 9.6, 0.50, size=14, color=C["muted"])
    for i, word in enumerate(["Problem", "Method", "Evidence", "Boundary"]):
        small_pill(s, word, 0.82 + i * 1.70, 4.55, [C["pale_blue"], C["pale_teal"], C["pale_amber"], C["pale_red"]][i])
    note(s, f"章节过渡：{heading}。这一部分重点讲 {subtitle}")
    slides.append(s)
    return s


def add_rich_note_box(slide, title_text, lines, x=8.95, y=1.52, w=3.45, h=4.65, accent=C["teal"]):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["white"]
    bg.line.color.rgb = C["line"]
    bg.line.width = Pt(0.7)
    base.add_textbox(slide, title_text, x + 0.18, y + 0.15, w - 0.36, 0.28, size=11, color=accent, bold=True)
    top = y + 0.58
    for line in lines:
        base.add_textbox(slide, "• " + line, x + 0.18, top, w - 0.32, 0.40, size=9.3, color=DARK)
        top += 0.50
    return bg


def make_deck(src, figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []
    total = 34

    def slide(main="", sub="", section="", accent=C["teal"]):
        s = prs.slides.add_slide(blank)
        add_bg(s, accent)
        if main:
            title(s, main, sub, section, accent)
        slides.append(s)
        return s

    # 1 Cover
    s = slide(accent=C["teal"])
    base.add_textbox(s, "FZYC-Mol", 0.68, 0.82, 3.4, 0.42, size=17, color=C["teal"], bold=True)
    base.add_textbox(s, "面向可靠性审计的验证集治理\n分子性质预测框架", 0.68, 1.38, 9.7, 1.22, size=30, bold=True, color=DARK)
    base.add_textbox(s, "完整流程美化版汇报 PPT | 基于初稿-7", 0.72, 2.86, 6.6, 0.30, size=13, color=C["muted"])
    # visual spine
    x, y = 0.82, 4.22
    labels = [("输入", "数据/SMILES"), ("候选", "多专家池"), ("选择", "验证门控"), ("测试", "冻结评估"), ("输出", "可靠性证据")]
    for i, (h, b) in enumerate(labels):
        process_node(s, str(i + 1), h, b, x + i * 2.45, y, 1.75, 0.72, C["white"], C["teal"])
        if i < 4:
            base.add_arrow(s, x + i * 2.45 + 1.82, y + 0.36, x + (i + 1) * 2.45 - 0.08, y + 0.36, C["teal"])
    base.add_textbox(s, "汇报目标：把论文中的每个流程讲成一条可复核证据链，而不是只展示指标。", 0.82, 5.72, 10.8, 0.42, size=14, color=DARK, bold=True)
    note(s, "开场：说明本 PPT 是完整流程版。建议用 25-35 分钟讲，重点是把 FZYC-Mol 的每一步为什么存在、如何运行、输出什么证据讲清楚。")

    # 2 One-slide thesis
    s = slide("一页讲清：FZYC-Mol 解决的不是“更高分”，而是“更可信”", "核心是把模型选择过程做成可复核流程", "Executive summary", C["blue"])
    base.add_box(s, "传统汇报方式\n模型 A 比模型 B 高多少？", 0.82, 1.65, 3.1, 0.82, C["pale_blue"], C["blue"], 13, True)
    base.add_arrow(s, 4.10, 2.06, 5.20, 2.06, C["muted"])
    base.add_box(s, "本文汇报方式\n为什么被选中、何时可信？", 5.42, 1.65, 3.35, 0.82, C["pale_teal"], C["teal"], 13, True)
    base.add_arrow(s, 8.95, 2.06, 10.05, 2.06, C["muted"])
    base.add_box(s, "最终输出\n性能 + 适用域 + 风险 + 边界", 10.25, 1.65, 2.45, 0.82, C["pale_amber"], C["amber"], 12, True)
    base.add_bullets(s, [
        "候选策略必须先登记，不能测试后补救",
        "验证集决定接受、拒绝或保留，测试集只做最终评估",
        "可靠性审计解释何时可信、何时复核、何时拒用",
        "负结果和失败案例进入主线，减少选择性报告风险",
    ], 1.00, 3.35, 9.2, 1.75, 12.7, bullet_color=C["blue"])
    add_rich_note_box(s, "讲解重点", ["先讲可信流程", "再讲指标提升", "最后讲边界", "不要承诺普遍最优"], 9.65, 3.15, 2.75, 2.25, C["blue"])
    base.add_takeaway(s, "一句话：FZYC-Mol 是验证集治理驱动的可靠性报告框架。")
    note(s, "这一页是整套 PPT 的中心句。你可以先用它告诉导师：本文不是只在刷榜，而是在解决模型开发中最容易被审稿人质疑的选择偏差与可靠性问题。")

    # 3 detailed talk map
    s = slide("完整汇报路线：八个环节串成一整套流程", "每个环节都回答一个审稿人可能追问的问题", "Roadmap", C["teal"])
    steps = [
        ("问题定义", "为什么平均指标不够"),
        ("数据与划分", "任务难度如何设计"),
        ("多视图表示", "信息从哪里来"),
        ("候选登记", "模型怎么进入候选池"),
        ("验证集治理", "谁能进入最终策略"),
        ("可靠性审计", "预测何时可信"),
        ("压力测试", "边界在哪里"),
        ("决策输出", "如何复用与投稿"),
    ]
    for i, (h, b) in enumerate(steps):
        row, col = divmod(i, 4)
        x = 0.76 + col * 3.08
        y = 1.55 + row * 1.55
        process_node(s, f"{i+1}", h, b, x, y, 2.32, 0.88, C["white"], C["teal" if i < 5 else "blue"])
    base.add_textbox(s, "建议讲法：每到一个流程，先说输入，再说处理规则，最后说输出证据。", 1.0, 5.25, 10.8, 0.42, size=14, color=DARK, bold=True)
    base.add_takeaway(s, "这不是论文目录，而是一条从问题到决策输出的流程线。")
    note(s, "这一页告诉听众后续结构。你可以强调每个流程都有输入、处理规则和输出，不是单纯按论文段落读。")

    section_slide(prs, slides, "01", "为什么需要验证集治理？", "从真实药物发现风险切入，而不是从模型堆叠开始", C["blue"])

    # 5 Problem expanded
    s = slide("真实药物发现中，模型错误会改变实验优先级", "平均分高并不自动意味着这一次预测可用", "Problem", C["blue"])
    base.add_bullets(s, [
        "早期筛选面对大量候选分子，实验资源有限，错误排序会改变后续合成与实验排队。",
        "ADMET 和毒性终点常有类别不平衡，高 ROC-AUC 可能掩盖阳性召回不足。",
        "新骨架、低相似度分子和 bRo5 化学空间容易超出训练分布。",
        "活性悬崖使结构相似分子出现性质跳变，平均 RMSE 难以反映样本级风险。",
    ], 0.95, 1.55, 7.9, 2.55, 12.2, bullet_color=C["blue"])
    add_rich_note_box(s, "审稿人会问", ["测试集是否被反复调参？", "低相似度分子表现如何？", "毒性阳性是否召回？", "失败案例在哪里？"], 9.25, 1.55, 3.0, 2.65, C["blue"])
    base.add_box(s, "本文回答：用验证集治理和可靠性审计，让每个结论都能追溯到证据。", 1.05, 5.15, 10.7, 0.64, C["pale_blue"], C["blue"], 15, True)
    base.add_takeaway(s, "问题不是“模型不够复杂”，而是“模型选择和使用边界不够透明”。")
    note(s, "本页扩展背景。强调模型错误对实际药物发现流程的影响，不要只说 benchmark。")

    # 6 failure matrix
    s = slide("四类典型失效场景决定评价协议", "FZYC-Mol 的数据和分析模块围绕这些风险展开", "Problem", C["blue"])
    matrix = [
        ["失效场景", "表现形式", "本文对应分析"],
        ["结构外推", "新骨架、低相似度", "scaffold / structure-separated"],
        ["不平衡毒性", "阳性少、召回不稳", "PR-AUC / Brier / fixed precision"],
        ["活性悬崖", "相似结构性质跳变", "MoleculeACE / cliff-pair"],
        ["bRo5 空间", "规则五以外化学空间", "公共压力测试 / 适用域审计"],
    ]
    base.add_native_table(s, matrix, 0.85, 1.45, 8.25, 3.25, font_size=10.0, header_fill=C["pale_blue"])
    add_rich_note_box(s, "讲解方式", ["先讲风险", "再讲对应证据", "最后讲边界", "避免泛化过强"], 9.55, 1.45, 2.8, 2.75, C["blue"])
    base.add_takeaway(s, "评价协议的每个模块都对应一个实际失效问题。")
    note(s, "用矩阵把“为什么要做这么多实验”讲清楚。每个实验不是为了堆数量，而是为了覆盖一种真实风险。")

    section_slide(prs, slides, "02", "FZYC-Mol 怎么运行？", "从数据输入、表示、候选池到验证门控逐步拆解", C["teal"])

    # 8 full pipeline visual
    s = slide("端到端流程：输入、候选、选择、测试、报告", "这页是全流程地图，后续逐步放大每个环节", "Full workflow", C["teal"])
    nodes = [
        ("1", "公开数据", "MoleculeNet\nTDC / ACE / bRo5"),
        ("2", "数据划分", "train / val / test\nscaffold / low-sim"),
        ("3", "多视图表示", "图、指纹、描述符\n片段、冻结表征"),
        ("4", "候选专家", "强基线、图模型\n融合、补救头"),
        ("5", "验证门控", "接受 / 拒绝 / 保留\n冻结权重阈值"),
        ("6", "最终测试", "一次性报告\n不再调参"),
        ("7", "可靠性审计", "风险、校准、保形\nOOD、失败案例"),
    ]
    for i, (num, h, b) in enumerate(nodes):
        x = 0.48 + (i % 4) * 3.15
        y = 1.40 + (i // 4) * 1.65
        process_node(s, num, h, b, x, y, 2.35, 0.95, C["white"], C["teal"])
        if i < 6 and i % 4 != 3:
            base.add_arrow(s, x + 2.42, y + 0.48, x + 2.98, y + 0.48, C["teal"])
    base.add_arrow(s, 10.03, 1.88, 10.03, 2.73, C["teal"])
    base.add_textbox(s, "关键边界：测试集只在第 6 步出现，任何候选生成、权重拟合和阈值设定都必须在此前完成。", 0.78, 5.28, 11.4, 0.45, size=13.5, color=DARK, bold=True)
    base.add_takeaway(s, "FZYC-Mol 的流程价值来自“先冻结，再测试”。")
    note(s, "这页是完整流程总图。建议慢一点讲，因为后续所有内容都围绕这七步展开。")

    # 9 data input
    s = slide("流程 1：数据输入与任务定义先固定", "公开数据、任务类型、主指标和边界要在训练前定义", "Flow 1", C["teal"])
    base.add_native_table(s, base.extract_source()["tables"][0], 0.58, 1.30, 12.05, 3.90, font_size=7.2, header_fill=C["pale_teal"])
    base.add_bullets(s, ["回归任务以 RMSE 为主指标", "分类任务以 ROC-AUC 为主指标", "bRo5 仅作为公共压力测试，非独立盲测"], 0.85, 5.55, 10.4, 0.76, 10.5, bullet_color=C["teal"])
    base.add_takeaway(s, "数据用途先定义，后续结论才不会越界。")
    note(s, "讲表1：MoleculeNet 是主面板，TDC 是外部 ADMET，MoleculeACE 是活性悬崖，bRo5 是公共压力测试。")

    # 10 split
    s = slide("流程 2：划分策略决定评价难度", "Random、scaffold、structure-separated 和 low-similarity 对应不同使用情境", "Flow 2", C["teal"])
    rows = [
        ("Random split", "同分布插值能力", "不能单独代表真实外推"),
        ("Scaffold split", "新骨架外推", "更接近药物发现筛选"),
        ("Structure-separated", "强化结构隔离", "测试更严格"),
        ("Low-similarity subset", "近邻不足场景", "暴露适用域风险"),
    ]
    for i, (h, b, c) in enumerate(rows):
        y = 1.40 + i * 0.88
        base.add_box(s, h, 0.82, y, 2.35, 0.45, C["pale_teal"], C["teal"], 11, True)
        base.add_textbox(s, b, 3.45, y + 0.05, 3.0, 0.25, size=11.5, color=DARK, bold=True)
        base.add_textbox(s, c, 6.72, y + 0.05, 4.9, 0.25, size=11, color=C["muted"])
    base.add_picture_fit(s, base.extract_figures()[10], 0.90, 5.02, 10.85, 1.28)
    base.add_textbox(s, "Source: 图11（局部示意）", 0.95, 6.36, 1.7, 0.15, size=6.5, color=C["muted"])
    base.add_takeaway(s, "划分不是技术细节，而是结论外推范围的边界。")
    note(s, "这页讲划分策略，说明不同划分不只是重复实验，而是模拟不同真实难度。")

    # 11 representations
    s = slide("流程 3：多视图表示覆盖不同化学信息层次", "每类表示进入候选池后仍需验证集决定是否保留", "Flow 3", C["teal"])
    views = [
        ("分子图", "原子-键连接\nGNN / D-MPNN"),
        ("指纹", "Morgan / MACCS\n局部子结构"),
        ("描述符", "RDKit 二维性质\n快速表格基线"),
        ("片段/骨架", "BRICS / Murcko\n解释与富集"),
        ("冻结表征", "ChemBERTa / MolT5\n预训练上下文"),
    ]
    for i, (h, b) in enumerate(views):
        process_node(s, str(i + 1), h, b, 0.58 + i * 2.52, 1.48, 2.02, 1.05, C["white"], C["teal"])
    base.add_picture_fit(s, base.extract_figures()[0], 1.05, 3.08, 10.90, 2.78)
    base.add_textbox(s, "Source: 图1", 1.10, 5.96, 1.0, 0.15, size=6.5, color=C["muted"])
    base.add_takeaway(s, "多视图设计不是为了堆复杂度，而是为了让表示选择可审计。")
    note(s, "讲清楚五类表示。注意不要说复杂模型一定更优，而是说每类表示有不同证据角色。")

    # 12 candidates
    s = slide("流程 4：候选池包含单专家、融合、门控和补救", "候选越多，越需要预登记和验证门控", "Flow 4", C["teal"])
    groups = [
        ("单专家", "RF / ExtraTrees / XGBoost\nLightGBM / CatBoost / Chemprop"),
        ("目标变换", "log1p / 分位数正态化\n截尾目标 / 类别权重"),
        ("融合", "Top-K 均值\n堆叠 / 自适应共识"),
        ("可靠性模块", "适用域门控\n不确定性加权 / 补救头"),
    ]
    for i, (h, b) in enumerate(groups):
        base.add_box(s, h + "\n" + b, 0.85 + i * 3.05, 1.55, 2.35, 1.28, [C["white"], C["pale_blue"], C["pale_teal"], C["pale_amber"]][i], C["line"], 10.5, True)
    base.add_picture_fit(s, base.extract_figures()[2], 1.15, 3.40, 10.65, 2.30)
    base.add_takeaway(s, "候选策略不是越多越好，关键是能否被验证集证据接受。")
    note(s, "这一页讲候选池组成。提醒：候选策略必须在测试前登记。")

    # 13 governance
    s = slide("流程 5：验证集治理负责接受、拒绝或保留候选", "这是防止测试集开发化的关键", "Flow 5", C["teal"])
    flow = ["固定任务", "训练专家", "验证预测", "构造融合", "选择/冻结", "一次测试"]
    for i, h in enumerate(flow):
        x = 0.75 + i * 2.05
        base.add_box(s, f"{i+1}\n{h}", x, 1.55, 1.35, 0.85, C["white"], C["teal"], 12, True)
        if i < len(flow) - 1:
            base.add_arrow(s, x + 1.42, 1.98, x + 1.93, 1.98, C["teal"])
    add_rich_note_box(s, "规则解释", [
        "验证集用于策略选择",
        "测试集只在冻结后评估",
        "失败模块记录为负结果",
        "选择器风险另行审计",
    ], 0.95, 3.25, 4.65, 2.35, C["teal"])
    add_rich_note_box(s, "输出文件应包括", [
        "候选登记表",
        "验证/测试预测",
        "随机种子与划分",
        "负结果和失败案例",
    ], 6.15, 3.25, 4.65, 2.35, C["blue"])
    base.add_takeaway(s, "验证集治理让“为什么选这个模型”有证据可查。")
    note(s, "重点讲验证集治理。这里是整篇论文最核心的方法学贡献。")

    # 14 reliability flow
    s = slide("流程 6：可靠性审计把预测值变成决策证据", "风险、校准、保形预测和解释共同构成输出层", "Flow 6", C["teal"])
    rel = [
        ("适用域", "最近邻相似度\n骨架距离"),
        ("风险分数", "模型分歧\n预测偏差"),
        ("校准", "Brier / ECE\n可靠性曲线"),
        ("保形预测", "80/90/95%\n覆盖率"),
        ("解释", "基序 / 片段\n最近邻案例"),
    ]
    for i, (h, b) in enumerate(rel):
        process_node(s, str(i + 1), h, b, 0.58 + i * 2.52, 1.50, 2.02, 1.05, C["white"], C["blue" if i < 2 else "teal"])
    base.add_box(s, "最终输出：点预测 + 适用域证据 + 风险分位 + 校准/覆盖 + 拒用理由", 1.15, 3.45, 10.75, 0.72, C["pale_teal"], C["teal"], 15, True)
    base.add_bullets(s, ["风险分数用于复核排序，不是自动纠错", "保形区间提供覆盖率意义，不替代实验验证", "解释模块只作关联证据，不作因果结论"], 1.05, 5.12, 9.6, 0.88, 11.8)
    base.add_takeaway(s, "可靠性审计回答“这次预测能不能用”。")
    note(s, "讲每个可靠性模块的输出含义和边界，避免把风险分数说成自动纠错机制。")

    section_slide(prs, slides, "03", "主要证据如何支持框架？", "先看主结果，再看强基线、外部验证和可靠性审计", C["blue"])

    # 16 Fig2 workflow large
    s = slide("图2总览：FZYC-Mol 的流程从数据一直到可靠性输出", "作为后续结果页的视觉总索引", "Workflow evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[1], 0.62, 1.30, 11.9, 5.25)
    base.add_textbox(s, "Source: 图2", 0.72, 6.56, 1.0, 0.15, size=6.5, color=C["muted"])
    base.add_takeaway(s, "图2说明每个结果模块都挂接在同一治理流程上。")
    note(s, "这页用来承接从方法到结果。提醒听众：后面的结果不是散点，而是图2中各输出模块的证据。")

    # 17 main result table
    s = slide("主结果表：选择性增益比“全面最优”更可信", "MoleculeNet 主面板显示终点依赖的接受模式", "Main evidence", C["blue"])
    base.add_native_table(s, base.extract_source()["tables"][1], 0.55, 1.25, 12.15, 3.35, font_size=7.3, header_fill=C["pale_blue"])
    base.add_bullets(s, ["ESOL / BACE：验证选择与测试观测最优一致", "Lipophilicity：补救头被验证集接受", "FreeSolv：物理相互作用相关边界案例", "ClinTox：需结合 PR-AUC、校准与固定精度召回"], 0.88, 4.95, 11.0, 1.0, 10.8)
    base.add_takeaway(s, "主结果支持“受验证集约束的选择性增益”。")
    note(s, "讲表2，强调每个保留策略都经过验证集门控。")

    # 18 Figure ranking
    s = slide("图4：不同终点偏好的模型家族不同", "候选池存在必要性，单一主干模型难以覆盖所有终点", "Main evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[3], 0.62, 1.24, 8.8, 5.24)
    add_rich_note_box(s, "读图要点", ["不同终点最优家族不同", "简单模型有时更稳定", "候选选择需透明", "结果应按终点解释"], 9.70, 1.55, 2.75, 2.65, C["blue"])
    base.add_textbox(s, "Source: 图4", 0.72, 6.54, 1.0, 0.15, size=6.5, color=C["muted"])
    base.add_takeaway(s, "终点差异是验证集治理的直接动机。")
    note(s, "这一页从图4讲模型家族差异，说明为什么不能只押一个模型。")

    # 19 Fig5
    s = slide("图5：主性能比较需要和冻结规则一起解释", "性能图只回答结果，不能单独回答选择是否公平", "Main evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[4], 0.60, 1.22, 11.95, 4.9)
    base.add_textbox(s, "Source: 图5", 0.72, 6.22, 1.0, 0.15, size=6.5, color=C["muted"])
    base.add_box(s, "解释顺序：先说验证集是否接受，再说测试集表现，最后说边界案例。", 1.05, 6.48, 10.6, 0.32, C["white"], C["line"], 10.5, True)
    note(s, "这页视觉上以图为主。讲的时候不要只读柱子，要把选择规则带回来。")

    # 20 rescue
    s = slide("图6：补救头只有在验证集接受时才进入最终策略", "定向补救是受控增强，不是测试后调参", "Rescue evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[5], 0.60, 1.25, 7.6, 4.9)
    add_rich_note_box(s, "流程解释", ["先进入候选池", "验证集判断是否接受", "测试前冻结策略", "未通过即记录负结果"], 8.65, 1.55, 3.55, 2.55, C["blue"])
    base.add_box(s, "关键案例：Lipophilicity 的补救头被验证集接受，其他终点保留原策略。", 8.65, 4.55, 3.55, 0.70, C["pale_blue"], C["blue"], 11.5, True, "left")
    base.add_takeaway(s, "补救模块的可信度来自门控，而不是来自事后挑选。")
    note(s, "解释定向补救。强调补救头没有在所有任务强行进入最终结果。")

    # 21 fusion
    s = slide("图7：多方法融合在 BBBP 和 ClinTox 形成选择性增益", "不平衡分类任务必须额外解释校准和召回", "Fusion evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[6], 0.60, 1.25, 8.5, 4.95)
    add_rich_note_box(s, "ClinTox 讲法", ["ROC-AUC 高", "阳性样本少", "需看 PR-AUC", "需看固定精度召回"], 9.45, 1.55, 2.85, 2.55, C["red"])
    base.add_takeaway(s, "分类结果的可信解释必须超越单一 ROC-AUC。")
    note(s, "讲 ClinTox 的时候一定加上不平衡标签和假阴性风险。")

    # 22 external
    s = slide("外部 ADMET：22 个终点支持选择性而非普遍增益", "win/tie/loss = 5/17/0，说明框架不会强制复杂化", "External evidence", C["blue"])
    base.add_native_table(s, base.extract_source()["tables"][2], 0.70, 1.30, 5.65, 2.50, font_size=8.5, header_fill=C["pale_blue"])
    base.add_picture_fit(s, base.extract_figures()[7], 6.70, 1.20, 5.7, 2.35)
    base.add_picture_fit(s, base.extract_figures()[8], 6.70, 3.72, 5.7, 2.35)
    base.add_takeaway(s, "外部验证的重点是受控接受，不是制造更多胜利。")
    note(s, "讲 TDC 结果：五个终点增强，其余保持原策略，这恰好证明框架克制。")

    # 23 reliability
    s = slide("可靠性分析：分类错误更容易被风险分数识别", "回归高误差仍受局部跳变和物理因素限制", "Reliability evidence", C["blue"])
    base.add_picture_fit(s, base.extract_figures()[9], 0.62, 1.25, 7.15, 5.02)
    base.add_native_table(s, base.extract_source()["tables"][3], 8.05, 1.30, 4.25, 3.15, font_size=7.5, header_fill=C["pale_teal"])
    base.add_box(s, "解释边界：风险分数是复核排序，不是自动拒用规则。", 8.05, 4.85, 4.25, 0.58, C["pale_teal"], C["teal"], 11.5, True, "left")
    base.add_takeaway(s, "可靠性输出让模型知道“自己什么时候不稳”。")
    note(s, "讲风险分数：分类错误 AUROC 中位 0.788，回归高误差 0.652，说明两类任务可靠性结构不同。")

    # 24 conformal/ranking
    s = slide("保形预测与排名审计：模型可信度也要审计选择器", "覆盖率接近目标，但验证集第一名不等于测试最优", "Audit evidence", C["blue"])
    cards = [
        ("分类覆盖率", "0.814 / 0.918 / 0.956", "目标 80/90/95%"),
        ("回归覆盖率", "0.823 / 0.925 / 0.962", "目标 80/90/95%"),
        ("Spearman", "0.667", "验证-测试中等相关"),
        ("Top-3 / Top-1", "0.295 / 0.135", "排序不稳定"),
    ]
    for i, (h, v, b) in enumerate(cards):
        x = 0.90 + (i % 2) * 5.75
        y = 1.55 + (i // 2) * 1.50
        base.add_box(s, f"{h}\n{v}\n{b}", x, y, 4.60, 1.0, C["white"], C["line"], 13, True)
    base.add_textbox(s, "讲解逻辑：保形预测检查输出覆盖率，排名审计检查候选选择是否稳定；两者分别回答“预测是否有边界”和“选择是否可相信”。", 1.0, 5.10, 10.8, 0.58, size=13, color=DARK)
    base.add_takeaway(s, "FZYC-Mol 不只审计模型，还审计模型选择过程。")
    note(s, "这一页可以慢讲，因为它很适合回答审稿人关于验证集过拟合的质疑。")

    section_slide(prs, slides, "04", "边界、失败和复用", "把负结果、压力测试和决策输出讲成方法可信度的一部分", C["amber"])

    # 26 OOD
    s = slide("图11：低相似度和结构分离暴露外推风险", "性能随相似度下降而退化，是模型边界而非异常现象", "Boundary evidence", C["amber"])
    base.add_picture_fit(s, base.extract_figures()[10], 0.62, 1.25, 8.45, 5.0)
    add_rich_note_box(s, "讲解重点", ["不要说模型消除 OOD", "应说模型标记风险区域", "活性悬崖需样本级复核", "平均指标不能替代边界分析"], 9.40, 1.55, 2.95, 2.95, C["amber"])
    base.add_takeaway(s, "压力测试告诉我们模型何时该谨慎使用。")
    note(s, "讲 OOD 与活性悬崖：这部分不是削弱模型，而是明确方法边界。")

    # 27 negative results
    s = slide("负结果：未通过验证门控的模块仍然有信息价值", "失败模块帮助界定方法适用范围", "Negative evidence", C["amber"])
    base.add_native_table(s, base.extract_source()["tables"][4], 0.75, 1.35, 6.10, 3.15, font_size=8.3, header_fill=C["pale_amber"])
    base.add_picture_fit(s, base.extract_figures()[11], 7.10, 1.35, 5.05, 3.25)
    base.add_bullets(s, ["FreeSolv：物理相互作用边界", "bRo5：公共压力测试，非独立盲测", "基序归因：关联解释，非因果机制"], 0.95, 5.15, 10.5, 0.85, 11.5, bullet_color=C["amber"])
    base.add_takeaway(s, "负结果进入主文，能减少选择性报告风险。")
    note(s, "讲负结果的重要性。特别强调 bRo5 不能讲成独立实验盲测。")

    # 28 interpretation
    s = slide("图13：化学解释用于复核，不用于证明机制", "基序、片段和最近邻案例必须和统计校正一起解释", "Interpretability", C["amber"])
    base.add_picture_fit(s, base.extract_figures()[12], 0.62, 1.25, 8.65, 5.05)
    add_rich_note_box(s, "安全表述", ["提供关联证据", "提示高误差样本", "辅助化学复核", "不推出因果机制"], 9.55, 1.55, 2.7, 2.55, C["amber"])
    base.add_takeaway(s, "解释模块越接近化学结论，越要控制主张强度。")
    note(s, "讲这页时避免说“证明某片段导致性质变化”。改说“与高误差或模型行为相关”。")

    # 29 decision card
    s = slide("最终复用形式：把每个预测写成决策卡", "这能把算法输出转化为药物化学可读证据", "Reuse")
    fields = [
        ("点预测", "性质值或类别概率"),
        ("适用域", "最近邻相似度与骨架距离"),
        ("风险分位", "复核优先级"),
        ("校准/覆盖", "ECE、Brier、保形区间"),
        ("解释证据", "片段、基序、最近邻案例"),
        ("拒用理由", "超出适用域或证据不足"),
    ]
    for i, (h, b) in enumerate(fields):
        x = 0.90 + (i % 3) * 4.05
        y = 1.45 + (i // 3) * 1.55
        base.add_box(s, h + "\n" + b, x, y, 3.05, 0.92, C["white"], C["teal" if i < 3 else "blue"], 11.5, True)
    base.add_box(s, "读法：信任不是一个二元判断，而是性能、适用域、风险和解释共同形成的证据状态。", 1.0, 5.15, 10.95, 0.65, C["pale_teal"], C["teal"], 14, True, "left")
    base.add_takeaway(s, "FZYC-Mol 最终服务的是“信、复核、拒用”的决策。")
    note(s, "这页是复用场景。可以说未来软件化时，每个分子应该输出这样的决策卡，而不是只给一个分数。")

    # 30 reproducibility
    s = slide("投稿与复现材料：让审稿人能追踪每个结论", "流程可信还需要文件可信", "Reproducibility")
    items = [
        ("候选登记表", "所有专家、融合、补救和门控策略"),
        ("split seeds", "训练/验证/测试划分与随机种子"),
        ("预测文件", "验证集和测试集逐样本预测"),
        ("source data", "每张图和表的原始数据"),
        ("统计脚本", "paired difference、bootstrap、Wilcoxon"),
        ("环境文件", "包版本、运行命令、硬件说明"),
    ]
    for i, (h, b) in enumerate(items):
        x = 0.80 + (i % 2) * 5.85
        y = 1.35 + (i // 2) * 1.15
        base.add_box(s, h + "\n" + b, x, y, 4.80, 0.70, C["white"], C["line"], 10.6, True, "left")
    base.add_takeaway(s, "论文要能被读懂，也要能被复核。")
    note(s, "这页强调投稿前材料准备。尤其是候选登记和逐样本预测，对回应审稿质疑很重要。")

    # 31 limitations
    s = slide("投稿前必须守住的四条边界", "边界清楚不是削弱论文，而是保护论文", "Limitations")
    pairs = [
        ("不要写", "所有终点普遍提升", "应写", "部分终点选择性增益"),
        ("不要写", "bRo5 独立盲测", "应写", "公开压力测试"),
        ("不要写", "风险分数自动纠错", "应写", "复核排序提示"),
        ("不要写", "基序解释因果机制", "应写", "关联性化学证据"),
    ]
    for i, (a, b, c, d) in enumerate(pairs):
        y = 1.45 + i * 0.92
        base.add_box(s, a, 0.88, y, 1.20, 0.40, C["pale_red"], C["red"], 10, True)
        base.add_textbox(s, b, 2.25, y + 0.06, 3.6, 0.22, size=11.5, color=C["red"])
        base.add_arrow(s, 6.10, y + 0.20, 6.85, y + 0.20, C["muted"])
        base.add_box(s, c, 7.10, y, 1.20, 0.40, C["pale_teal"], C["teal"], 10, True)
        base.add_textbox(s, d, 8.45, y + 0.06, 3.6, 0.22, size=11.5, color=C["teal"], bold=True)
    base.add_takeaway(s, "Nature 风格更看重证据强度与主张边界是否匹配。")
    note(s, "这页可以直接用于导师讨论和投稿前自查。逐条说明哪些话不能说过头。")

    # 32 oral script
    s = slide("讲解顺序建议：每个流程都按三句话展开", "让听众跟得上，不被模型名和指标淹没", "Presentation script")
    recipe = [
        ("一句话定义", "这个流程解决什么问题？"),
        ("一个关键规则", "输入是什么，谁做决策？"),
        ("一个输出证据", "结果如何支持或限制主张？"),
    ]
    for i, (h, b) in enumerate(recipe):
        base.add_box(s, f"{i+1}\n{h}\n{b}", 1.05 + i * 3.85, 1.70, 2.85, 1.25, [C["pale_blue"], C["pale_teal"], C["pale_amber"]][i], C["line"], 13, True)
    base.add_bullets(s, [
        "流程页：先讲输入和规则，不急着讲数字",
        "结果页：先讲验证门控，再讲测试表现",
        "边界页：主动说明未完成验证和失败案例",
        "总结页：回到“可复核模型治理”而非“万能模型”",
    ], 1.15, 4.05, 9.8, 1.25, 12.5)
    base.add_takeaway(s, "好的汇报不是信息最多，而是每个流程都有清楚因果线。")
    note(s, "这是给你讲 PPT 时用的提示页。可根据场合保留或隐藏。")

    # 33 conclusion
    s = slide("最终总结：FZYC-Mol 的贡献是可靠性流程", "性能、选择过程、适用域和失败模式一起构成论文价值", "Conclusion", C["teal"])
    base.add_textbox(s, "FZYC-Mol 将分子性质预测从“谁分数更高”推进到“谁的选择过程更可信”。", 0.88, 1.45, 11.0, 0.55, size=23, bold=True, color=DARK)
    base.add_bullets(s, [
        "方法贡献：候选登记、验证门控、测试冻结",
        "证据贡献：MoleculeNet、TDC ADMET、MoleculeACE、bRo5 公共压力测试",
        "可靠性贡献：风险分数、保形预测、验证-测试排名审计",
        "边界贡献：负结果、失败案例和非盲测说明",
    ], 1.05, 2.75, 9.7, 1.65, 13, bullet_color=C["teal"])
    base.add_box(s, "结尾建议：把论文定位为“可靠性审计框架”，不要定位为“所有 ADMET 任务的万能预测器”。", 1.05, 5.45, 10.8, 0.60, C["pale_teal"], C["teal"], 13, True, "left")
    note(s, "最后总结贡献。重申所有流程都是为了让模型选择和可靠性边界可复核。")

    # 34 backup Q&A
    s = slide("备答：可能被问到的 6 个问题", "用于答辩或组会追问", "Q&A", PURPLE)
    qa = [
        ("为什么不用单一大模型？", "终点依赖明显，候选治理比押注单模型更稳。"),
        ("验证集会不会过拟合？", "报告 ranking audit、Top-3、optimism gap 和 nested validation。"),
        ("bRo5 是不是盲测？", "不是，只能写成公开压力测试。"),
        ("风险分数能否拒用样本？", "可作复核排序，不能自动替代实验判断。"),
        ("基序解释是否因果？", "不是，属于关联证据，需要统计和案例复核。"),
        ("论文最大卖点是什么？", "可复核的模型选择和可靠性报告流程。"),
    ]
    for i, (q, a) in enumerate(qa):
        x = 0.82 + (i % 2) * 6.05
        y = 1.35 + (i // 2) * 1.42
        base.add_box(s, q + "\n" + a, x, y, 5.10, 0.92, C["white"], PURPLE, 10.8, True, "left")
    base.add_takeaway(s, "主动回答边界问题，能显著降低审稿和答辩风险。")
    note(s, "最后这一页可作为备答页。讲正式报告时可以不讲，问答时使用。")

    for idx, s in enumerate(slides, 1):
        footer(s, idx, total)
    return prs, slides


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    for i, s in enumerate(prs.slides, 1):
        try:
            if s.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        for sh in s.shapes:
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > prs.slide_width + 2 or sh.top + sh.height > prs.slide_height + 2:
                issues.append(f"high: slide {i} shape out of bounds")
            if getattr(sh, "has_text_frame", False):
                txt = sh.text.strip()
                if len(txt) > 220 and sh.width < Inches(5.2):
                    issues.append(f"medium: slide {i} text box dense ({len(txt)} chars)")
                if "xxxx" in txt.lower() or "lorem" in txt.lower():
                    issues.append(f"high: slide {i} placeholder remains")
    with ZipFile(path) as z:
        media = sum(1 for n in z.namelist() if n.startswith("ppt/media/"))
        slide_xml = sum(1 for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        note_xml = sum(1 for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml"))
    return {"slides": len(prs.slides), "media": media, "notes": notes, "slide_xml": slide_xml, "note_xml": note_xml, "issues": issues}


def write_report(stats_before, stats_after):
    lines = [
        "# 初稿-7 完整流程美化版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt 工作流 + python-pptx。",
        "- 改进方向：更完整流程、更强视觉层次、更多讲解文字与讲者备注。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats_after['slides']} 页",
        f"- 媒体文件：{stats_after['media']} 个",
        f"- 讲者备注：{stats_after['notes']} 页",
        f"- slide XML：{stats_after['slide_xml']} 个",
        f"- note XML：{stats_after['note_xml']} 个",
        "",
        "## 自审结果",
    ]
    remaining = stats_after["issues"]
    if remaining:
        lines += [f"- {x}" for x in remaining]
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines += [
        "",
        "## 说明",
        "- 当前环境无可靠 headless 渲染器，因此未输出逐页截图预览；已进行 PPTX 包结构、媒体、备注、边界和文本密度检查。",
        "- 科学图像来自 `初稿-7.docx` 内嵌原图，未修改原始数据图。",
        "- 与上一版相比，本版从 25 页扩展到 34 页，加入完整流程拆解、章节页、讲解栏、备答页和更详细讲者备注。",
    ]
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    src = base.extract_source()
    figs = base.extract_figures()
    base.create_contact_sheet(figs)
    prs, slides = make_deck(src, figs)
    prs.save(OUT_PPTX)
    stats_before = audit(OUT_PPTX)
    # Corrective normalization pass.
    prs2 = Presentation(OUT_PPTX)
    prs2.save(OUT_PPTX)
    stats_after = audit(OUT_PPTX)
    write_report(stats_before, stats_after)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats_after)


if __name__ == "__main__":
    main()
