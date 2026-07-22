# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
FLOW_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_overall_flow_largefont_20260616.py"
spec = importlib.util.spec_from_file_location("flow_ppt", FLOW_SCRIPT)
flow = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(flow)

base = flow.base
WORK = ROOT / "reports" / "ppt_from_draft7_overall_flow_with_figs_20260616"
WORK.mkdir(parents=True, exist_ok=True)

OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_整体流程详解配图版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_整体流程详解配图版PPT_QA报告.md"


def src_label(slide, body):
    flow.tx(slide, body, 0.62, 6.33, 11.8, 0.16, size=6.8, color=flow.C["muted"])


def build_deck(figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 12

    # 1
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.tx(s, "FZYC-Mol", 0.66, 0.46, 4.8, 0.56, size=30, bold=True)
    flow.tx(s, "整体实验流程详解：从问题定义到可审计输出", 0.70, 1.04, 9.6, 0.36, size=15.5, color=flow.SLATE)
    flow.flow_row(s, ["问题定义", "数据登记", "划分冻结", "表示生成", "专家训练", "验证选择"], 0.62, 1.88, 12.05, 0.76, accent=flow.TEAL)
    flow.flow_row(s, ["冻结测试", "可靠性审计", "OOD压力", "负结果归档", "论文输出", "答辩复核"], 0.62, 3.02, 12.05, 0.76, accent=flow.BLUE)
    flow.big_box(s, "原则 1：先登记再比较", "数据集、划分、候选专家和评价指标在实验前固定。", 0.82, 4.42, 3.55, 0.92, accent=flow.TEAL)
    flow.big_box(s, "原则 2：验证集负责选择", "selector 只读取验证集，测试集在策略冻结后使用。", 4.86, 4.42, 3.55, 0.92, accent=flow.AMBER)
    flow.big_box(s, "原则 3：输出边界证据", "性能、风险、失败案例和负结果共同限制主张。", 8.90, 4.42, 3.28, 0.92, accent=flow.RED)
    flow.takeaway(s, "这版只加入流程/架构/治理图，不加入性能结果图。")
    flow.footer(s, 1, total)
    flow.note(s, "开场说明：这版以整体实验流程为主，并加入少量概念图帮助理解。所有结果类图仍不放入。")

    # 2 workflow figure
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "总览图：整体 workflow 串起所有实验阶段", "先看全局路径，再逐步拆解每个阶段怎么执行。", flow.TEAL)
    base.add_picture_fit(s, figs[1], 0.58, 1.30, 7.55, 4.80, frame=True)
    flow.big_box(s, "图怎么读", "从任务输入开始，依次经过数据划分、多视图表示、候选专家、验证集选择和证据输出。", 8.52, 1.42, 3.62, 0.96, accent=flow.TEAL)
    flow.big_box(s, "讲解重点", "每个箭头代表一个可复核文件或脚本，而不是只在图上画出的概念。", 8.52, 2.60, 3.62, 0.96, accent=flow.BLUE)
    flow.big_box(s, "流程边界", "未通过验证门控的候选进入负结果，不再回到测试集调参。", 8.52, 3.78, 3.62, 0.96, accent=flow.RED)
    src_label(s, "Source: Fig. 2 from 初稿-7.docx；本页为流程理解图，不展示性能结果。")
    flow.takeaway(s, "整体 workflow 的核心是：每一步都能追溯到输入、操作、质控和输出。")
    flow.footer(s, 2, total)
    flow.note(s, "这一页讲总流程图。不要逐字读图内小字，只讲从输入到输出的大逻辑，以及为什么验证集选择是中枢。")

    # 3 problem and claims
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.BLUE)
    flow.title(s, "阶段 0：先定义研究问题和允许主张", "在实验开始前确定论文要证明什么、不能证明什么。", flow.BLUE)
    flow.flow_row(s, ["ADMET需求", "可靠性瓶颈", "核心假设", "可检验问题", "主张边界"], 0.78, 1.58, 11.80, 0.86, accent=flow.BLUE)
    flow.two_col_steps(
        s,
        "输入与决策",
        ["输入：分子性质预测任务、ADMET 场景、公开基准与外部任务。",
         "决策：主张限定为验证治理、选择性增益和可靠性审计。",
         "不预设：某个模型家族一定最优，或融合一定提升。"],
        "质控与输出",
        ["质控：每个主张必须对应后续实验模块。",
         "输出：实验登记表、候选模块清单、评价指标清单。",
         "风险控制：不能证明的内容提前写成边界。"],
        accent=flow.BLUE,
    )
    flow.takeaway(s, "流程第一步不是跑模型，而是固定问题、主张和证据标准。", accent=flow.BLUE)
    flow.footer(s, 3, total)
    flow.note(s, "这一页说明研究定位。要强调论文目标不是万能预测器，而是可审计的验证治理框架。")

    # 4 data + split
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.AMBER)
    flow.title(s, "阶段 1-2：数据登记、清洗与划分冻结", "数据版本和划分协议必须先固定，后续实验只能引用冻结版本。", flow.AMBER)
    flow.flow_row(s, ["导入数据", "字段检查", "SMILES标准化", "标签审计", "划分生成", "版本冻结"], 0.78, 1.54, 11.80, 0.86, accent=flow.AMBER)
    flow.two_col_steps(
        s,
        "数据清洗",
        ["记录 MoleculeNet、TDC、MoleculeACE 等来源。",
         "检查 SMILES、标签、单位、阳性定义和缺失值。",
         "处理重复分子、无效 SMILES 和冲突标签。",
         "输出 cleaned.csv、dataset_registry.csv、label_audit.md。"],
        "划分冻结",
        ["生成 random、scaffold、structure-separated 和低相似度分层。",
         "记录 seed、train/valid/test index、阳性率和 scaffold 报告。",
         "冻结 split_index.json，不能按结果好坏更改划分。",
         "不同划分结果不能混在同一性能主张里。"],
        accent=flow.AMBER,
    )
    flow.takeaway(s, "数据和划分决定外推结论能否成立。", accent=flow.AMBER)
    flow.footer(s, 4, total)
    flow.note(s, "这一页合并数据和划分。强调所有清洗规则与划分文件都要能回溯。")

    # 5 architecture figure
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.PURPLE)
    flow.title(s, "阶段 3-4：表示生成与候选专家训练", "模型结构图帮助说明多视图输入如何进入专家池和证据输出。", flow.PURPLE)
    base.add_picture_fit(s, figs[0], 0.58, 1.30, 7.55, 4.70, frame=True)
    flow.big_box(s, "表示生成", "从同一冻结数据读取 SMILES，生成描述符、指纹、图结构和任务上下文。", 8.52, 1.42, 3.62, 0.92, accent=flow.PURPLE)
    flow.big_box(s, "专家训练", "强基线、图模型、指纹模型、表格模型和补救模块先作为候选登记。", 8.52, 2.54, 3.62, 0.92, accent=flow.BLUE)
    flow.big_box(s, "训练输出", "保存 valid/test 预测、模型配置、训练日志和 candidate_registry。", 8.52, 3.66, 3.62, 0.92, accent=flow.TEAL)
    flow.big_box(s, "注意", "候选专家不直接进入最终结论，必须等待 selector 判断。", 8.52, 4.78, 3.62, 0.70, accent=flow.RED)
    src_label(s, "Source: Fig. 1 from 初稿-7.docx；本页为模型结构图，不展示实验结果。")
    flow.takeaway(s, "表示与训练阶段只产生候选和验证证据，最终策略由 selector 决定。", accent=flow.BLUE)
    flow.footer(s, 5, total)
    flow.note(s, "这一页讲模型结构图。重点解释多源表示、专家预测矩阵和证据输出，不讲具体性能。")

    # 6 selector figure
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "阶段 5：selector 与门控是整套流程的中枢", "治理图说明候选如何被接受、拒绝或转入负结果。", flow.TEAL)
    base.add_picture_fit(s, figs[2], 0.58, 1.30, 7.65, 4.35, frame=True)
    flow.big_box(s, "接受", "验证性能改善，且校准、风险或适用域信号没有明显恶化。", 8.60, 1.42, 3.55, 0.82, accent=flow.TEAL)
    flow.big_box(s, "拒绝", "未通过门控的候选不进入主性能结论，保留为负结果。", 8.60, 2.42, 3.55, 0.82, accent=flow.RED)
    flow.big_box(s, "冻结", "输出 selected_strategy.json 后，才允许进入测试集评估。", 8.60, 3.42, 3.55, 0.82, accent=flow.BLUE)
    flow.big_box(s, "审计", "selector_audit.md 记录每个候选的进入、拒绝和原因。", 8.60, 4.42, 3.55, 0.82, accent=flow.AMBER)
    src_label(s, "Source: Fig. 3 from 初稿-7.docx；本页为治理流程图，不展示性能结果。")
    flow.takeaway(s, "selector 不是为了让结果好看，而是为了让模型选择过程可审计。")
    flow.footer(s, 6, total)
    flow.note(s, "这一页是核心。要把接受、拒绝、冻结和审计四件事讲清楚。")

    # 7 frozen testing
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.AMBER)
    flow.title(s, "阶段 6：冻结测试与排名审计流程", "测试集只在策略冻结后使用，用于一次性外部检验。", flow.AMBER)
    flow.flow_row(s, ["读取冻结策略", "运行test预测", "计算主指标", "rank audit", "optimism gap", "锁定主表"], 0.78, 1.54, 11.80, 0.86, accent=flow.AMBER)
    flow.two_col_steps(
        s,
        "执行步骤",
        ["读取 selected_strategy.json，不允许临时替换候选。",
         "在 test split 上一次性计算回归或分类主指标。",
         "比较 valid rank 与 test rank，记录排名错位。",
         "计算 optimism gap，判断验证优势是否过于乐观。"],
        "质控输出",
        ["输出 test_metrics.csv、rank_audit.csv、optimism_gap.md。",
         "若测试表现与验证选择冲突，必须在讨论或限制中解释。",
         "主表只来自已锁定策略，不能再回到 selector 调整。",
         "测试结果用于验证流程，不用于重新设计流程。"],
        accent=flow.AMBER,
    )
    flow.takeaway(s, "冻结测试流程回答“结果是否由测试集调出来”的关键质疑。", accent=flow.AMBER)
    flow.footer(s, 7, total)
    flow.note(s, "讲冻结测试与 rank audit。注意这里仍不展示结果，只讲如何产生和审计结果。")

    # 8 reliability
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.GREEN)
    flow.title(s, "阶段 7：可靠性审计输出流程", "性能之外，还要输出模型何时可靠、何时需要谨慎。", flow.GREEN)
    flow.flow_row(s, ["最终预测", "AD距离", "不确定性", "校准误差", "risk score", "conformal覆盖"], 0.78, 1.54, 11.80, 0.86, accent=flow.GREEN)
    flow.two_col_steps(
        s,
        "执行步骤",
        ["基于训练分布计算适用域距离或相似度信号。",
         "从模型方差、专家分歧或校准误差生成 uncertainty。",
         "构建 risk score，并生成 risk-coverage 分析文件。",
         "对分类或回归输出进行校准与保形覆盖评估。"],
        "质控输出",
        ["输出 ad_score.csv、uncertainty.csv、risk_coverage.csv。",
         "风险分数只作为审计证据，不作为单独性能胜利。",
         "高风险样本进入失败案例和边界分析。",
         "可靠性模块用于限定适用域，不替代主指标。"],
        accent=flow.GREEN,
    )
    flow.takeaway(s, "可靠性阶段把单个预测扩展为“预测值 + 风险 + 适用域 + 覆盖率”。", accent=flow.GREEN)
    flow.footer(s, 8, total)
    flow.note(s, "这一页讲可靠性文件如何产生。不放 risk-coverage 结果图，只讲生成流程和作用。")

    # 9 OOD
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.RED)
    flow.title(s, "阶段 8：OOD、低相似度与失败案例流程", "所有外推风险进入同一条错误分析链。", flow.RED)
    flow.flow_row(s, ["scaffold压力", "结构分离", "Tanimoto bin", "MoleculeACE", "失败样本", "风险富集"], 0.78, 1.54, 11.80, 0.86, accent=flow.RED)
    flow.two_col_steps(
        s,
        "执行步骤",
        ["比较 random、scaffold、structure-separated 的协议差异。",
         "按 >0.7、0.5-0.7、<0.5 建立互斥低相似度分层。",
         "对 MoleculeACE 活性悬崖样本单独记录误差和预测差异。",
         "将 ClinTox 假阴性、高风险 ADME、低相似度失败等归档。"],
        "质控输出",
        ["输出 ood_metrics.csv、tanimoto_bins.csv、cliff_pairs.csv。",
         "失败案例不能只选好看的样本，类别要覆盖低相似度和活性悬崖。",
         "风险分数若能富集失败样本，才支持可靠性审计。",
         "OOD 结果用来限定边界，不用来夸大外推能力。"],
        accent=flow.RED,
    )
    flow.takeaway(s, "失败分析告诉读者哪些预测不能直接相信。", accent=flow.RED)
    flow.footer(s, 9, total)
    flow.note(s, "讲外推边界和失败样本流程。注意要说这是保护论文主张的环节。")

    # 10 ablation
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.PURPLE)
    flow.title(s, "阶段 9：消融、负结果与解释流程", "模块贡献、失败候选和化学解释都要记录，但不能越界表述。", flow.PURPLE)
    flow.flow_row(s, ["Full策略", "w/o模块", "单专家", "简单平均", "负结果表", "解释审计"], 0.78, 1.54, 11.80, 0.86, accent=flow.PURPLE)
    flow.two_col_steps(
        s,
        "执行步骤",
        ["统一消融矩阵：Full、best single、simple mean、w/o selector、w/o fusion、w/o AD gate、w/o uncertainty。",
         "将未通过门控的 rescue head、3D-lite、bRo5 扩展等纳入负结果。",
         "对 motif attribution 和 fragment enrichment 只做关联解释。",
         "检查每个解释是否有支持度、效应量和 p/FDR。"],
        "质控输出",
        ["输出 ablation_matrix.csv、negative_results.md、motif_report.md。",
         "解释结果不能写成因果机制证明。",
         "负结果不能包装成主性能成功。",
         "主文保留摘要，补充材料放完整长表。"],
        accent=flow.PURPLE,
    )
    flow.takeaway(s, "消融和负结果证明作者没有只挑成功结果。", accent=flow.BLUE)
    flow.footer(s, 10, total)
    flow.note(s, "讲消融和负结果流程。提醒用户：解释只讲关联，不讲因果。")

    # 11 manuscript output
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.BLUE)
    flow.title(s, "阶段 10：论文输出、复现包与答辩复核流程", "把所有阶段产物汇总成可审稿、可复核、可答辩的材料。", flow.BLUE)
    flow.flow_row(s, ["主文流程图", "主表/补充表", "数据声明", "代码声明", "QA清单", "答辩问题"], 0.78, 1.54, 11.80, 0.86, accent=flow.BLUE)
    flow.big_box(s, "论文输出", "方法按数据、划分、表示、专家、selector、冻结测试、可靠性审计顺序写。", 0.78, 3.02, 3.70, 1.15, accent=flow.BLUE)
    flow.big_box(s, "复现输出", "打包 registry、split index、feature manifest、model config、metrics source 和脚本。", 4.78, 3.02, 3.70, 1.15, accent=flow.TEAL)
    flow.big_box(s, "答辩输出", "准备测试集调参、是否全面优于基线、外推边界、解释是否因果四类问题。", 8.78, 3.02, 3.38, 1.15, accent=flow.AMBER)
    flow.big_box(s, "最终检查", "每个主张必须能回指到流程产物；不能核验的内容写为限制或后续工作。", 1.05, 4.72, 11.10, 0.82, accent=flow.RED, body_size=12.2)
    flow.takeaway(s, "整体流程的终点是一套可追溯、可复核、不过度外推的证据链。", accent=flow.BLUE)
    flow.footer(s, 11, total)
    flow.note(s, "讲论文与复现输出。强调所有文件都是为了让审稿人能复核。")

    # 12 summary
    s = prs.slides.add_slide(blank); flow.add_bg(s, flow.TEAL)
    flow.title(s, "最终汇报口径：用整体流程解释可信度", "用流程而不是堆结果图来回答审稿人关心的问题。", flow.TEAL)
    flow.flow_row(s, ["先固定问题", "再冻结数据", "再训练候选", "再验证选择", "再一次测试", "最后审计边界"], 0.78, 1.58, 11.80, 0.88, accent=flow.TEAL)
    flow.big_box(s, "能证明", "验证治理流程、选择性增益、风险审计和负结果记录。", 0.90, 3.02, 3.50, 1.05, accent=flow.TEAL)
    flow.big_box(s, "不能夸大", "不能说所有终点全面最优，不能把解释结果说成因果机制。", 4.86, 3.02, 3.50, 1.05, accent=flow.RED)
    flow.big_box(s, "答辩重点", "围绕测试集是否调参、外推边界、失败案例和复现材料回答。", 8.82, 3.02, 3.28, 1.05, accent=flow.AMBER)
    flow.big_box(s, "一句话结论", "FZYC-Mol 是面向可靠性审计的验证集治理框架，而不是替代实验验证的万能预测器。", 1.10, 4.72, 11.0, 0.88, accent=flow.BLUE, body_size=12.8)
    flow.takeaway(s, "讲完整流程，才能让听众相信每个结果都有来源、每个边界都被记录。")
    flow.footer(s, 12, total)
    flow.note(s, "最后总结。强调本 PPT 适合讲整体流程，有 3 张概念图辅助理解，但不放结果图。")

    prs.save(OUT_PPTX)


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    text_lengths, shape_counts = [], []
    sw, sh = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        total_text = 0
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > sw + 1000 or shp.top + shp.height > sh + 1000:
                issues.append(f"Slide {si}: shape outside slide bounds")
            if hasattr(shp, "text") and shp.text.strip():
                total_text += len(shp.text.strip())
        if total_text < 120:
            issues.append(f"Slide {si}: low text density")
        if total_text > 850:
            issues.append(f"Slide {si}: high text density ({total_text} chars)")
        text_lengths.append(total_text)
        shape_counts.append(len(slide.shapes))
    with ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        notes_xml = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")]
    if len(media) != 3:
        issues.append(f"Expected 3 conceptual figures, found {len(media)} media files")
    return {
        "slides": len(prs.slides),
        "media": len(media),
        "notes": notes,
        "slide_xml": len(slide_xml),
        "note_xml": len(notes_xml),
        "min_text": min(text_lengths),
        "max_text": max(text_lengths),
        "min_shapes": min(shape_counts),
        "max_shapes": max(shape_counts),
        "issues": issues,
    }


def write_report(stats):
    lines = [
        "# 初稿-7 整体流程详解配图版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 按用户要求：在整体流程详解大字版基础上适当加入图；仅加入 workflow、模型结构和 selector 治理 3 张概念图。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个（3 张概念图，无结果图/性能图）",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 自动检查提示如下：")
        for item in stats["issues"]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 本版只插入概念图：整体 workflow、模型结构、selector 治理。",
            "- 未插入 MoleculeNet 主结果图、TDC 结果图、risk-coverage 曲线、OOD 压力测试图或解释结果图。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页截图预览；已完成 PPTX 包结构、媒体、备注、边界和文本密度检查。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    figs = base.extract_figures()
    if len(figs) < 3:
        raise RuntimeError(f"Expected at least 3 figures, got {len(figs)}")
    build_deck(figs)
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
