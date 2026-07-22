# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DENSE_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_simplified_dense_20260616.py"
spec = importlib.util.spec_from_file_location("dense_ppt", DENSE_SCRIPT)
dense = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(dense)

base = dense.base
WORK = ROOT / "reports" / "ppt_from_draft7_overall_flow_largefont_20260616"
WORK.mkdir(parents=True, exist_ok=True)

OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_整体流程详解大字版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_整体流程详解大字版PPT_QA报告.md"


C = dense.C
DARK = dense.DARK
SLATE = dense.SLATE
BLUE = dense.BLUE
TEAL = dense.TEAL
AMBER = dense.AMBER
RED = dense.RED
GREEN = dense.GREEN
PURPLE = RGBColor(94, 83, 158)
PALE_PURPLE = RGBColor(238, 236, 249)


def add_bg(slide, accent=TEAL):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = dense.LIGHT
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(7.18), Inches(12.36), Inches(0.018))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = C["line"]
    bottom.line.fill.background()


def tx(slide, body, x, y, w, h, size=13, color=DARK, bold=False, align="left", valign="top"):
    return base.add_textbox(
        slide, body, x, y, w, h, size=size, color=color, bold=bold,
        align=align, valign=valign, margin=0.06,
    )


def title(slide, main, sub="", accent=TEAL):
    tx(slide, main, 0.55, 0.34, 11.8, 0.52, size=25, color=DARK, bold=True)
    if sub:
        tx(slide, sub, 0.58, 0.88, 11.2, 0.30, size=12.0, color=C["muted"])
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.22), Inches(1.22), Inches(0.04))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()


def footer(slide, idx, total):
    tx(slide, "FZYC-Mol | 整体实验流程详解", 0.55, 7.22, 4.0, 0.15, size=7.2, color=C["muted"])
    tx(slide, f"{idx:02d}/{total}", 12.12, 7.22, 0.72, 0.15, size=7.2, color=C["muted"], align="right")


def big_box(slide, head, body, x, y, w, h, accent=TEAL, fill=C["white"], head_size=13, body_size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = C["line"]
    box.line.width = Pt(0.8)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    tx(slide, head, x + 0.18, y + 0.10, w - 0.32, 0.26, size=head_size, color=accent, bold=True)
    tx(slide, body, x + 0.18, y + 0.44, w - 0.32, h - 0.52, size=body_size, color=DARK)
    return box


def arrow(slide, x1, y1, x2, y2, color=C["muted"]):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.6)
    line.line.end_arrowhead = True
    return line


def flow_node(slide, label, x, y, w, h, accent=TEAL, fill=C["white"], size=12.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = DARK
    return shape


def flow_row(slide, items, x, y, w, h, accent=TEAL):
    gap = 0.15
    item_w = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        nx = x + i * (item_w + gap)
        flow_node(slide, item, nx, y, item_w, h, accent=accent, size=11.7)
        if i < len(items) - 1:
            arrow(slide, nx + item_w + 0.02, y + h / 2, nx + item_w + gap - 0.04, y + h / 2, color=accent)


def two_col_steps(slide, left_title, left_items, right_title, right_items, accent=TEAL):
    big_box(slide, left_title, "\n".join(left_items), 0.72, 3.08, 5.65, 2.80, accent=accent, body_size=11.5)
    big_box(slide, right_title, "\n".join(right_items), 6.78, 3.08, 5.65, 2.80, accent=BLUE if accent != BLUE else TEAL, body_size=11.5)


def takeaway(slide, body, accent=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(6.55), Inches(12.20), Inches(0.43))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C["pale_teal"] if accent == TEAL else C["pale_blue"]
    bar.line.fill.background()
    tx(slide, "Takeaway  " + body, 0.75, 6.66, 11.75, 0.20, size=11.0, color=DARK, bold=True)


def note(slide, body):
    base.add_notes(slide, body)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 12

    # 1
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    tx(s, "FZYC-Mol", 0.66, 0.46, 4.8, 0.56, size=30, bold=True)
    tx(s, "整体实验流程详解：从问题定义到可审计输出", 0.70, 1.04, 9.6, 0.36, size=15.5, color=SLATE)
    flow_row(
        s,
        ["问题定义", "数据登记", "划分冻结", "表示生成", "专家训练", "验证选择"],
        0.62, 1.88, 12.05, 0.76, accent=TEAL,
    )
    flow_row(
        s,
        ["冻结测试", "可靠性审计", "OOD压力", "负结果归档", "论文输出", "答辩复核"],
        0.62, 3.02, 12.05, 0.76, accent=BLUE,
    )
    big_box(s, "流程原则 1：先登记再比较", "所有数据集、划分、候选专家和评价指标在实验前固定，避免后验挑选。", 0.82, 4.42, 3.55, 0.92, accent=TEAL)
    big_box(s, "流程原则 2：验证集负责选择", "selector 只读取验证集，测试集只在策略冻结后使用一次。", 4.86, 4.42, 3.55, 0.92, accent=AMBER)
    big_box(s, "流程原则 3：输出边界证据", "性能、校准、风险、失败案例和负结果共同决定论文主张边界。", 8.90, 4.42, 3.28, 0.92, accent=RED)
    takeaway(s, "这一版 PPT 不放结果图，只讲完整实验流程和每个阶段的输入、操作、质控与输出。")
    footer(s, 1, total)
    note(s, "开场说明：这版不是结果汇报，而是流程汇报。按从问题到输出的顺序讲，每一阶段都说明输入、处理、质量控制和产出。")

    # 2
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "阶段 0：先定义研究问题和允许主张", "在实验开始前确定论文要证明什么、不能证明什么。", BLUE)
    flow_row(s, ["临床/ADMET需求", "模型可靠性瓶颈", "核心假设", "可检验问题", "主张边界"], 0.78, 1.58, 11.80, 0.86, accent=BLUE)
    two_col_steps(
        s,
        "输入与决策",
        ["输入：分子性质预测任务、ADMET 场景、公开基准与外部任务。",
         "决策：论文主张限定为验证治理、选择性增益和可靠性审计。",
         "不决策：不预设某个模型家族一定最优，不预设融合一定提升。"],
        "质控与输出",
        ["质控：每个主张必须对应后续实验模块，不能用单一指标支撑全部结论。",
         "输出：实验登记表、候选模块清单、评价指标清单、负结果记录规则。",
         "风险控制：把“不能证明”的内容提前写成边界，而不是后期补救。"],
        accent=BLUE,
    )
    takeaway(s, "流程第一步不是跑模型，而是固定问题、主张和证据标准。", accent=BLUE)
    footer(s, 2, total)
    note(s, "这一页讲为什么要做验证集治理。重点说明本研究不是追求所有任务最优，而是建立可审计的预测流程。")

    # 3
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "阶段 1：数据登记与清洗流程", "所有数据先进入统一 registry，后续实验只引用冻结版本。", TEAL)
    flow_row(s, ["下载/导入", "字段检查", "SMILES标准化", "重复处理", "标签审计", "版本冻结"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 记录来源：MoleculeNet、TDC、MoleculeACE 或其他公开数据。",
         "2. 检查字段：SMILES、任务标签、单位、分类阳性定义、缺失值。",
         "3. 标准化分子：盐形式、重复项、无效 SMILES、冲突标签。",
         "4. 生成 registry：每个数据集保存版本号、样本量、任务类型和主指标。"],
        "质控输出",
        ["输出：cleaned.csv、dataset_registry.csv、label_audit.md。",
         "质控：任何剔除样本必须有原因；同一分子冲突标签不能静默保留。",
         "冻结：后续划分、训练、评估均引用同一数据版本。",
         "写作：数据清洗只写可核验规则，不写模糊处理。"],
        accent=TEAL,
    )
    takeaway(s, "数据登记让后续每个结果都能追溯到同一个清洗版本。")
    footer(s, 3, total)
    note(s, "这一页讲数据如何进入实验。强调所有公开数据都要先清洗和登记，避免后续出现同一数据集多个版本导致结果不一致。")

    # 4
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "阶段 2：划分策略与泄漏控制流程", "划分不是技术细节，而是决定外推结论能否成立的核心协议。", AMBER)
    flow_row(s, ["读取 registry", "生成 seeds", "random split", "scaffold split", "结构分离", "低相似度分层"], 0.78, 1.54, 11.80, 0.86, accent=AMBER)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 固定随机种子和划分脚本，生成 train/valid/test 三部分。",
         "2. 对 MoleculeNet 同时保留 random、scaffold 和结构压力划分。",
         "3. 对低相似度样本按 Tanimoto bin 建立互斥分层。",
         "4. 对外部 ADMET 使用官方划分或明确记录重新划分规则。"],
        "质控输出",
        ["输出：split_index.json、seed_list.txt、scaffold_report.md。",
         "质控：训练集与测试集不能共享泄漏分子；划分后任务阳性率要记录。",
         "冻结：划分一旦确定，不因模型结果好坏而更改。",
         "写作：不同划分结果不能混在同一性能主张里。"],
        accent=AMBER,
    )
    takeaway(s, "划分流程决定“模型是否真能外推”，必须在训练前冻结。", accent=AMBER)
    footer(s, 4, total)
    note(s, "讲这一页时强调 scaffold、结构分离和低相似度分层的作用。它们不是额外装饰，而是防止随机划分给出过度乐观结论。")

    # 5
    s = prs.slides.add_slide(blank); add_bg(s, PURPLE)
    title(s, "阶段 3：分子表示与特征生成流程", "同一分子生成多视图表示，供不同专家模型读取。", PURPLE)
    flow_row(s, ["SMILES输入", "RDKit描述符", "Morgan/指纹", "图结构", "任务元信息", "特征缓存"], 0.78, 1.54, 11.80, 0.86, accent=PURPLE)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 从冻结 clean data 读取 SMILES，不从临时文件重新取数。",
         "2. 生成传统描述符、分子指纹、图神经网络输入和可选任务上下文。",
         "3. 对每类表示记录维度、缺失率、失败分子和计算版本。",
         "4. 将表示写入 feature_store，训练阶段只读取缓存。"],
        "质控输出",
        ["输出：fingerprint.npy、descriptor.parquet、graph_cache.pt、feature_manifest.md。",
         "质控：特征生成失败的分子要回写 registry；不能只在训练时报错。",
         "冻结：同一 split 下所有候选专家读取同一表示版本。",
         "写作：多视图表示是候选池来源，不直接等同于性能提升。"],
        accent=PURPLE,
    )
    takeaway(s, "表示生成阶段的目标是保证不同专家的输入一致、可追踪、可复算。", accent=BLUE)
    footer(s, 5, total)
    note(s, "这一页讲多视图表示。要说明指纹、描述符和图结构都是候选专家的输入视角，不能在论文中把多视图本身夸大为已证明的提升来源。")

    # 6
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "阶段 4：候选专家训练流程", "所有模型先作为候选登记，不能直接进入最终结论。", BLUE)
    flow_row(s, ["读取 split", "读取特征", "训练专家", "valid评估", "保存预测", "登记候选"], 0.78, 1.54, 11.80, 0.86, accent=BLUE)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 为每个数据集和 seed 建立训练任务，区分回归与分类。",
         "2. 专家池包括强基线、图模型、指纹模型、表格模型和可控补救模块。",
         "3. 每个专家输出 valid_pred、test_pred 以及训练日志。",
         "4. 候选性能先进入 valid leaderboard，不直接看测试排名。"],
        "质控输出",
        ["输出：candidate_registry.csv、valid_predictions、training_log、model_config.yaml。",
         "质控：失败训练也要记录；超参数范围必须在测试前固定。",
         "冻结：测试预测可以保存，但不能用于选择候选。",
         "写作：候选池越大，越需要 selector 和 rank audit 控制乐观偏差。"],
        accent=BLUE,
    )
    takeaway(s, "训练阶段只产生候选和验证证据，最终策略由下一阶段 selector 决定。", accent=BLUE)
    footer(s, 6, total)
    note(s, "这一页讲模型训练，但不要讲结果图。说明每个模型只是候选，训练结束后进入选择器，测试集不参与决策。")

    # 7
    s = prs.slides.add_slide(blank); add_bg(s, TEAL)
    title(s, "阶段 5：验证集 selector 与门控流程", "selector 是整套流程的中枢：决定保留、拒绝或转入负结果。", TEAL)
    flow_row(s, ["候选登记表", "valid指标", "校准/风险", "门控判断", "策略冻结", "负结果归档"], 0.78, 1.54, 11.80, 0.86, accent=TEAL)
    big_box(s, "选择规则", "回归任务以 RMSE/MAE 为主；分类任务同时检查 ROC-AUC、PR-AUC、Brier、ECE 和固定精度召回。", 0.78, 3.02, 3.70, 1.15, accent=TEAL)
    big_box(s, "门控规则", "只有当验证性能改善且校准、风险或适用域信号没有明显恶化时，候选才进入最终策略。", 4.78, 3.02, 3.70, 1.15, accent=AMBER)
    big_box(s, "拒绝规则", "未通过门控的候选不进入主性能结论，保留为负结果、边界案例或后续验证接口。", 8.78, 3.02, 3.38, 1.15, accent=RED)
    big_box(s, "冻结产物", "输出 selected_strategy.json、rejected_candidates.csv、selector_audit.md，之后才能进入测试集评估。", 1.05, 4.72, 11.10, 0.82, accent=BLUE, body_size=12.2)
    takeaway(s, "selector 不是让结果更好看的工具，而是防止测试集后验选择的治理环节。")
    footer(s, 7, total)
    note(s, "这一页是核心。要强调 selector 只看验证集。接受和拒绝都要有记录，这样审稿人可以看到没有选择性展示。")

    # 8
    s = prs.slides.add_slide(blank); add_bg(s, AMBER)
    title(s, "阶段 6：冻结测试与排名审计流程", "测试集只在策略冻结后使用，用于一次性外部检验。", AMBER)
    flow_row(s, ["读取冻结策略", "运行test预测", "计算主指标", "rank audit", "optimism gap", "锁定主表"], 0.78, 1.54, 11.80, 0.86, accent=AMBER)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 读取 selected_strategy.json，不允许临时替换候选模型。",
         "2. 在 test split 上一次性计算回归或分类主指标。",
         "3. 比较 valid rank 与 test rank，记录是否存在排名错位。",
         "4. 计算 optimism gap，判断验证集优势是否过于乐观。"],
        "质控输出",
        ["输出：test_metrics.csv、rank_audit.csv、optimism_gap.md、main_table_source.xlsx。",
         "质控：若测试表现与验证选择冲突，必须在讨论或限制中解释。",
         "冻结：主表只来自已锁定策略，不能再回到 selector 调整。",
         "写作：测试结果用于验证流程，不用于重新设计流程。"],
        accent=AMBER,
    )
    takeaway(s, "冻结测试流程回答“结果是否由测试集调出来”的关键质疑。", accent=AMBER)
    footer(s, 8, total)
    note(s, "这一页讲测试集如何使用。重点是一次性评估、rank audit 和 optimism gap。")

    # 9
    s = prs.slides.add_slide(blank); add_bg(s, GREEN)
    title(s, "阶段 7：可靠性审计输出流程", "性能之外，还要输出模型何时可靠、何时需要谨慎。", GREEN)
    flow_row(s, ["最终预测", "AD距离", "不确定性", "校准误差", "risk score", "conformal覆盖"], 0.78, 1.54, 11.80, 0.86, accent=GREEN)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 基于训练分布计算适用域距离或相似度信号。",
         "2. 从模型方差、专家分歧或校准误差生成 uncertainty。",
         "3. 构建 risk score，并绘制 risk-coverage 分析。",
         "4. 对分类或回归输出进行校准与保形覆盖评估。"],
        "质控输出",
        ["输出：ad_score.csv、uncertainty.csv、risk_coverage.csv、conformal_report.md。",
         "质控：风险分数只作为审计证据，不作为单独性能胜利。",
         "解释：高风险样本进入失败案例和边界分析。",
         "写作：可靠性模块用于限定适用域，不替代主指标。"],
        accent=GREEN,
    )
    takeaway(s, "可靠性阶段把单个预测扩展为“预测值 + 风险 + 适用域 + 覆盖率”。", accent=GREEN)
    footer(s, 9, total)
    note(s, "这一页讲 AD、uncertainty、risk-coverage、calibration 和 conformal。不要展示结果曲线，只讲这些文件和指标怎么产生。")

    # 10
    s = prs.slides.add_slide(blank); add_bg(s, RED)
    title(s, "阶段 8：OOD、低相似度与失败案例流程", "所有外推风险必须进入同一条错误分析链。", RED)
    flow_row(s, ["scaffold压力", "结构分离", "Tanimoto bin", "MoleculeACE", "失败样本", "风险富集"], 0.78, 1.54, 11.80, 0.86, accent=RED)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 比较 random、scaffold、structure-separated 的协议差异。",
         "2. 按 >0.7、0.5-0.7、<0.5 建立互斥低相似度分层。",
         "3. 对 MoleculeACE 活性悬崖样本单独记录误差和预测差异。",
         "4. 将 ClinTox 假阴性、高风险 ADME、低相似度失败等归档。"],
        "质控输出",
        ["输出：ood_metrics.csv、tanimoto_bins.csv、cliff_pairs.csv、failure_casebook.md。",
         "质控：失败案例不能只选好看的样本，类别要覆盖低相似度和活性悬崖。",
         "解释：风险分数若能富集失败样本，才支持可靠性审计。",
         "写作：OOD 结果用来限定边界，不用来夸大外推能力。"],
        accent=RED,
    )
    takeaway(s, "失败分析是整体流程的保护层：它告诉读者哪些预测不能直接相信。", accent=RED)
    footer(s, 10, total)
    note(s, "这一页讲外推边界。强调低相似度和活性悬崖必须有流程和文件，不只是正文一句话。")

    # 11
    s = prs.slides.add_slide(blank); add_bg(s, PURPLE)
    title(s, "阶段 9：消融、负结果与解释流程", "模块贡献、失败候选和化学解释都要被记录，但不能越界表述。", PURPLE)
    flow_row(s, ["Full策略", "w/o模块", "单专家", "简单平均", "负结果表", "解释审计"], 0.78, 1.54, 11.80, 0.86, accent=PURPLE)
    two_col_steps(
        s,
        "执行步骤",
        ["1. 统一消融矩阵：Full、best single、simple mean、w/o selector、w/o fusion、w/o AD gate、w/o uncertainty。",
         "2. 将未通过验证门控的 rescue head、3D-lite、bRo5 扩展等纳入负结果。",
         "3. 对 motif attribution 和 fragment enrichment 只做关联解释。",
         "4. 检查每个解释是否有支持度、效应量和 p/FDR。"],
        "质控输出",
        ["输出：ablation_matrix.csv、negative_results.md、motif_report.md、fragment_fdr.csv。",
         "质控：解释结果不能写成因果机制证明；负结果不能包装成主性能成功。",
         "主文取舍：主文保留摘要，补充材料放完整长表。",
         "写作：消融回答贡献来源，负结果回答边界。"],
        accent=PURPLE,
    )
    takeaway(s, "消融和负结果让论文更可信：它们证明作者没有只挑成功结果。", accent=BLUE)
    footer(s, 11, total)
    note(s, "这一页讲消融、负结果和解释分析。要强调解释是关联，负结果要完整保留。")

    # 12
    s = prs.slides.add_slide(blank); add_bg(s, BLUE)
    title(s, "阶段 10：论文输出、复现包与答辩复核流程", "最后把所有阶段产物汇总成可审稿、可复核、可答辩的材料。", BLUE)
    flow_row(s, ["主文流程图", "主表/补充表", "数据声明", "代码声明", "QA清单", "答辩问题"], 0.78, 1.54, 11.80, 0.86, accent=BLUE)
    big_box(s, "论文输出", "方法部分按数据、划分、表示、专家、selector、冻结测试、可靠性审计顺序写；结果部分按证据强弱组织。", 0.78, 3.02, 3.70, 1.15, accent=BLUE)
    big_box(s, "复现输出", "打包 dataset registry、split index、feature manifest、model config、metrics source 和 figure scripts。", 4.78, 3.02, 3.70, 1.15, accent=TEAL)
    big_box(s, "答辩输出", "准备四类问题：测试集调参、是否全面优于基线、外推边界、解释是否因果。", 8.78, 3.02, 3.38, 1.15, accent=AMBER)
    big_box(s, "最终检查", "每个主张必须能回指到流程产物；不能核验的内容写为限制或后续工作。", 1.05, 4.72, 11.10, 0.82, accent=RED, body_size=12.2)
    takeaway(s, "整体流程的终点不是一张结果图，而是一套可追溯、可复核、不过度外推的证据链。", accent=BLUE)
    footer(s, 12, total)
    note(s, "最后一页收束全流程。强调所有产物都用于支撑论文和答辩，任何不能核验的主张都要降调。")

    prs.save(OUT_PPTX)


def audit(path: Path):
    prs = Presentation(path)
    issues = []
    notes = 0
    sw, sh = prs.slide_width, prs.slide_height
    text_lengths, shape_counts = [], []
    for si, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                notes += 1
        except Exception:
            pass
        total_text = 0
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > sw + 1000 or shp.top + shp.height > sh + 1000:
                issues.append(f"Slide {si}: shape outside slide bounds")
            if hasattr(shp, "text") and shp.text.strip():
                total_text += len(shp.text.strip())
        if total_text < 120:
            issues.append(f"Slide {si}: low text density")
        if total_text > 900:
            issues.append(f"Slide {si}: high text density ({total_text} chars)")
        text_lengths.append(total_text)
        shape_counts.append(len(slide.shapes))
    with ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        notes_xml = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")]
    if media:
        issues.append(f"Unexpected media files found: {len(media)}")
    return {
        "slides": len(prs.slides),
        "media": len(media),
        "notes": notes,
        "slide_xml": len(slide_xml),
        "note_xml": len(notes_xml),
        "min_text": min(text_lengths),
        "max_text": max(text_lengths),
        "min_shapes": min(shape_counts),
        "max_shapes": max(shape_counts),
        "issues": issues,
    }


def write_report(stats):
    lines = [
        "# 初稿-7 整体流程详解大字版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 按用户要求：不放结果图；只讲完整整体实验流程；字体整体放大；使用原生流程图和步骤表。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个（应为 0，表示未插入结果图）",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 自动检查提示如下：")
        for item in stats["issues"]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 本版未插入任何论文结果图或性能图，全部用 PPT 原生形状重建流程。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页截图预览；已完成 PPTX 包结构、媒体、备注、边界和文本密度检查。",
            "- 12 页用于完整讲清全流程；若必须严格 10 页，可进一步合并阶段 2/3 与阶段 8/9。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    build_deck()
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
