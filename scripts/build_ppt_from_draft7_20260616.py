# -*- coding: utf-8 -*-
from __future__ import annotations

import math
import shutil
from pathlib import Path
from zipfile import ZipFile

from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
WORK = ROOT / "reports" / "ppt_from_draft7_20260616"
FIG_DIR = WORK / "assets" / "figures"
FIG_DIR.mkdir(parents=True, exist_ok=True)


def locate_docx() -> Path:
    candidates = [p for p in (Path.home() / "Desktop").rglob("初稿-7.docx")]
    if not candidates:
        raise FileNotFoundError("Could not locate 初稿-7.docx.")
    return max(candidates, key=lambda p: p.stat().st_mtime)


SRC_DOCX = locate_docx()
OUT_PPTX = SRC_DOCX.parent / "初稿-7_详细汇报PPT.pptx"
QA_MD = SRC_DOCX.parent / "初稿-7_详细汇报PPT_QA报告.md"
OUTLINE_MD = WORK / "ppt_outline_cn.md"
MANIFEST_MD = WORK / "asset_manifest.md"


COLORS = {
    "bg": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(20, 32, 45),
    "muted": RGBColor(86, 100, 115),
    "line": RGBColor(202, 211, 222),
    "blue": RGBColor(35, 93, 155),
    "teal": RGBColor(27, 124, 118),
    "green": RGBColor(61, 145, 97),
    "amber": RGBColor(190, 132, 35),
    "red": RGBColor(176, 67, 67),
    "pale_blue": RGBColor(226, 238, 248),
    "pale_teal": RGBColor(224, 244, 241),
    "pale_amber": RGBColor(249, 238, 214),
    "pale_red": RGBColor(249, 226, 226),
}


def cjk_count(text: str) -> int:
    return sum(1 for c in text if 0x4E00 <= ord(c) <= 0x9FFF)


def extract_figures() -> list[Path]:
    paths: list[Path] = []
    with ZipFile(SRC_DOCX) as z:
        media = sorted(
            [n for n in z.namelist() if n.startswith("word/media/")],
            key=lambda s: int("".join(ch for ch in Path(s).stem if ch.isdigit()) or 0),
        )
        for i, name in enumerate(media, 1):
            ext = Path(name).suffix.lower() or ".png"
            dst = FIG_DIR / f"fig{i:02d}{ext}"
            dst.write_bytes(z.read(name))
            paths.append(dst)
    return paths


def extract_source() -> dict:
    doc = Document(SRC_DOCX)
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    tables = []
    for t in doc.tables:
        tables.append([[cell.text.strip() for cell in row.cells] for row in t.rows])
    return {"paragraphs": paras, "tables": tables, "doc": doc}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(0.7)


def add_textbox(slide, text, x, y, w, h, size=14, color=None, bold=False, align="left", valign="top", margin=0.04):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "mid": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return tb


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_textbox(slide, section, 0.48, 0.22, 1.4, 0.22, size=8.5, color=COLORS["teal"], bold=True)
    add_textbox(slide, title, 0.48, 0.38, 10.2, 0.48, size=22, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 0.50, 0.88, 9.8, 0.30, size=9.5, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.50), Inches(1.16), Inches(1.0), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()


def add_footer(slide, idx, total, source="FZYC-Mol 初稿-7"):
    add_textbox(slide, source, 0.50, 7.20, 3.0, 0.18, size=6.5, color=COLORS["muted"])
    add_textbox(slide, f"{idx}/{total}", 12.25, 7.20, 0.55, 0.18, size=6.5, color=COLORS["muted"], align="right")


def add_bullets(slide, items, x, y, w, h, size=13, color=None, bullet_color=None, gap=0.12):
    top = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(top + 0.08), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color or COLORS["teal"]
        dot.line.fill.background()
        add_textbox(slide, item, x + 0.18, top, w - 0.18, 0.36, size=size, color=color or COLORS["ink"])
        top += 0.42 + gap


def add_box(slide, text, x, y, w, h, fill, stroke=None, size=12, bold=False, align="center"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke or COLORS["line"]
    shp.line.width = Pt(0.8)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}[align]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = COLORS["ink"]
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=None):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color or COLORS["muted"]
    line.line.width = Pt(1.4)
    line.line.end_arrowhead = True
    return line


def add_takeaway(slide, text, color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.72), Inches(12.32), Inches(0.38))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color or COLORS["pale_teal"]
    shp.line.fill.background()
    add_textbox(slide, "Takeaway  " + text, 0.66, 6.80, 11.95, 0.20, size=9.5, color=COLORS["ink"], bold=True)


def add_picture_fit(slide, img_path, x, y, w, h, frame=True):
    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = min(w / (iw / 96), h / (ih / 96))  # rough inch conversion; final uses EMU below
    # Use actual aspect ratio instead of DPI assumptions.
    aspect = iw / ih
    box_aspect = w / h
    if aspect >= box_aspect:
        disp_w = w
        disp_h = w / aspect
    else:
        disp_h = h
        disp_w = h * aspect
    px = x + (w - disp_w) / 2
    py = y + (h - disp_h) / 2
    if frame:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS["white"]
        bg.line.color.rgb = COLORS["line"]
        bg.line.width = Pt(0.7)
    pic = slide.shapes.add_picture(str(img_path), Inches(px), Inches(py), width=Inches(disp_w), height=Inches(disp_h))
    return pic


def add_native_table(slide, data, x, y, w, h, font_size=8.5, header_fill=COLORS["pale_blue"]):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if i == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(font_size if i > 0 else font_size + 0.3)
                    r.font.bold = i == 0
                    r.font.color.rgb = COLORS["ink"]
    return shape


def add_notes(slide, text):
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = text
    except Exception:
        pass


def create_contact_sheet(fig_paths: list[Path]) -> Path:
    thumbs = []
    for path in fig_paths:
        im = Image.open(path).convert("RGB")
        im.thumbnail((360, 240))
        canvas = Image.new("RGB", (380, 280), "white")
        canvas.paste(im, ((380 - im.width) // 2, 10))
        d = ImageDraw.Draw(canvas)
        d.text((10, 252), path.stem, fill=(20, 32, 45))
        thumbs.append(canvas)
    cols = 4
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 380, rows * 280), (245, 247, 250))
    for i, im in enumerate(thumbs):
        sheet.paste(im, ((i % cols) * 380, (i // cols) * 280))
    out = WORK / "assets" / "figure_contact_sheet.jpg"
    out.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out, quality=90)
    return out


def build_deck(src: dict, fig_paths: list[Path]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    total = 25

    def new_slide(title=None, subtitle=None, section=None):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = COLORS["bg"]
        if title:
            add_title(slide, title, subtitle, section)
        slides.append(slide)
        return slide

    # 1 Cover
    s = new_slide()
    add_textbox(s, "FZYC-Mol", 0.62, 0.82, 3.2, 0.48, size=16, color=COLORS["teal"], bold=True)
    add_textbox(s, "面向可靠性审计的验证集治理\n分子性质预测框架", 0.62, 1.38, 8.5, 1.22, size=28, bold=True)
    add_textbox(s, "基于初稿-7生成的详细汇报 PPT", 0.68, 2.70, 4.8, 0.30, size=12, color=COLORS["muted"])
    x0, y0 = 0.78, 4.35
    steps = ["候选登记", "验证门控", "测试冻结", "可靠性审计", "失败边界"]
    for i, st in enumerate(steps):
        add_box(s, st, x0 + i * 2.25, y0, 1.42, 0.52, COLORS["white"], COLORS["teal"], 11, True)
        if i < len(steps) - 1:
            add_arrow(s, x0 + i * 2.25 + 1.48, y0 + 0.26, x0 + (i + 1) * 2.25 - 0.12, y0 + 0.26, COLORS["teal"])
    add_textbox(s, "核心主张：FZYC-Mol 不追求“所有任务普遍最优”，而是让模型选择、适用域边界和失败模式可追踪、可复核。", 0.74, 5.40, 10.8, 0.55, size=13, color=COLORS["ink"])
    add_textbox(s, "汇报结构：问题 → 方法流程 → 评测证据 → 可靠性审计 → 边界与结论", 0.74, 6.18, 10.2, 0.28, size=10, color=COLORS["muted"])
    add_notes(s, "开场说明：这份汇报围绕 FZYC-Mol 的方法学主线展开。不要先讲模型堆叠，而要先讲为什么分子性质预测需要验证集治理和可靠性审计。")

    # 2 Agenda
    s = new_slide("报告路线：从真实风险到可复核模型治理", "每个部分只回答一个问题，避免变成论文逐段复述", "Overview")
    blocks = [
        ("1", "为什么需要", "公开基准高分不等于真实化学决策可靠"),
        ("2", "怎么做", "候选登记、验证门控、测试冻结"),
        ("3", "怎么评", "MoleculeNet、TDC、bRo5、MoleculeACE"),
        ("4", "是否可信", "风险分数、保形预测、排名审计"),
        ("5", "边界在哪", "负结果、失败案例、未来验证"),
    ]
    for i, (num, head, body) in enumerate(blocks):
        y = 1.55 + i * 0.86
        add_textbox(s, num, 0.85, y, 0.35, 0.28, size=14, color=COLORS["teal"], bold=True, align="center")
        add_textbox(s, head, 1.35, y - 0.02, 1.25, 0.28, size=14, bold=True)
        add_textbox(s, body, 2.75, y, 7.9, 0.28, size=13, color=COLORS["muted"])
    add_takeaway(s, "汇报的主线是“过程可信”，不是单纯罗列所有指标。")
    add_notes(s, "这一页交代听众接下来会听到什么。强调 PPT 会把每个流程讲清楚：先讲瓶颈，再讲框架，再讲评价和边界。")

    # 3 Problem
    s = new_slide("痛点：平均指标掩盖真实药物发现风险", "随机划分下的分数不能直接回答“这次预测能不能信”", "Problem")
    add_bullets(s, [
        "新骨架外推：测试分子可能远离训练分布",
        "不平衡毒性：高 ROC-AUC 不等于阳性召回可靠",
        "活性悬崖：相似结构可能出现性质跳变",
        "bRo5 空间：规则五以外分子更容易超出适用域",
    ], 0.8, 1.55, 5.2, 2.3, 13)
    add_box(s, "传统问题\n谁的平均分更高？", 7.15, 1.55, 2.0, 0.78, COLORS["pale_blue"], COLORS["blue"], 12, True)
    add_arrow(s, 9.35, 1.94, 10.35, 1.94, COLORS["muted"])
    add_box(s, "本文问题\n何时可信，何时拒用？", 10.55, 1.55, 2.0, 0.78, COLORS["pale_teal"], COLORS["teal"], 12, True)
    add_box(s, "模型错误的代价不是一个数字偏差，\n而是合成优先级、实验排队和毒性风险判断被改变。", 7.15, 3.25, 5.40, 1.05, COLORS["white"], COLORS["line"], 14, False, "left")
    add_takeaway(s, "FZYC-Mol 把“预测分数”扩展为“预测是否值得相信”的证据链。")
    add_notes(s, "这里要把问题讲得贴近药物发现：模型错误会影响资源投入和风险筛查，不只是一个 benchmark 指标不够好。")

    # 4 Core idea
    s = new_slide("核心想法：把模型开发变成可审计流程", "候选先登记，验证集做决策，测试集只做最终报告", "Core idea")
    stages = [
        ("候选池", "多视图表示\n强基线专家\n补救头"),
        ("验证集治理", "只用验证集\n选择/拒绝/保留"),
        ("冻结测试", "权重、阈值\n规则不再调整"),
        ("可靠性输出", "风险、校准\n适用域、失败案例"),
    ]
    x = 0.78
    for i, (h, b) in enumerate(stages):
        add_box(s, h + "\n" + b, x + i * 3.0, 1.75, 2.15, 1.15, COLORS["white"], COLORS["teal" if i == 1 else "line"], 12, True)
        if i < 3:
            add_arrow(s, x + i * 3.0 + 2.22, 2.32, x + (i + 1) * 3.0 - 0.18, 2.32, COLORS["teal"])
    add_textbox(s, "治理规则", 0.82, 3.65, 1.1, 0.25, size=12, color=COLORS["blue"], bold=True)
    add_bullets(s, [
        "测试集不参与候选生成、权重拟合和阈值设定",
        "未通过验证集门控的模块作为负结果或后续接口保留",
        "最终性能与选择器风险同时报告",
    ], 0.95, 4.02, 7.8, 1.45, 12.5, bullet_color=COLORS["blue"])
    add_takeaway(s, "方法创新点在于选择过程，而不是盲目堆叠更多模型。")
    add_notes(s, "讲清楚这不是一个单一神经网络结构，而是验证集治理框架。听众要记住三个词：候选登记、测试冻结、可靠性输出。")

    # 5 Figure 1
    s = new_slide("模型结构：多源表示进入同一候选证据池", "图1展示从分子表示到专家预测矩阵和证据输出的结构", "Method")
    add_picture_fit(s, fig_paths[0], 0.55, 1.30, 8.8, 5.10)
    add_bullets(s, ["多视图表示覆盖不同化学层次", "专家预测矩阵支撑融合与选择", "输出不只给分数，还给风险证据"], 9.65, 1.55, 3.0, 1.8, 11.5)
    add_textbox(s, "Source: 图1", 0.65, 6.42, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "FZYC-Mol 的结构核心是“多证据候选池”。")
    add_notes(s, "逐步讲图：左侧是输入表示，中间是专家矩阵与验证选择器，右侧是性能、适用域和风险证据输出。")

    # 6 Figure 2 workflow
    s = new_slide("整体工作流：从数据划分到可靠性报告", "图2是汇报后续所有流程的总地图", "Workflow")
    add_picture_fit(s, fig_paths[1], 0.60, 1.25, 11.95, 5.25)
    add_textbox(s, "读图顺序：数据划分 → 多视图表示 → 候选专家 → 验证选择 → 风险/校准/失败输出", 0.75, 6.52, 10.8, 0.20, size=8.5, color=COLORS["muted"])
    add_takeaway(s, "每一步都有清楚的输入、决策依据和输出证据。")
    add_notes(s, "这一页作为流程总览，后面按模块拆解。提醒听众：测试集只在策略冻结后出现。")

    # 7 Data design
    s = new_slide("评测设计：覆盖插值、外推、外部 ADMET 和压力测试", "数据集不只是越多越好，而是对应不同应用难度", "Evaluation")
    table1 = src["tables"][0]
    add_native_table(s, table1, 0.55, 1.28, 12.10, 4.20, font_size=7.2)
    add_bullets(s, ["MoleculeNet：主性能面板", "TDC ADMET：外部终点迁移", "MoleculeACE/bRo5：活性悬崖与化学空间压力"], 0.90, 5.78, 10.3, 0.72, 10.5, bullet_color=COLORS["blue"])
    add_takeaway(s, "不同数据模块对应不同证据问题，不能混成一个平均指标。")
    add_notes(s, "逐行说明：MoleculeNet 是主面板，TDC 是外部 ADMET，MoleculeACE 看活性悬崖，bRo5 只作为公共压力测试，不能说成独立盲测。")

    # 8 Multi-view representation
    s = new_slide("多视图表示：不是堆特征，而是覆盖化学层次", "每类表示承担不同证据角色", "Method")
    views = [
        ("分子图", "原子-键连接\nGNN / D-MPNN"),
        ("指纹", "Morgan / MACCS\n局部子结构"),
        ("RDKit 描述符", "二维拓扑\n物化属性摘要"),
        ("片段/骨架", "BRICS / Murcko\n解释与富集"),
        ("冻结表征", "ChemBERTa / MolT5\n上下文表示"),
    ]
    for i, (h, b) in enumerate(views):
        add_box(s, h + "\n" + b, 0.62 + i * 2.45, 1.55, 1.95, 1.05, COLORS["white"], COLORS["line"], 10.5, True)
        if i < 4:
            add_arrow(s, 0.62 + i * 2.45 + 1.98, 2.08, 0.62 + (i + 1) * 2.45 - 0.10, 2.08, COLORS["line"])
    add_box(s, "进入同一候选池\n由验证集决定是否保留", 4.60, 3.40, 4.15, 0.82, COLORS["pale_teal"], COLORS["teal"], 15, True)
    add_bullets(s, ["复杂模型不预设更优", "同一终点可接受简单稳定专家", "解释模块只作关联证据"], 1.0, 5.10, 9.4, 0.92, 12.5)
    add_takeaway(s, "表示选择本身也是可审计经验问题。")
    add_notes(s, "这一页讲方法设计的合理性。五类表示分别覆盖结构、指纹、描述符、片段和预训练语义。")

    # 9 Candidate registration
    s = new_slide("候选登记：所有策略先登记，再比较", "候选池大时，流程约束比单次分数更重要", "Governance")
    categories = [
        ("单专家", "RF / XGBoost / CatBoost\nChemprop / GNN"),
        ("目标变换", "log1p\n分位数正态化\n截尾目标"),
        ("融合策略", "Top-K 均值\n堆叠\n不确定性加权"),
        ("门控/补救", "适用域门控\n定向补救头\n平局规则"),
    ]
    for i, (h, b) in enumerate(categories):
        add_box(s, h + "\n" + b, 0.82 + i * 3.05, 1.55, 2.25, 1.36, [COLORS["pale_blue"], COLORS["white"], COLORS["pale_teal"], COLORS["pale_amber"]][i], COLORS["line"], 11.2, True)
    add_picture_fit(s, fig_paths[2], 1.05, 3.35, 10.95, 2.60)
    add_textbox(s, "Source: 图3", 1.10, 6.03, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "候选扩展必须被验证集治理约束，否则容易形成事后选择偏差。")
    add_notes(s, "先讲四类候选，再用图3说明候选登记、验证选择、外部测试和负结果记录之间的关系。")

    # 10 Validation governance
    s = new_slide("验证集治理：六步冻结流程防止测试集被开发化", "这是 FZYC-Mol 方法学边界的核心", "Governance")
    flow = [
        "固定任务\n指标/种子",
        "训练单专家\n仅用训练集",
        "验证预测矩阵\n构造融合",
        "验证门控\n接受/拒绝",
        "冻结权重\n阈值/规则",
        "一次性测试\n记录负结果",
    ]
    for i, txt in enumerate(flow):
        x = 0.62 + (i % 3) * 4.05
        y = 1.45 + (i // 3) * 1.60
        add_box(s, txt, x, y, 2.55, 0.82, COLORS["white"], COLORS["teal" if i in [3, 4, 5] else "line"], 12, True)
        if i % 3 < 2:
            add_arrow(s, x + 2.62, y + 0.41, x + 3.88, y + 0.41, COLORS["teal"])
    add_arrow(s, 9.25, 1.86, 9.25, 2.68, COLORS["teal"])
    add_arrow(s, 9.25, 3.46, 6.90, 3.46, COLORS["teal"])
    add_arrow(s, 5.15, 3.46, 2.88, 3.46, COLORS["teal"])
    add_bullets(s, ["未通过门控：负结果或后续接口", "通过门控：进入最终保留策略", "测试集：只做最终一次性报告"], 1.05, 5.20, 7.5, 0.92, 12.5, bullet_color=COLORS["blue"])
    add_takeaway(s, "验证集不是小测试集，而是唯一的策略选择依据。")
    add_notes(s, "逐步讲六个步骤。重点强调：测试集不能参与权重拟合、阈值设定和平局规则。")

    # 11 Reliability modules
    s = new_slide("可靠性模块：从单点预测走向决策证据", "风险、校准、保形预测和解释模块回答“何时可信”", "Reliability")
    modules = [
        ("适用域", "最近邻相似度\n骨架距离\n低相似度分层"),
        ("风险分数", "模型分歧\n预测偏差\n重构误差"),
        ("校准/保形", "Brier / ECE\n80/90/95% 覆盖"),
        ("化学解释", "基序归因\n片段富集\n最近邻复核"),
    ]
    for i, (h, b) in enumerate(modules):
        add_box(s, h + "\n" + b, 0.75 + i * 3.05, 1.55, 2.25, 1.30, COLORS["white"], COLORS["line"], 11.4, True)
    add_box(s, "输出形式\n点预测 + 风险分位 + 适用域证据 + 拒用理由", 2.25, 3.75, 8.85, 0.78, COLORS["pale_teal"], COLORS["teal"], 15, True)
    add_bullets(s, ["风险分数是复核提示，不是自动纠错机制", "基序/片段解释是关联证据，不是因果机制", "保形预测提供覆盖率意义，但不替代校准"], 1.15, 5.25, 9.8, 0.92, 12.2)
    add_takeaway(s, "可靠性输出把模型结果转化为可讨论的决策卡。")
    add_notes(s, "这一页讲四个可靠性模块。每个模块都要说明它做什么，以及它不能证明什么。")

    # 12 Main MoleculeNet table
    s = new_slide("主结果：FZYC-Mol 在部分 MoleculeNet 终点形成选择性增益", "不要讲成所有任务统一提升", "Main results")
    add_native_table(s, src["tables"][1], 0.55, 1.30, 12.10, 3.25, font_size=7.1, header_fill=COLORS["pale_teal"])
    add_bullets(s, ["ESOL / BACE：验证选择与测试观测最优基本一致", "Lipophilicity：补救头被验证集接受", "FreeSolv：保留为物理相互作用边界案例"], 0.85, 4.95, 10.8, 1.05, 11.5)
    add_takeaway(s, "结果模式是选择性接受，而不是全面追求复杂模型。")
    add_notes(s, "讲表2时按回归、分类分开。强调 FreeSolv 是边界案例，不要把所有终点都包装成胜利。")

    # 13 Fig4 ranking
    s = new_slide("终点内模型家族排名揭示任务依赖性", "不同终点偏好的专家家族并不相同", "Main results")
    add_picture_fit(s, fig_paths[3], 0.55, 1.25, 9.0, 5.25)
    add_bullets(s, ["稳定候选由验证集识别", "复杂模型并非每个终点都占优", "排名差异提示需要候选治理"], 9.85, 1.65, 2.65, 1.25, 11.5)
    add_textbox(s, "Source: 图4", 0.65, 6.47, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "候选池的价值在于让终点差异被显式比较。")
    add_notes(s, "用这页说明为什么要多专家候选池。不同任务可能由不同模型家族主导。")

    # 14 Fig5 main comparison
    s = new_slide("主性能比较显示：增益集中在少数可解释终点", "结果需要结合验证门控和边界案例读", "Main results")
    add_picture_fit(s, fig_paths[4], 0.62, 1.25, 11.75, 4.85)
    add_textbox(s, "Source: 图5", 0.72, 6.15, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "性能图不是独立结论，必须和验证集治理一起解释。")
    add_notes(s, "讲图5时不要只读数字。说清楚每个最终策略进入结果的前提是验证集接受。")

    # 15 Rescue head
    s = new_slide("定向补救：只让验证集接受的补救头进入最终策略", "补救不是大规模重训练，而是受控增强", "Rescue")
    add_picture_fit(s, fig_paths[5], 0.55, 1.25, 7.95, 4.95)
    add_box(s, "进入最终策略\nLipophilicity", 9.05, 1.58, 2.75, 0.78, COLORS["pale_teal"], COLORS["teal"], 14, True)
    add_box(s, "未通过门控\n其他终点保留原策略", 9.05, 2.78, 2.75, 0.78, COLORS["white"], COLORS["line"], 13, True)
    add_bullets(s, ["避免测试后补救", "负结果进入审计", "减少选择性报告"], 9.10, 4.25, 2.8, 0.92, 11)
    add_textbox(s, "Source: 图6", 0.65, 6.24, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "补救头只有在验证集证据支持时才被接受。")
    add_notes(s, "这页要讲清楚补救头的门控逻辑。不是看到测试好就加入，而是验证集先接受。")

    # 16 Fusion gain
    s = new_slide("多方法融合：BBBP 与 ClinTox 出现选择性增益", "不平衡毒性任务需同时看 ROC-AUC、PR-AUC、校准和召回", "Fusion")
    add_picture_fit(s, fig_paths[6], 0.55, 1.25, 8.65, 5.05)
    add_bullets(s, ["BBBP：融合候选进入最终保留", "ClinTox：ROC-AUC 高但需警惕假阴性", "增益幅度小，也必须经验证门控"], 9.50, 1.65, 2.95, 1.35, 11.2)
    add_textbox(s, "Source: 图7", 0.65, 6.35, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "分类任务的可靠性不能只靠 ROC-AUC 判断。")
    add_notes(s, "突出 ClinTox 的解释边界：样本不平衡时，高 ROC-AUC 也不等于筛查可靠。")

    # 17 External ADMET
    s = new_slide("外部 TDC ADMET：选择性增益支持审计框架定位", "22 个终点最终为 win/tie/loss = 5/17/0", "External validation")
    add_native_table(s, src["tables"][2], 0.70, 1.35, 5.7, 2.45, font_size=8.2, header_fill=COLORS["pale_blue"])
    add_picture_fit(s, fig_paths[7], 6.75, 1.25, 5.75, 2.35)
    add_picture_fit(s, fig_paths[8], 6.75, 3.78, 5.75, 2.35)
    add_textbox(s, "Source: 表3、图8、图9", 0.75, 6.23, 2.1, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "多数终点保持原策略，说明框架不会强制复杂化模型。")
    add_notes(s, "这一页把外部 ADMET 的表格和两张图放在一起。讲 win/tie/loss = 5/17/0，重点是选择性而非普遍提升。")

    # 18 Reliability and risk coverage
    s = new_slide("风险分数与 risk-coverage：分类错误更容易被识别", "可靠性输出定位为复核提示，不是自动拒用规则", "Reliability")
    add_picture_fit(s, fig_paths[9], 0.55, 1.25, 7.0, 5.15)
    add_native_table(s, src["tables"][3], 7.80, 1.35, 4.65, 3.15, font_size=7.6, header_fill=COLORS["pale_teal"])
    add_bullets(s, ["分类错误 AUROC 中位 0.788", "回归高误差 AUROC 中位 0.652", "回归任务仍需更强物理或局部结构证据"], 7.95, 4.86, 4.1, 0.85, 10.5)
    add_textbox(s, "Source: 图10、表4", 0.65, 6.45, 1.6, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "风险分数能排序复核优先级，但不能替代实验判断。")
    add_notes(s, "解释 risk-coverage：高风险样本优先复核后，保留样本的错误风险可能下降。")

    # 19 Conformal + selector audit
    s = new_slide("保形预测与排名审计：验证集有用但不完美", "覆盖率接近目标，选择器仍有排序不稳定", "Audit")
    metrics = [
        ("分类覆盖率", "0.814 / 0.918 / 0.956", "对应 80/90/95%"),
        ("回归覆盖率", "0.823 / 0.925 / 0.962", "对应 80/90/95%"),
        ("Spearman", "中位 0.667", "验证-测试中等相关"),
        ("Top-3 / Top-1", "0.295 / 0.135", "第一名不等于最优"),
    ]
    for i, (h, val, cap) in enumerate(metrics):
        x = 0.85 + (i % 2) * 5.85
        y = 1.55 + (i // 2) * 1.75
        add_box(s, h + "\n" + val + "\n" + cap, x, y, 4.45, 1.10, COLORS["white"], COLORS["line"], 13, True)
    add_bullets(s, ["保形预测用于输出不确定性边界", "排名审计用于检查选择器本身", "nested validation 是候选池过拟合的补充诊断"], 1.10, 5.30, 9.7, 0.85, 12)
    add_takeaway(s, "可靠性证据不仅评价模型，也评价选择过程。")
    add_notes(s, "这一页讲两个审计层面：保形预测看输出覆盖率，排名审计看验证集选择是否稳定。")

    # 20 OOD and activity cliff
    s = new_slide("低相似度与活性悬崖：平均指标外的样本级风险", "结构相似不保证性质相近，远离训练分布也更容易失败", "Stress test")
    add_picture_fit(s, fig_paths[10], 0.55, 1.22, 8.15, 5.25)
    add_bullets(s, ["相似度下降时性能通常下降", "活性悬崖暴露局部标签跳变", "FZYC-Mol 标记高风险区域而非自动修正"], 9.05, 1.68, 3.15, 1.35, 11.5)
    add_textbox(s, "Source: 图11", 0.65, 6.50, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "OOD 与活性悬崖决定了模型边界，而不是附属分析。")
    add_notes(s, "讲清楚：这部分不是为了证明模型消除了活性悬崖，而是证明框架能把这类风险显性化。")

    # 21 Ablation and negative results
    s = new_slide("消融与负结果：复杂模块不一定稳定带来收益", "负结果是模型治理流程的一部分", "Ablation")
    add_native_table(s, src["tables"][4], 0.70, 1.28, 5.85, 3.10, font_size=8.0, header_fill=COLORS["pale_amber"])
    add_picture_fit(s, fig_paths[11], 6.85, 1.25, 5.55, 3.40)
    add_bullets(s, ["FreeSolv：物理相互作用边界", "bRo5：公共压力测试，非独立盲测", "基序富集：关联解释，不是因果证明"], 0.92, 5.05, 9.6, 0.85, 11.2)
    add_textbox(s, "Source: 表5、图12", 0.75, 6.28, 1.5, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "可信方法需要报告失败模块，而不是只展示胜利样本。")
    add_notes(s, "强调负结果的价值。表5中的 bRo5 行非常关键：不能把公共压力测试讲成独立实验盲测。")

    # 22 Chemical interpretation
    s = new_slide("化学解释：基序、片段与最近邻只支持关联证据", "解释模块帮助复核错误样本，但不等于机制证明", "Interpretability")
    add_picture_fit(s, fig_paths[12], 0.55, 1.25, 8.50, 5.10)
    add_bullets(s, ["基序归因：定位模型关注结构", "片段富集：需报告效应量和 FDR", "最近邻案例：辅助判断外推风险"], 9.35, 1.65, 2.95, 1.35, 11.5)
    add_textbox(s, "Source: 图13", 0.65, 6.38, 1.0, 0.15, size=6.5, color=COLORS["muted"])
    add_takeaway(s, "解释越接近化学结论，越需要统计和案例双重约束。")
    add_notes(s, "讲解时避免说“发现了机制”。应该说：这些结果为样本复核提供关联性证据。")

    # 23 How to explain the whole process
    s = new_slide("从预测到决策卡：FZYC-Mol 的最终输出应这样读", "模型输出不只是一个性质值，而是一组可复核证据", "Use")
    rows = [
        ("点预测", "模型给出的性质或类别概率"),
        ("适用域", "最近邻相似度、骨架距离、低相似度分层"),
        ("风险分位", "错误或高误差复核优先级"),
        ("校准/覆盖", "Brier、ECE、保形区间或集合"),
        ("解释证据", "基序、片段、最近邻案例"),
        ("拒用理由", "超出适用域或证据不足时明确标记"),
    ]
    for i, (h, b) in enumerate(rows):
        x = 0.85 + (i % 2) * 5.75
        y = 1.35 + (i // 2) * 1.25
        add_box(s, h + "\n" + b, x, y, 4.65, 0.82, COLORS["white"], COLORS["line"], 11.5, True, "left")
    add_takeaway(s, "可复核输出能帮助药物化学家决定“信、复核、拒用”。")
    add_notes(s, "这页是应用视角。把前面模块综合成一个实际使用时的决策卡。")

    # 24 Limitations
    s = new_slide("边界与投稿前风险：哪些话不能说过头", "Nature 风格更重视证据边界，而不是绝对化结论", "Limitations")
    add_box(s, "不能说", 0.85, 1.35, 1.45, 0.42, COLORS["pale_red"], COLORS["red"], 12, True)
    add_box(s, "应当说", 6.95, 1.35, 1.45, 0.42, COLORS["pale_teal"], COLORS["teal"], 12, True)
    pairs = [
        ("所有 ADMET 终点普遍提升", "部分终点选择性增益"),
        ("bRo5 是独立盲测", "bRo5 是公共压力测试"),
        ("风险分数自动识别全部错误", "风险分数提供复核排序"),
        ("基序解释证明因果机制", "基序解释提供关联证据"),
    ]
    for i, (bad, good) in enumerate(pairs):
        y = 2.05 + i * 0.78
        add_textbox(s, bad, 0.95, y, 4.65, 0.30, size=12, color=COLORS["red"])
        add_arrow(s, 5.75, y + 0.15, 6.55, y + 0.15, COLORS["muted"])
        add_textbox(s, good, 7.05, y, 4.8, 0.30, size=12, color=COLORS["teal"], bold=True)
    add_takeaway(s, "边界写清楚，反而能提高审稿人对方法可靠性的信任。")
    add_notes(s, "这一页适合结尾前强调学术诚信和投稿风险。讲的时候要直说：这不是削弱论文，而是保护论文。")

    # 25 Conclusion
    s = new_slide("总结：FZYC-Mol 的贡献是一套可靠性报告流程", "性能、选择过程、适用域和失败模式必须一起报告", "Conclusion")
    add_textbox(s, "一句话结论", 0.80, 1.35, 2.0, 0.28, size=12, color=COLORS["teal"], bold=True)
    add_textbox(s, "FZYC-Mol 通过验证集治理和适用域审计，让分子性质预测从“分数比较”走向“可复核决策证据”。", 0.80, 1.78, 10.8, 0.70, size=20, bold=True)
    add_bullets(s, [
        "方法层面：候选登记、验证门控、测试冻结",
        "证据层面：MoleculeNet、TDC、bRo5、MoleculeACE",
        "可靠性层面：风险分数、保形预测、排名审计",
        "边界层面：负结果、失败案例、非盲测说明",
    ], 1.00, 3.20, 9.8, 1.70, 13, bullet_color=COLORS["teal"])
    add_box(s, "最终提醒\n不要把选择性增益讲成普遍最优", 8.65, 5.25, 3.2, 0.78, COLORS["pale_amber"], COLORS["amber"], 13, True)
    add_notes(s, "最后收束：贡献是可靠性流程，而不是一个万能模型。适合回答审稿人或导师：本文到底新在哪里，边界在哪里。")

    for idx, slide in enumerate(slides, 1):
        add_footer(slide, idx, total)

    return prs, slides


def audit_pptx(pptx_path: Path, planned_notes: bool = True) -> list[str]:
    prs = Presentation(pptx_path)
    issues: list[str] = []
    media_count = 0
    with ZipFile(pptx_path) as z:
        media_count = sum(1 for n in z.namelist() if n.startswith("ppt/media/"))
    note_count = 0
    for i, slide in enumerate(prs.slides, 1):
        try:
            if slide.notes_slide.notes_text_frame.text.strip():
                note_count += 1
        except Exception:
            pass
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width + 1 or shape.top + shape.height > prs.slide_height + 1:
                issues.append(f"high: slide {i} shape out of bounds")
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if len(text) > 170 and shape.width < Inches(4.0):
                    issues.append(f"medium: slide {i} narrow text box may overflow ({len(text)} chars)")
                if any(tok in text.lower() for tok in ["lorem", "xxxx"]):
                    issues.append(f"high: slide {i} placeholder text remains")
    if planned_notes and note_count < len(prs.slides):
        issues.append(f"medium: only {note_count}/{len(prs.slides)} slides have notes")
    return [f"slides={len(prs.slides)}", f"media={media_count}", f"notes={note_count}"] + issues


def write_outline_and_manifest(fig_paths: list[Path]) -> None:
    OUTLINE_MD.write_text(
        "# 初稿-7 详细汇报 PPT 结构\n\n"
        "类型：methods / AI / tool / algorithm。\n\n"
        "核心论点：FZYC-Mol 通过候选登记、验证集治理、测试集冻结、适用域审计和负结果记录，使分子性质预测从单一分数比较转向可复核的可靠性报告。\n\n"
        "结构：问题背景 -> 方法框架 -> 数据与候选登记 -> 验证集治理 -> 主结果 -> 外部验证 -> 可靠性审计 -> 压力测试 -> 负结果 -> 边界与总结。\n",
        encoding="utf-8",
    )
    lines = ["# Asset manifest\n"]
    captions = [
        "FZYC-Mol 模型结构",
        "FZYC-Mol 整体工作流",
        "强基线与选择器治理",
        "MoleculeNet 终点内模型家族排名",
        "MoleculeNet 主性能比较",
        "定向补救门控",
        "MoleculeNet 多方法融合门控",
        "PyTDC 官方 ADMET scaffold 划分压力",
        "外部 TDC 融合门控",
        "risk-coverage 曲线",
        "随机划分到 scaffold/结构分离压力测试",
        "最终策略整合与候选选择摘要",
        "基序归因与片段富集",
    ]
    for i, p in enumerate(fig_paths, 1):
        with Image.open(p) as im:
            size = f"{im.width}x{im.height}"
        lines.append(f"- `{p.name}` | 图{i}: {captions[i-1]} | {size} | extracted from `{SRC_DOCX.name}` | used in PPT.\n")
    MANIFEST_MD.write_text("\n".join(lines), encoding="utf-8")


def write_qa(report_lines: list[str], defects_before: list[str], defects_after: list[str], contact_sheet: Path):
    content = [
        "# 初稿-7 详细汇报 PPT QA 报告",
        "",
        f"- 源文档：`{SRC_DOCX}`",
        f"- 输出 PPT：`{OUT_PPTX}`",
        f"- 图像资产：`{FIG_DIR}`",
        f"- 图像 contact sheet：`{contact_sheet}`",
        "",
        "## Verification",
    ]
    content.extend([f"- {line}" for line in report_lines])
    content.extend(["", "## Self-review defects before revision/check", *[f"- {x}" for x in defects_before]])
    content.extend(["", "## Remaining issues after verification", *[f"- {x}" for x in defects_after if not x.startswith(('slides=', 'media=', 'notes='))]])
    if not any(not x.startswith(("slides=", "media=", "notes=")) for x in defects_after):
        content.append("- 未发现高/中严重度结构性问题。")
    content.extend([
        "",
        "## Notes",
        "- 当前环境没有可用的 headless PowerPoint/LibreOffice 渲染器，因此未输出逐页渲染预览；已使用 python-pptx 进行结构、边界、媒体和备注检查。",
        "- PPT 中的原始科学图像来自 docx 内嵌图片，未重新绘制或改动数据。",
        "- `ppt-master` 插件当前不可用，本 PPT 使用 nature-paper2ppt 工作流和 python-pptx 生成。",
    ])
    QA_MD.write_text("\n".join(content), encoding="utf-8")


def main():
    src = extract_source()
    fig_paths = extract_figures()
    contact_sheet = create_contact_sheet(fig_paths)
    write_outline_and_manifest(fig_paths)

    prs, slides = build_deck(src, fig_paths)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)

    defects_first = audit_pptx(OUT_PPTX)

    # Reopen/save to validate package and normalize relationships.
    prs2 = Presentation(OUT_PPTX)
    prs2.save(OUT_PPTX)
    defects_after = audit_pptx(OUT_PPTX)

    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    report_lines = [
        "PPTX 已生成并可由 python-pptx 重新打开。",
        f"幻灯片数量：{len(Presentation(OUT_PPTX).slides)}。",
        f"嵌入图像资产数量：{len(fig_paths)} 张源图已提取；PPT 媒体数见结构检查。",
        "每页均写入中文讲者备注。",
        "布局采用问题-方案-流程-证据-边界的 methods arc。",
    ]
    write_qa(report_lines, defects_first, defects_after, contact_sheet)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print("\\n".join(defects_after))


if __name__ == "__main__":
    main()
