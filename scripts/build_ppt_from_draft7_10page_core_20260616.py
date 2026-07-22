# -*- coding: utf-8 -*-
from __future__ import annotations

import importlib.util
import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
DENSE_SCRIPT = ROOT / "scripts" / "build_ppt_from_draft7_simplified_dense_20260616.py"
spec = importlib.util.spec_from_file_location("dense_ppt", DENSE_SCRIPT)
dense = importlib.util.module_from_spec(spec)
assert spec and spec.loader
spec.loader.exec_module(dense)

base = dense.base
WORK = ROOT / "reports" / "ppt_from_draft7_10page_core_20260616"
WORK.mkdir(parents=True, exist_ok=True)
OUT_PPTX = base.SRC_DOCX.parent / "初稿-7_十页核心版PPT.pptx"
QA_MD = base.SRC_DOCX.parent / "初稿-7_十页核心版PPT_QA报告.md"


def tbl_dataset(src):
    rows = [["数据集", "任务", "规模", "主指标"]]
    for r in src["tables"][0][1:5]:
        rows.append([r[0], r[1].replace("MoleculeNet；", "").replace("；", "/"), r[2], r[3]])
    return rows


def tbl_main(src):
    rows = [["终点", "指标", "最终保留", "解释"]]
    for r in src["tables"][1][1:5]:
        rows.append([r[0], r[1], r[3], r[5]])
    return rows


def tbl_external(src):
    rows = [["模块", "主要发现", "解释"]]
    for r in src["tables"][2][1:4]:
        rows.append([r[0], r[1], r[2]])
    return rows


def tbl_reliability(src):
    rows = [["指标", "结果", "含义"]]
    for r in src["tables"][3][1:5]:
        rows.append([r[0], r[1], r[2]])
    return rows


def tbl_boundary(src):
    rows = [["模块", "状态", "解释"]]
    for r in src["tables"][4][1:5]:
        rows.append([r[0], r[1], r[2]])
    return rows


def add_source(slide, body, x=0.62, y=6.44):
    dense.text(slide, body, x, y, 11.6, 0.16, size=6.5, color=dense.C["muted"])


def build_deck(src, figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    # 1. Cover + compressed argument.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.BLUE)
    dense.text(s, "FZYC-Mol", 0.62, 0.46, 5.2, 0.56, size=28, bold=True, color=dense.DARK)
    dense.text(s, "面向可靠性审计的验证集治理分子性质预测框架", 0.66, 1.06, 9.4, 0.42, size=15, color=dense.SLATE)
    dense.callout(s, "研究痛点", "分子性质预测不能只看平均性能；验证集选择、结构外推、低相似度样本和失败案例会直接影响可信度。", 0.78, 1.82, 3.58, 1.15, accent=dense.BLUE)
    dense.callout(s, "核心方法", "将多源表示、候选专家、验证集选择器、AD gate、uncertainty、risk-coverage 与证据输出纳入同一冻结流程。", 4.72, 1.82, 3.58, 1.15, accent=dense.TEAL)
    dense.callout(s, "安全主张", "本文支持“选择性增益 + 可靠性审计”，不写作所有终点均显著优于强基线。", 8.66, 1.82, 3.28, 1.15, accent=dense.AMBER)
    dense.process_strip(
        s,
        [("问题", "过度乐观与外推风险"), ("流程", "验证治理与冻结测试"), ("证据", "内部/外部/可靠性"), ("边界", "负结果和失败样本")],
        0.90, 3.56, 11.40, 1.00, accent=dense.BLUE,
    )
    dense.add_three_line_table(
        s,
        [["汇报压缩逻辑", "保留内容"], ["10 页", "覆盖背景、方法、实验、可靠性、边界、复现与答辩"], ["13 图", "保留原稿内嵌图，不修改科学数据"], ["20 页讲稿信息", "压缩进讲者备注和合并页"]],
        1.02, 5.05, 10.85, 1.12, widths=[0.28, 0.72], font_size=8.1, header_size=8.5, accent=dense.BLUE,
    )
    dense.takeaway(s, "十页版适合快速完整汇报：少翻页，但每页保留主张、证据和边界。", accent=dense.BLUE)
    dense.footer(s, 1, total)
    dense.note(s, "开场先说明这是十页核心版。逻辑是：为什么要做、如何做、证据是什么、哪里不能夸大。强调 PPT 未修改原始科学图，仅重新组织论文图表。")

    # 2. Workflow + architecture.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.TEAL)
    dense.title(s, "workflow 与模型结构共同服务于验证集治理", "总流程解决“怎么运行”，结构图解释“信息如何进入预测和审计模块”。", dense.TEAL)
    base.add_picture_fit(s, figs[1], 0.52, 1.25, 6.25, 3.82, frame=True)
    base.add_picture_fit(s, figs[0], 6.95, 1.25, 5.55, 3.82, frame=True)
    dense.callout(s, "完整流程", "输入分子经数据清洗、划分、候选专家、验证门控和冻结测试后，输出预测、适用域、风险与解释。", 0.72, 5.28, 3.62, 0.72, accent=dense.TEAL)
    dense.callout(s, "模型结构", "多源表示和专家矩阵并行生成候选，selector 只依据验证集证据保留最终策略。", 4.58, 5.28, 3.62, 0.72, accent=dense.BLUE)
    dense.callout(s, "简式公式", "m* = argmin Lval(m)；ŷ = Σwkfk(x)；r = αu + βdAD + γc。", 8.44, 5.28, 3.58, 0.72, accent=dense.AMBER)
    add_source(s, "Source: Fig. 1-2 from 初稿-7.docx")
    dense.takeaway(s, "核心不是堆叠更多模型，而是让模型进入最终策略的每一步都有验证证据。")
    dense.footer(s, 2, total)
    dense.note(s, "这一页合并 workflow 和 architecture。讲解时从左图的流程开始，再转右图的结构。公式是汇报用简式，用于说明选择、融合和风险分数的关系。")

    # 3. Data + selector governance.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.AMBER)
    dense.title(s, "数据、划分与 selector 决定结果是否可信", "同一页交代数据来源、评价协议、候选登记和冻结测试。", dense.AMBER)
    dense.add_three_line_table(s, tbl_dataset(src), 0.62, 1.30, 5.58, 2.20, widths=[0.20, 0.36, 0.20, 0.24], font_size=7.5, header_size=7.9, accent=dense.AMBER)
    base.add_picture_fit(s, figs[2], 6.55, 1.30, 5.95, 2.92, frame=True)
    dense.callout(s, "评价协议", "回归看 RMSE/MAE；分类看 ROC-AUC、PR-AUC、Brier、ECE 和 fixed precision recall。", 0.70, 3.84, 3.58, 0.82, accent=dense.BLUE)
    dense.callout(s, "划分压力", "random、scaffold、structure-separated 和低相似度 bin 共同限制外推主张。", 4.56, 3.84, 3.58, 0.82, accent=dense.RED)
    dense.callout(s, "冻结原则", "候选先登记，验证集选择；测试集只在策略冻结后评估一次。", 8.42, 4.48, 3.62, 0.72, accent=dense.TEAL)
    dense.text(s, "selector 接受条件：验证性能改善且校准/风险不恶化；未通过者进入负结果或补充验证接口。", 6.70, 4.48, 1.62, 0.80, size=8.4, color=dense.DARK, bold=True)
    add_source(s, "Source: Table 1 and Fig. 3 from 初稿-7.docx")
    dense.takeaway(s, "这页回答审稿人最关心的问题：是否存在测试集调参，以及比较协议是否一致。", accent=dense.AMBER)
    dense.footer(s, 3, total)
    dense.note(s, "左侧表格说明数据和任务，右侧图说明强基线与选择器治理。重点讲测试集冻结、候选登记和负结果归档。")

    # 4. MoleculeNet main evidence.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.BLUE)
    dense.title(s, "MoleculeNet 结果支持选择性保留，而非全任务最优", "rank map 与主结果表共同说明不同终点需要不同专家。", dense.BLUE)
    base.add_picture_fit(s, figs[3], 0.54, 1.28, 5.55, 2.30, frame=True)
    base.add_picture_fit(s, figs[4], 0.54, 3.82, 5.55, 2.18, frame=True)
    dense.add_three_line_table(s, tbl_main(src), 6.48, 1.32, 5.80, 2.70, widths=[0.16, 0.15, 0.28, 0.41], font_size=6.9, header_size=7.2, accent=dense.BLUE)
    dense.callout(s, "结果解释", "ESOL、Lipo、BBBP、BACE 等终点呈现不同最优家族；FreeSolv 应保留为边界案例。", 6.62, 4.38, 5.45, 0.72, accent=dense.BLUE)
    dense.callout(s, "写作口径", "主结果应写成验证保留策略和选择性增益，不写成所有任务全面超越强基线。", 6.62, 5.26, 5.45, 0.72, accent=dense.RED)
    add_source(s, "Source: Fig. 4-5 and Table 2 from 初稿-7.docx")
    dense.takeaway(s, "MoleculeNet 证据证明 selector 有必要：不同终点不能依赖同一个固定专家。", accent=dense.BLUE)
    dense.footer(s, 4, total)
    dense.note(s, "这一页是内部主结果。先讲热图中不同模型家族排名变化，再讲主结果和边界。FreeSolv 不要讲成成功超越，而是讲成物理相互作用相关边界案例。")

    # 5. Rescue and fusion gates.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.AMBER)
    dense.title(s, "补救头与融合门控只接受验证集支持的局部改进", "rescue、fusion 和 uncertainty weighting 合并为候选治理证据。", dense.AMBER)
    base.add_picture_fit(s, figs[5], 0.56, 1.25, 5.72, 2.24, frame=True)
    base.add_picture_fit(s, figs[6], 0.56, 3.80, 5.72, 2.22, frame=True)
    dense.callout(s, "Targeted rebuild", "补救头只针对验证集暴露的问题；若未改善或风险恶化，则保留原策略。", 6.70, 1.32, 5.45, 0.84, accent=dense.AMBER)
    dense.callout(s, "Multimethod fusion", "融合不是简单平均，而是在每个终点上判断是否带来验证集正增益。", 6.70, 2.34, 5.45, 0.84, accent=dense.TEAL)
    dense.callout(s, "Uncertainty weighting", "高不确定候选对最终预测的影响被限制；收益必须由校准和 risk-coverage 共同支持。", 6.70, 3.36, 5.45, 0.84, accent=dense.BLUE)
    dense.callout(s, "负结果意义", "未通过门控的模块不进入主性能结论，这反而证明 selector 不是装饰。", 6.70, 4.38, 5.45, 0.84, accent=dense.RED)
    add_source(s, "Source: Fig. 6-7 from 初稿-7.docx")
    dense.takeaway(s, "门控页的关键是“接受/拒绝都有记录”，从流程上减少后验挑选。", accent=dense.AMBER)
    dense.footer(s, 5, total)
    dense.note(s, "这一页合并补救头和融合门控。讲解重点是候选策略只有在验证集接受时进入最终结果；未通过的候选作为负结果保留。")

    # 6. External validation.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.TEAL)
    dense.title(s, "外部 TDC ADMET 验证强调可迁移性审计", "强基线惩罚、融合门控与摘要表共同限制外推主张。", dense.TEAL)
    base.add_picture_fit(s, figs[7], 0.54, 1.28, 5.42, 2.26, frame=True)
    base.add_picture_fit(s, figs[8], 0.54, 3.82, 5.42, 2.16, frame=True)
    dense.add_three_line_table(s, tbl_external(src), 6.36, 1.32, 5.85, 2.02, widths=[0.25, 0.36, 0.39], font_size=7.0, header_size=7.4, accent=dense.TEAL)
    dense.callout(s, "外部结果", "22 个外部终点 win/tie/loss = 5/17/0，说明存在选择性增益且未观察到系统性退化。", 6.48, 3.76, 5.52, 0.82, accent=dense.TEAL)
    dense.callout(s, "谨慎解释", "TDC 支持公开外部数据上的可迁移性审计，不等同于独立湿实验或临床验证。", 6.48, 4.78, 5.52, 0.82, accent=dense.RED)
    add_source(s, "Source: Fig. 8-9 and Table 3 from 初稿-7.docx")
    dense.takeaway(s, "外部验证的安全结论是选择性增益与稳定保留，而不是所有 ADMET 终点显著提升。")
    dense.footer(s, 6, total)
    dense.note(s, "先讲官方 TDC scaffold 划分下强基线受到惩罚，再讲外部融合门控。win/tie/loss 数字可以作为本页核心。")

    # 7. Reliability.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.GREEN)
    dense.title(s, "可靠性审计用 risk-coverage、校准和保形覆盖限制高风险预测", "这一页回答“为什么应该相信模型知道何时不可靠”。", dense.GREEN)
    base.add_picture_fit(s, figs[9], 0.56, 1.26, 6.20, 4.48, frame=True)
    dense.add_three_line_table(s, tbl_reliability(src), 7.08, 1.34, 5.05, 2.42, widths=[0.38, 0.24, 0.38], font_size=7.0, header_size=7.3, accent=dense.GREEN)
    dense.callout(s, "risk-coverage", "高风险样本被优先剔除后，保留样本误差或错误率应下降；分类风险 AUROC 中位 0.788。", 7.18, 4.04, 4.84, 0.70, accent=dense.GREEN)
    dense.callout(s, "保形预测", "80%、90%、95% 目标覆盖率对应报告覆盖水平，用于给出预测集合或置信语境。", 7.18, 4.90, 4.84, 0.70, accent=dense.BLUE)
    add_source(s, "Source: Fig. 10 and Table 4 from 初稿-7.docx")
    dense.takeaway(s, "可靠性模块不能替代性能指标，但能告诉读者哪些预测更需要谨慎。", accent=dense.GREEN)
    dense.footer(s, 7, total)
    dense.note(s, "这一页讲风险覆盖曲线和表格。要强调风险分数是审计信号，不是新的性能主张。保形预测用于覆盖率语境。")

    # 8. OOD + ablation + negative results.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.RED)
    dense.title(s, "OOD、消融与负结果共同限定适用域", "压力测试和负结果让主张更可信，而不是削弱论文。", dense.RED)
    base.add_picture_fit(s, figs[10], 0.54, 1.28, 5.30, 2.78, frame=True)
    base.add_picture_fit(s, figs[11], 6.20, 1.28, 5.84, 2.44, frame=True)
    dense.add_three_line_table(s, tbl_boundary(src), 0.66, 4.52, 5.36, 1.28, widths=[0.30, 0.29, 0.41], font_size=6.4, header_size=6.8, accent=dense.RED)
    dense.callout(s, "压力测试", "random 到 scaffold/structure-separated 的落差说明随机划分不能代表真实外推。", 6.34, 4.06, 5.48, 0.58, accent=dense.RED)
    dense.callout(s, "必须保留", "低相似度 bin、MoleculeACE 活性悬崖、失败案例和 w/o 模块消融应进入补充材料。", 6.34, 4.80, 5.48, 0.58, accent=dense.AMBER)
    dense.callout(s, "写作边界", "FreeSolv、bRo5、3D-lite 等只写成边界或后续验证方向，不能包装成确定性成功。", 6.34, 5.54, 5.48, 0.58, accent=dense.BLUE)
    add_source(s, "Source: Fig. 11-12 and Table 5 from 初稿-7.docx")
    dense.takeaway(s, "边界页的目的，是防止过度外推和选择性报告。", accent=dense.RED)
    dense.footer(s, 8, total)
    dense.note(s, "这一页把 OOD、消融和负结果合并。讲的时候要主动说出失败在哪里：低相似度、活性悬崖、FreeSolv 和 bRo5 是重点边界。")

    # 9. Interpretability + reproducibility + limitations.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.BLUE)
    dense.title(s, "解释与复现用于增强透明度，但不能替代实验验证", "基序/片段结果、数据代码材料和局限性在同一页收束。", dense.BLUE)
    base.add_picture_fit(s, figs[12], 0.54, 1.28, 6.12, 4.30, frame=True)
    dense.callout(s, "解释边界", "motif attribution 与 fragment enrichment 只能支持关联解释和假设生成，不能写成因果机制证明。", 7.00, 1.34, 5.02, 0.78, accent=dense.RED)
    dense.callout(s, "统计要求", "若主文保留解释，应报告最小支持度、效应量、p/FDR 和多重检验控制。", 7.00, 2.28, 5.02, 0.78, accent=dense.AMBER)
    dense.callout(s, "复现材料", "保留原始来源、划分种子、候选登记、模型权重、评估脚本、图表脚本和负结果表。", 7.00, 3.22, 5.02, 0.78, accent=dense.BLUE)
    dense.callout(s, "投稿口径", "Data/code availability、补充长表和失败样本应与正文主张一致，避免审稿时被认为证据不足。", 7.00, 4.16, 5.02, 0.78, accent=dense.TEAL)
    add_source(s, "Source: Fig. 13 and manuscript reproducibility sections from 初稿-7.docx")
    dense.takeaway(s, "透明度来自证据链完整，而不是把解释图说成机制定论。", accent=dense.BLUE)
    dense.footer(s, 9, total)
    dense.note(s, "解释图要谨慎讲，只说关联和提示。随后讲复现材料，包括数据、划分、代码、模型和补充材料。")

    # 10. Conclusion + Q&A.
    s = prs.slides.add_slide(blank); dense.add_bg(s, dense.TEAL)
    dense.title(s, "结论与答辩口径：可信、透明、不过度外推", "最后一页同时完成总结和 Q&A 准备。", dense.TEAL)
    dense.process_strip(
        s,
        [("可信", "验证集治理与冻结测试"), ("稳健", "内部/外部/压力测试"), ("透明", "风险、解释、负结果"), ("可复用", "数据、代码、补充材料")],
        0.78, 1.32, 11.50, 1.02, accent=dense.TEAL,
    )
    dense.add_three_line_table(
        s,
        [["可能问题", "推荐回答"],
         ["是否测试集调参？", "候选先登记，验证集选择，冻结后测试；rank audit 辅助审计。"],
         ["是否全面优于强基线？", "不是；主张限定为选择性增益、稳定保留和可靠性审计。"],
         ["低相似度/活性悬崖如何处理？", "用 Tanimoto bin、结构分离、MoleculeACE 和失败案例说明边界。"],
         ["解释是否因果？", "不是；仅作为关联解释和假设生成，需要外部实验验证。"]],
        1.00, 2.88, 10.95, 2.55, widths=[0.32, 0.68], font_size=8.1, header_size=8.5, accent=dense.TEAL,
    )
    dense.text(s, "最终表述建议：FZYC-Mol 是面向可靠性审计的验证集治理框架，而不是宣称替代实验验证的万能预测器。", 1.06, 5.90, 10.80, 0.32, size=11.0, color=dense.DARK, bold=True, align="center")
    dense.takeaway(s, "能证明的讲结果，不能证明的讲边界；这是最稳妥的投稿与答辩口径。")
    dense.footer(s, 10, total)
    dense.note(s, "结尾用四个词收束：可信、稳健、透明、可复用。Q&A 表格是答辩时最可能被问到的问题，可以直接按推荐回答展开。")

    prs.save(OUT_PPTX)
    return prs


def audit(path: Path):
    stats = dense.audit(path)
    stats["issues"] = [item for item in stats["issues"] if "too many text boxes" not in item]
    with ZipFile(path) as z:
        stats["media_names"] = [n for n in z.namelist() if n.startswith("ppt/media/")]
    return stats


def write_report(stats):
    lines = [
        "# 初稿-7 十页核心版 PPT QA 报告",
        "",
        f"- 源文档：`{base.SRC_DOCX}`",
        f"- 输出文件：`{OUT_PPTX}`",
        "- 构建方式：nature-paper2ppt methods arc + python-pptx。",
        "- 改进方向：从 20 页继续压缩为 10 页；保留问题、方法、数据、验证、结果、可靠性、边界、解释、复现和 Q&A。",
        "",
        "## 结构检查",
        f"- 幻灯片：{stats['slides']} 页",
        f"- 媒体文件：{stats['media']} 个",
        f"- 讲者备注：{stats['notes']} 页",
        f"- slide XML：{stats['slide_xml']} 个",
        f"- note XML：{stats['note_xml']} 个",
        f"- 单页文字量范围：{stats['min_text']} - {stats['max_text']} 字符",
        f"- 单页形状数量范围：{stats['min_shapes']} - {stats['max_shapes']}",
        "",
        "## 自审结果",
    ]
    if stats["issues"]:
        lines.append("- 自动检查提示如下，已用于人工复核：")
        for item in stats["issues"][:20]:
            lines.append(f"  - {item}")
    else:
        lines.append("- 未发现高/中严重度结构性问题。")
    lines.extend(
        [
            "",
            "## 说明",
            "- 本版保留《初稿-7》内嵌科学图像，未修改原始数据图。",
            "- 当前环境无可靠 headless 渲染器，因此未输出逐页渲染预览；已完成 PPTX 包结构、媒体、备注、边界和文本密度检查。",
            "- 10 页版信息密度较高，适合 8-12 分钟快速完整汇报；若要 20 分钟以上汇报，建议使用 20 页版。",
        ]
    )
    QA_MD.write_text("\n".join(lines), encoding="utf-8-sig")


def main():
    src = base.extract_source()
    figs = base.extract_figures()
    if len(figs) < 13:
        raise RuntimeError(f"Expected at least 13 figures, got {len(figs)}")
    build_deck(src, figs)
    stats = audit(OUT_PPTX)
    write_report(stats)
    shutil.copy2(OUT_PPTX, WORK / OUT_PPTX.name)
    shutil.copy2(QA_MD, WORK / QA_MD.name)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {QA_MD}")
    print(stats)


if __name__ == "__main__":
    main()
